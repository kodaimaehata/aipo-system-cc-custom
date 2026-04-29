"""Codex-driven image anonymization for pptx pictures.

For each picture in a pptx, we ask the Codex CLI (with the image attached via
``-i``) whether the image contains any of the approved original identifiers.
When it does, Codex produces an anonymized PNG which we splice back into the
package by **cloning** the affected image part — never by mutating a shared
blob, which would ripple into every slide that referenced the same media.

Contract summary (documented in full in ``references/image-codex-flow.md``):

  - codex command:      ``codex exec --skip-git-repo-check --sandbox read-only -i <png> -``
  - model default:      ``gpt-5-codex`` (override via ``model=`` kwarg)
  - per-image timeout:  300 s
  - retries:            2 (total = 3 attempts)
  - max images per run: 20 (raises ``CodexBudgetExceeded`` beyond that)
  - failure policy:     log + keep original image; never abort the pipeline
"""

from __future__ import annotations

import hashlib
import os
import re
import shutil
import subprocess
import time
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Mapping


DEFAULT_MODEL = "gpt-5-codex"
DEFAULT_TIMEOUT_SEC = 300
DEFAULT_RETRIES = 2
DEFAULT_MAX_IMAGES = 20


class CodexBudgetExceeded(RuntimeError):
    """Raised when the per-run image budget would be exceeded."""


@dataclass
class ImageOutcome:
    image_id: str
    source_hash: str
    status: str                         # "unchanged" | "replaced" | "failed" | "skipped"
    slides_affected: list[int] = field(default_factory=list)
    anonymized_path: Path | None = None
    notes: str = ""


def anonymize_pptx_images(
    *,
    pptx_path: Path,
    approved_mapping: Mapping[str, str],
    work_dir: Path,
    model: str = DEFAULT_MODEL,
    timeout_sec: int = DEFAULT_TIMEOUT_SEC,
    retries: int = DEFAULT_RETRIES,
    max_images: int = DEFAULT_MAX_IMAGES,
) -> list[ImageOutcome]:
    """Walk every pptx picture and, when appropriate, substitute an anonymized copy.

    ``work_dir`` must already exist and have 0o700 permissions. This function
    writes extracted originals, codex inputs/outputs, and a sibling
    ``image_codex_log.md`` into ``work_dir``.
    """
    from pptx import Presentation  # local import: optional dep

    work_dir = Path(work_dir)
    work_dir.mkdir(parents=True, exist_ok=True)
    log_path = work_dir / "image_codex_log.md"

    prs = Presentation(str(pptx_path))
    picture_index = _enumerate_pictures(prs)
    if not picture_index:
        _append_log(log_path, ["No pictures found — nothing to do."])
        return []

    if len(picture_index) > max_images:
        raise CodexBudgetExceeded(
            f"pptx has {len(picture_index)} pictures; exceeds max_images={max_images}."
        )

    outcomes: list[ImageOutcome] = []
    for idx, (image_hash, shapes) in enumerate(picture_index.items(), start=1):
        outcome = _process_one_image(
            image_id=f"img-{idx:02d}",
            image_hash=image_hash,
            shapes=shapes,
            approved_mapping=approved_mapping,
            work_dir=work_dir,
            model=model,
            timeout_sec=timeout_sec,
            retries=retries,
        )
        outcomes.append(outcome)
        _append_log(log_path, [_format_outcome(outcome)])
        if outcome.status == "replaced" and outcome.anonymized_path:
            _swap_image_in_place(prs, shapes, outcome.anonymized_path)

    prs.save(str(pptx_path))
    return outcomes


def _enumerate_pictures(prs) -> dict[str, list]:
    """Return {blob_sha1 -> [(slide_idx, shape)]} for every picture in the deck."""
    picture_index: dict[str, list] = {}
    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "shape_type", None) != 13:  # PICTURE
                continue
            try:
                blob = shape.image.blob
            except Exception:
                continue
            digest = hashlib.sha1(blob).hexdigest()
            picture_index.setdefault(digest, []).append((slide_idx, shape))
    return picture_index


def _process_one_image(
    *,
    image_id: str,
    image_hash: str,
    shapes,
    approved_mapping: Mapping[str, str],
    work_dir: Path,
    model: str,
    timeout_sec: int,
    retries: int,
) -> ImageOutcome:
    slide_idx, example_shape = shapes[0]
    try:
        blob = example_shape.image.blob
        ext = (example_shape.image.ext or "png").lower()
    except Exception as exc:
        return ImageOutcome(
            image_id=image_id,
            source_hash=image_hash,
            status="failed",
            slides_affected=[s for s, _ in shapes],
            notes=f"blob extraction failed: {exc}",
        )

    source_path = work_dir / f"{image_id}_source.{ext}"
    source_path.write_bytes(blob)
    os.chmod(source_path, 0o600)

    anonymized_path = work_dir / f"{image_id}_anonymized.{ext}"
    prompt_path = _write_prompt(work_dir, image_id, approved_mapping)
    for attempt in range(retries + 1):
        try:
            findings = _invoke_codex(
                prompt_path=prompt_path,
                image_path=source_path,
                output_path=anonymized_path,
                model=model,
                timeout_sec=timeout_sec,
            )
        except subprocess.TimeoutExpired:
            if attempt == retries:
                return ImageOutcome(
                    image_id=image_id,
                    source_hash=image_hash,
                    status="failed",
                    slides_affected=[s for s, _ in shapes],
                    notes=f"codex timed out after {retries + 1} attempts",
                )
            time.sleep(2 ** attempt)
            continue
        except subprocess.CalledProcessError as exc:
            if attempt == retries:
                return ImageOutcome(
                    image_id=image_id,
                    source_hash=image_hash,
                    status="failed",
                    slides_affected=[s for s, _ in shapes],
                    notes=f"codex exited {exc.returncode}: {exc.stderr[:200] if exc.stderr else ''}",
                )
            time.sleep(2 ** attempt)
            continue
        break
    else:
        return ImageOutcome(
            image_id=image_id,
            source_hash=image_hash,
            status="failed",
            slides_affected=[s for s, _ in shapes],
            notes="exhausted retries",
        )

    if findings == 0 or not anonymized_path.exists():
        return ImageOutcome(
            image_id=image_id,
            source_hash=image_hash,
            status="unchanged",
            slides_affected=[s for s, _ in shapes],
            notes=f"codex reported 0 findings (first slide={slide_idx})",
        )
    if hashlib.sha1(anonymized_path.read_bytes()).hexdigest() == image_hash:
        return ImageOutcome(
            image_id=image_id,
            source_hash=image_hash,
            status="unchanged",
            slides_affected=[s for s, _ in shapes],
            notes="codex returned byte-identical image",
        )

    return ImageOutcome(
        image_id=image_id,
        source_hash=image_hash,
        status="replaced",
        slides_affected=[s for s, _ in shapes],
        anonymized_path=anonymized_path,
        notes=f"codex reported {findings} replacement(s)",
    )


def _write_prompt(work_dir: Path, image_id: str, approved_mapping: Mapping[str, str]) -> Path:
    mapping_lines = [
        f"- `{original}` -> `{replacement}`"
        for original, replacement in approved_mapping.items()
    ]
    prompt_path = work_dir / f"{image_id}_prompt.md"
    prompt_path.write_text(
        "\n".join(
            [
                "# Image anonymization task",
                "",
                "You are given one image from a business presentation. Produce an",
                "anonymized PNG where only the substrings listed below are masked",
                "and redrawn with their replacement text. All other pixels must",
                "stay byte-identical to the source.",
                "",
                "Save the output to the path given on stdin as `output=<path>`.",
                "If nothing in the image matches the list, copy the source image",
                "to that path unchanged and print `findings=0`.",
                "",
                "Print exactly one line of the form `findings=N` to stdout before",
                "exiting, where N is the count of replacements you applied.",
                "",
                "## Mapping",
                "",
                *mapping_lines,
            ]
        ),
        encoding="utf-8",
    )
    os.chmod(prompt_path, 0o600)
    return prompt_path


def _invoke_codex(
    *,
    prompt_path: Path,
    image_path: Path,
    output_path: Path,
    model: str,
    timeout_sec: int,
) -> int:
    stdin_input = (
        prompt_path.read_text(encoding="utf-8")
        + f"\n\noutput={output_path}\n"
    )
    cmd = [
        "codex",
        "exec",
        "--skip-git-repo-check",
        "--sandbox",
        "read-only",
        "-i",
        str(image_path),
        "-m",
        model,
        "-",
    ]
    completed = subprocess.run(
        cmd,
        input=stdin_input,
        text=True,
        capture_output=True,
        timeout=timeout_sec,
        check=True,
    )
    return _parse_findings_count(completed.stdout)


def _parse_findings_count(stdout: str) -> int:
    match = re.search(r"findings\s*=\s*(\d+)", stdout)
    if not match:
        return 0
    try:
        return int(match.group(1))
    except ValueError:
        return 0


def _swap_image_in_place(prs, shapes, anonymized_path: Path) -> None:
    """Clone the image part per affected slide and rewrite the blip rId there only."""
    new_blob = anonymized_path.read_bytes()
    for _, shape in shapes:
        slide_part = shape.part
        try:
            image_part = slide_part.package.get_or_add_image_part(new_blob)
        except AttributeError:
            image_part = _fallback_add_image_part(slide_part, new_blob, anonymized_path)
        new_rid = slide_part.relate_to(image_part, _PIC_REL_TYPE)
        shape._element.blipFill_blip.rEmbed = new_rid  # type: ignore[attr-defined]


def _fallback_add_image_part(slide_part, blob: bytes, anonymized_path: Path):
    """Manually add an image part when the high-level helper is unavailable."""
    from pptx.opc.package import PartFactory  # local import
    from pptx.opc.packuri import PackURI
    from pptx.parts.image import ImagePart

    existing = [
        p.partname for p in slide_part.package.iter_parts()
        if isinstance(p, ImagePart)
    ]
    next_idx = 1
    while any(str(pn).endswith(f"/image{next_idx}{anonymized_path.suffix}") for pn in existing):
        next_idx += 1
    partname = PackURI(f"/ppt/media/image{next_idx}{anonymized_path.suffix}")
    image_part = ImagePart.new(slide_part.package, partname, blob)
    return image_part


_PIC_REL_TYPE = (
    "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image"
)


def _append_log(log_path: Path, lines: Iterable[str]) -> None:
    existed = log_path.exists()
    with log_path.open("a", encoding="utf-8") as fh:
        if not existed:
            fh.write("# Image Codex Log\n\n")
        for line in lines:
            fh.write(line + "\n")
    try:
        os.chmod(log_path, 0o600)
    except PermissionError:
        pass


def _format_outcome(outcome: ImageOutcome) -> str:
    affected = ",".join(str(s) for s in outcome.slides_affected)
    return (
        f"- {outcome.image_id} | sha1={outcome.source_hash[:8]} | "
        f"status={outcome.status} | slides={affected} | {outcome.notes}"
    )
