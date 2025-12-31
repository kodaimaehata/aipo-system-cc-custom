#!/usr/bin/env python3
"""Create sample template for pptx-from-template skill."""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pathlib import Path


def create_sample_template(output_path: str | Path) -> None:
    """Create a sample template with placeholders.

    Args:
        output_path: Path to save the template
    """
    prs = Presentation()

    # Slide 1: Title Slide (layout 0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "{{title}}"
    slide1.placeholders[1].text = "{{subtitle}}\n{{date}}"

    # Slide 2: Section Header (layout 2)
    slide2 = prs.slides.add_slide(prs.slide_layouts[2])
    slide2.shapes.title.text = "{{section_title}}"

    # Slide 3: Title and Content (layout 1)
    slide3 = prs.slides.add_slide(prs.slide_layouts[1])
    slide3.shapes.title.text = "{{content_title}}"
    body3 = slide3.placeholders[1]
    tf3 = body3.text_frame
    tf3.text = "{{bullet_1}}"
    p2 = tf3.add_paragraph()
    p2.text = "{{bullet_2}}"
    p3 = tf3.add_paragraph()
    p3.text = "{{bullet_3}}"

    # Slide 4: Title Only (layout 5) - for custom content
    slide4 = prs.slides.add_slide(prs.slide_layouts[5])
    slide4.shapes.title.text = "{{table_title}}"

    # Add a placeholder text box for table location hint
    txBox = slide4.shapes.add_textbox(Inches(0.5), Inches(2), Inches(9), Inches(0.5))
    tf = txBox.text_frame
    tf.paragraphs[0].text = "[表データはここに挿入されます]"
    tf.paragraphs[0].font.size = Pt(14)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)

    # Slide 5: Blank (layout 6) - for images
    slide5 = prs.slides.add_slide(prs.slide_layouts[6])

    # Add title text box
    title_box = slide5.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
    tf = title_box.text_frame
    tf.paragraphs[0].text = "{{image_title}}"
    tf.paragraphs[0].font.size = Pt(32)
    tf.paragraphs[0].font.bold = True

    # Add placeholder for image location
    img_box = slide5.shapes.add_textbox(Inches(1.5), Inches(2), Inches(7), Inches(3))
    tf = img_box.text_frame
    tf.paragraphs[0].text = "[画像はここに挿入されます]"
    tf.paragraphs[0].font.size = Pt(18)
    tf.paragraphs[0].font.color.rgb = RGBColor(0x80, 0x80, 0x80)
    tf.paragraphs[0].alignment = PP_ALIGN.CENTER

    # Save template
    prs.save(str(output_path))
    print(f"✓ テンプレート作成: {output_path}")


def create_sample_data(output_path: str | Path) -> None:
    """Create sample data file for testing.

    Args:
        output_path: Path to save the data file
    """
    import json

    data = {
        "metadata": {
            "title": "サンプルプレゼンテーション",
            "author": "Claude Code",
            "date": "2025-12-31"
        },
        "placeholders": {
            "title": "月次報告書",
            "subtitle": "2025年1月度",
            "date": "2025-01-15",
            "section_title": "売上報告",
            "content_title": "今月のハイライト",
            "bullet_1": "売上目標達成率: 105%",
            "bullet_2": "新規顧客獲得: 15社",
            "bullet_3": "リピート率向上: +3%",
            "table_title": "部門別売上",
            "image_title": "売上推移グラフ"
        },
        "slides": [
            {
                "layout": 0,
                "title": "月次報告書",
                "subtitle": "2025年1月度\n作成: 営業部"
            },
            {
                "layout": 2,
                "title": "売上報告"
            },
            {
                "layout": 1,
                "title": "今月のハイライト",
                "content": [
                    "売上目標達成率: 105%",
                    "新規顧客獲得: 15社",
                    "リピート率向上: +3%",
                    "顧客満足度: 4.5/5.0"
                ]
            },
            {
                "layout": 5,
                "title": "部門別売上",
                "table": {
                    "headers": ["部門", "目標", "実績", "達成率"],
                    "rows": [
                        ["営業第一部", "500万", "520万", "104%"],
                        ["営業第二部", "300万", "315万", "105%"],
                        ["オンライン", "200万", "225万", "112%"]
                    ]
                }
            },
            {
                "layout": 1,
                "title": "来月の計画",
                "content": [
                    "新製品キャンペーン開始",
                    "既存顧客フォローアップ強化",
                    "展示会出展準備"
                ]
            }
        ]
    }

    with open(output_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)

    print(f"✓ サンプルデータ作成: {output_path}")


if __name__ == "__main__":
    script_dir = Path(__file__).parent.parent
    templates_dir = script_dir / "templates"
    templates_dir.mkdir(exist_ok=True)

    create_sample_template(templates_dir / "sample_template.pptx")
    create_sample_data(templates_dir / "sample_data.json")
