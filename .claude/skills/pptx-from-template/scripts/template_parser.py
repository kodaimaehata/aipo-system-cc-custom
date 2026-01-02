"""Template parsing for pptx-from-template.

Handles template loading and placeholder replacement.
"""

import re
from pathlib import Path
from typing import Any

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor


class TemplateError(Exception):
    """Template handling error."""

    def __init__(self, code: str, message: str, details: str = ""):
        self.code = code
        self.message = message
        self.details = details
        super().__init__(f"{code}: {message}")


def load_template(template_path: str | Path) -> Presentation:
    """Load a PowerPoint template.

    Args:
        template_path: Path to the template file

    Returns:
        Presentation object

    Raises:
        TemplateError: If file not found or invalid
    """
    path = Path(template_path)

    if not path.exists():
        raise TemplateError(
            "E001",
            "テンプレートファイルが見つかりません",
            f"パス: {path}",
        )

    if path.suffix.lower() != ".pptx":
        raise TemplateError(
            "E006",
            "有効なPPTXファイルではありません",
            f"拡張子: {path.suffix}",
        )

    try:
        return Presentation(str(path))
    except Exception as e:
        raise TemplateError(
            "E006",
            "テンプレートの読み込みに失敗しました",
            str(e),
        )


def create_new_presentation() -> Presentation:
    """Create a new empty presentation.

    Returns:
        New Presentation object
    """
    return Presentation()


def replace_placeholders(prs: Presentation, replacements: dict[str, str]) -> None:
    """Replace {{key}} placeholders in all slides.

    Args:
        prs: Presentation object
        replacements: Dictionary of placeholder key -> replacement value
    """
    pattern = re.compile(r"\{\{(\w+)\}\}")

    for slide in prs.slides:
        for shape in slide.shapes:
            if not hasattr(shape, "text_frame"):
                continue

            for paragraph in shape.text_frame.paragraphs:
                # Try run-level replacement first (preserves formatting)
                for run in paragraph.runs:
                    for key, value in replacements.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in run.text:
                            run.text = run.text.replace(placeholder, value)

                # Fallback to paragraph-level if runs are empty
                if not paragraph.runs:
                    for key, value in replacements.items():
                        placeholder = f"{{{{{key}}}}}"
                        if placeholder in paragraph.text:
                            paragraph.text = paragraph.text.replace(placeholder, value)


def add_slide_from_data(prs: Presentation, slide_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a slide based on data specification.

    Args:
        prs: Presentation object
        slide_data: Slide specification dictionary
        warnings: List to append warnings to
    """
    # Get layout index (default: 1 = Title and Content)
    layout_idx = slide_data.get("layout", 1)
    if not isinstance(layout_idx, int) or layout_idx < 0:
        layout_idx = 1

    # Ensure layout index is within bounds
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = min(1, len(prs.slide_layouts) - 1)

    slide_layout = prs.slide_layouts[layout_idx]
    slide = prs.slides.add_slide(slide_layout)

    # Set title
    if "title" in slide_data and slide.shapes.title:
        slide.shapes.title.text = str(slide_data["title"])

    # Set subtitle (for title slide layout)
    if "subtitle" in slide_data:
        try:
            if len(slide.placeholders) > 1:
                slide.placeholders[1].text = str(slide_data["subtitle"])
        except (KeyError, IndexError):
            pass

    # Set content (bullet points)
    if "content" in slide_data:
        content = slide_data["content"]
        if isinstance(content, list) and len(slide.placeholders) > 1:
            try:
                body = slide.placeholders[1]
                tf = body.text_frame
                tf.clear()

                for i, item in enumerate(content):
                    if i == 0:
                        tf.paragraphs[0].text = str(item)
                    else:
                        p = tf.add_paragraph()
                        p.text = str(item)
                        p.level = 0
            except (KeyError, IndexError, AttributeError):
                pass
        elif isinstance(content, str):
            try:
                if len(slide.placeholders) > 1:
                    slide.placeholders[1].text = content
            except (KeyError, IndexError):
                pass

    # Add table
    if "table" in slide_data:
        _add_table_to_slide(slide, slide_data["table"], warnings)

    # Add image
    if "image" in slide_data:
        _add_image_to_slide(slide, slide_data["image"], warnings)

    # Add free-form shapes
    if "shapes" in slide_data:
        _add_shapes_to_slide(slide, slide_data["shapes"], warnings)

    # Add speaker notes
    if "notes" in slide_data:
        notes_text = str(slide_data["notes"])
        if notes_text:
            notes_slide = slide.notes_slide
            notes_slide.notes_text_frame.text = notes_text


def _add_table_to_slide(slide, table_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a table to the slide.

    Args:
        slide: Slide object
        table_data: Table specification with headers and rows
        warnings: List to append warnings to
    """
    headers = table_data.get("headers", [])
    rows = table_data.get("rows", [])

    if not headers and not rows:
        return

    # Calculate dimensions
    cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    row_count = len(rows) + (1 if headers else 0)

    if cols == 0 or row_count == 0:
        return

    # Add table shape
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9.0)
    height = Inches(0.8 * row_count)

    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table

    # Set column widths
    col_width = Inches(9.0 / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Fill headers
    if headers:
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            # Bold header
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    # Fill rows
    start_row = 1 if headers else 0
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            if col_idx < cols:
                table.cell(start_row + row_idx, col_idx).text = str(value)


def _add_image_to_slide(slide, image_path: str, warnings: list[str]) -> None:
    """Add an image to the slide.

    Args:
        slide: Slide object
        image_path: Path to the image file
        warnings: List to append warnings to
    """
    path = Path(image_path)
    if not path.exists():
        # Warning already added in validate_data
        return

    try:
        # Add image centered on slide
        left = Inches(1.5)
        top = Inches(2.0)
        width = Inches(7.0)
        slide.shapes.add_picture(str(path), left, top, width=width)
    except Exception as e:
        warnings.append(f"W001: 画像の追加に失敗しました: {e}")


def get_layout_name(prs: Presentation, layout_idx: int) -> str:
    """Get the name of a slide layout.

    Args:
        prs: Presentation object
        layout_idx: Layout index

    Returns:
        Layout name or "Unknown"
    """
    try:
        if 0 <= layout_idx < len(prs.slide_layouts):
            return prs.slide_layouts[layout_idx].name
    except Exception:
        pass
    return "Unknown"


# =============================================================================
# Free-form shapes support
# =============================================================================

# Shape type mapping
SHAPE_TYPE_MAP = {
    "rectangle": MSO_SHAPE.RECTANGLE,
    "rounded_rectangle": MSO_SHAPE.ROUNDED_RECTANGLE,
    "oval": MSO_SHAPE.OVAL,
    "circle": MSO_SHAPE.OVAL,
    "triangle": MSO_SHAPE.ISOSCELES_TRIANGLE,
    "right_arrow": MSO_SHAPE.RIGHT_ARROW,
    "left_arrow": MSO_SHAPE.LEFT_ARROW,
    "up_arrow": MSO_SHAPE.UP_ARROW,
    "down_arrow": MSO_SHAPE.DOWN_ARROW,
    "pentagon": MSO_SHAPE.PENTAGON,
    "hexagon": MSO_SHAPE.HEXAGON,
    "star": MSO_SHAPE.STAR_5_POINT,
    "callout": MSO_SHAPE.RECTANGULAR_CALLOUT,
    "cloud": MSO_SHAPE.CLOUD,
    "heart": MSO_SHAPE.HEART,
    "lightning": MSO_SHAPE.LIGHTNING_BOLT,
}


def _parse_color(color_str: str) -> RGBColor | None:
    """Parse color string to RGBColor.

    Args:
        color_str: Color in #RRGGBB format or color name

    Returns:
        RGBColor object or None if invalid
    """
    if not color_str:
        return None

    # Handle hex color
    if color_str.startswith("#"):
        try:
            hex_color = color_str.lstrip("#")
            if len(hex_color) == 6:
                r = int(hex_color[0:2], 16)
                g = int(hex_color[2:4], 16)
                b = int(hex_color[4:6], 16)
                return RGBColor(r, g, b)
        except ValueError:
            pass

    # Common color names
    color_names = {
        "red": RGBColor(255, 0, 0),
        "green": RGBColor(0, 128, 0),
        "blue": RGBColor(0, 0, 255),
        "yellow": RGBColor(255, 255, 0),
        "orange": RGBColor(255, 165, 0),
        "purple": RGBColor(128, 0, 128),
        "black": RGBColor(0, 0, 0),
        "white": RGBColor(255, 255, 255),
        "gray": RGBColor(128, 128, 128),
        "grey": RGBColor(128, 128, 128),
    }
    return color_names.get(color_str.lower())


def _add_shapes_to_slide(slide, shapes: list[dict[str, Any]], warnings: list[str]) -> None:
    """Add free-form shapes to the slide.

    Args:
        slide: Slide object
        shapes: List of shape specifications
        warnings: List to append warnings to
    """
    for i, shape_data in enumerate(shapes):
        shape_type = shape_data.get("type", "").lower()

        try:
            if shape_type == "textbox":
                _add_textbox_shape(slide, shape_data, warnings)
            elif shape_type == "image":
                _add_image_shape(slide, shape_data, warnings)
            elif shape_type == "shape":
                _add_basic_shape(slide, shape_data, warnings)
            elif shape_type == "table":
                _add_table_shape(slide, shape_data, warnings)
            elif shape_type == "line":
                _add_line_shape(slide, shape_data, warnings)
            else:
                warnings.append(f"W005: shapes[{i}]: 未知のshapeタイプ '{shape_type}'")
        except Exception as e:
            warnings.append(f"W005: shapes[{i}]: 図形の追加に失敗しました: {e}")


def _add_textbox_shape(slide, shape_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a textbox shape.

    Args:
        slide: Slide object
        shape_data: Shape specification
        warnings: List to append warnings to
    """
    left = Inches(shape_data.get("left", 0))
    top = Inches(shape_data.get("top", 0))
    width = Inches(shape_data.get("width", 2))
    height = Inches(shape_data.get("height", 1))

    textbox = slide.shapes.add_textbox(left, top, width, height)
    tf = textbox.text_frame
    tf.word_wrap = True

    text = shape_data.get("text", "")
    tf.text = str(text)

    # Apply text formatting
    font_size = shape_data.get("font_size")
    bold = shape_data.get("bold", False)
    italic = shape_data.get("italic", False)
    font_color = shape_data.get("font_color")
    align = shape_data.get("align", "left")

    for paragraph in tf.paragraphs:
        # Alignment
        if align == "center":
            paragraph.alignment = PP_ALIGN.CENTER
        elif align == "right":
            paragraph.alignment = PP_ALIGN.RIGHT
        else:
            paragraph.alignment = PP_ALIGN.LEFT

        for run in paragraph.runs:
            if font_size:
                run.font.size = Pt(font_size)
            run.font.bold = bold
            run.font.italic = italic
            if font_color:
                color = _parse_color(font_color)
                if color:
                    run.font.color.rgb = color

    # Background fill
    fill_color = shape_data.get("fill_color")
    if fill_color:
        color = _parse_color(fill_color)
        if color:
            textbox.fill.solid()
            textbox.fill.fore_color.rgb = color


def _add_image_shape(slide, shape_data: dict[str, Any], warnings: list[str]) -> None:
    """Add an image shape with custom position.

    Args:
        slide: Slide object
        shape_data: Shape specification
        warnings: List to append warnings to
    """
    image_path = shape_data.get("path", "")
    if not image_path:
        warnings.append("W001: 画像パスが指定されていません")
        return

    path = Path(image_path)
    if not path.exists():
        warnings.append(f"W001: 画像ファイルが見つかりません ({image_path})")
        return

    left = Inches(shape_data.get("left", 0))
    top = Inches(shape_data.get("top", 0))

    # Width and height - at least one should be specified
    width = shape_data.get("width")
    height = shape_data.get("height")

    try:
        if width and height:
            slide.shapes.add_picture(str(path), left, top, Inches(width), Inches(height))
        elif width:
            slide.shapes.add_picture(str(path), left, top, width=Inches(width))
        elif height:
            slide.shapes.add_picture(str(path), left, top, height=Inches(height))
        else:
            # Default width
            slide.shapes.add_picture(str(path), left, top, width=Inches(4))
    except Exception as e:
        warnings.append(f"W001: 画像の追加に失敗しました: {e}")


def _add_basic_shape(slide, shape_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a basic shape (rectangle, oval, arrow, etc.).

    Args:
        slide: Slide object
        shape_data: Shape specification
        warnings: List to append warnings to
    """
    shape_type_name = shape_data.get("shape_type", "rectangle").lower()
    shape_type = SHAPE_TYPE_MAP.get(shape_type_name, MSO_SHAPE.RECTANGLE)

    left = Inches(shape_data.get("left", 0))
    top = Inches(shape_data.get("top", 0))
    width = Inches(shape_data.get("width", 2))
    height = Inches(shape_data.get("height", 1))

    shape = slide.shapes.add_shape(shape_type, left, top, width, height)

    # Fill color
    fill_color = shape_data.get("fill_color")
    if fill_color:
        color = _parse_color(fill_color)
        if color:
            shape.fill.solid()
            shape.fill.fore_color.rgb = color

    # Line/border color
    line_color = shape_data.get("line_color")
    if line_color:
        color = _parse_color(line_color)
        if color:
            shape.line.color.rgb = color

    # Line width
    line_width = shape_data.get("line_width")
    if line_width:
        shape.line.width = Pt(line_width)

    # Text inside shape
    text = shape_data.get("text")
    if text:
        shape.text = str(text)
        tf = shape.text_frame
        tf.word_wrap = True

        font_size = shape_data.get("font_size")
        font_color = shape_data.get("font_color")
        align = shape_data.get("align", "center")

        for paragraph in tf.paragraphs:
            if align == "center":
                paragraph.alignment = PP_ALIGN.CENTER
            elif align == "right":
                paragraph.alignment = PP_ALIGN.RIGHT
            else:
                paragraph.alignment = PP_ALIGN.LEFT

            for run in paragraph.runs:
                if font_size:
                    run.font.size = Pt(font_size)
                if font_color:
                    color = _parse_color(font_color)
                    if color:
                        run.font.color.rgb = color


def _add_table_shape(slide, shape_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a table at custom position.

    Args:
        slide: Slide object
        shape_data: Shape specification with headers and rows
        warnings: List to append warnings to
    """
    headers = shape_data.get("headers", [])
    rows = shape_data.get("rows", [])

    if not headers and not rows:
        return

    cols = len(headers) if headers else (len(rows[0]) if rows else 0)
    row_count = len(rows) + (1 if headers else 0)

    if cols == 0 or row_count == 0:
        return

    left = Inches(shape_data.get("left", 0.5))
    top = Inches(shape_data.get("top", 2.0))
    width = Inches(shape_data.get("width", 9.0))
    height = Inches(shape_data.get("height", 0.8 * row_count))

    table = slide.shapes.add_table(row_count, cols, left, top, width, height).table

    # Set column widths
    col_width = int(width / cols)
    for i in range(cols):
        table.columns[i].width = col_width

    # Fill headers
    if headers:
        for i, header in enumerate(headers):
            cell = table.cell(0, i)
            cell.text = str(header)
            for paragraph in cell.text_frame.paragraphs:
                for run in paragraph.runs:
                    run.font.bold = True

    # Fill rows
    start_row = 1 if headers else 0
    for row_idx, row in enumerate(rows):
        for col_idx, value in enumerate(row):
            if col_idx < cols:
                table.cell(start_row + row_idx, col_idx).text = str(value)


def _add_line_shape(slide, shape_data: dict[str, Any], warnings: list[str]) -> None:
    """Add a line connector.

    Args:
        slide: Slide object
        shape_data: Shape specification with start and end points
        warnings: List to append warnings to
    """
    start_x = Inches(shape_data.get("start_x", 0))
    start_y = Inches(shape_data.get("start_y", 0))
    end_x = Inches(shape_data.get("end_x", 2))
    end_y = Inches(shape_data.get("end_y", 0))

    # Calculate width and height from start/end points
    left = min(start_x, end_x)
    top = min(start_y, end_y)
    width = abs(end_x - start_x)
    height = abs(end_y - start_y)

    # Use a line shape
    from pptx.enum.shapes import MSO_CONNECTOR
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT, start_x, start_y, end_x, end_y
    )

    # Line color
    line_color = shape_data.get("line_color")
    if line_color:
        color = _parse_color(line_color)
        if color:
            connector.line.color.rgb = color

    # Line width
    line_width = shape_data.get("line_width", 1)
    connector.line.width = Pt(line_width)
