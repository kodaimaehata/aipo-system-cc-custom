"""Heuristic candidate scanner for office-anonymizer.

Surfaces person-name / company-name / property-name candidates from pptx files
so the skill can propose a mapping *before* the SG5 detector runs. SG5's
body-text detector is hint-based and returns zero candidates when the caller
provides no hints; this scanner fills that gap.

Only minimal context (candidate + category + occurrences + location coordinates)
is persisted. Raw surrounding snippets are NOT stored by default; callers can
pass ``with_context=True`` to opt into snippet capture at their own risk.
"""

from __future__ import annotations

import re
from collections import Counter
from dataclasses import dataclass, field
from pathlib import Path
from typing import Iterable, Iterator


# 職位後続語 that signal a preceding proper noun is a person-name candidate.
# Ordered so longer forms match first to keep "統括次長" from being split.
# Title suffix alt list. "様" is excluded on purpose: "仕様" / "様式" / "様態"
# dominate engineering docs and produce too many false positives. "氏" stays
# but is paired with a negative-lookahead guard below so it does not match
# "氏名" / "氏族" / etc.
_PERSON_TITLE_SUFFIXES = (
    "統括次長",
    "総括次長",
    "副部長",
    "副社長",
    "役員",
    "部長",
    "次長",
    "課長",
    "係長",
    "主任",
    "リーダー",
    "マネージャー",
    "ディレクター",
    "フェロー",
    "顧問",
    "参与",
    "氏(?![一-龯々〆ヶ])",
    "さん",
)

_KANJI_CLASS = "一-龯々〆ヶ"
_KATAKANA_CLASS = "ァ-ヴー"
_TITLE_ALT = "|".join(_PERSON_TITLE_SUFFIXES)

# Non-greedy so "中瀬統括次長" yields "中瀬" instead of "中瀬統括".
_PERSON_KANJI_SURNAME_RE = re.compile(
    rf"(?P<name>[{_KANJI_CLASS}]{{1,4}}?)(?={_TITLE_ALT})"
)
_PERSON_KATAKANA_SURNAME_RE = re.compile(
    rf"(?P<name>[{_KATAKANA_CLASS}]{{2,6}}?)(?={_TITLE_ALT})"
)

# Company suffix patterns, longest first so "不動産レジリース" matches before "不動産".
_COMPANY_SUFFIXES = (
    "ホールディングス",
    "コーポレーション",
    "リアルエステートサービス",
    "リアルエステート",
    "不動産レジリース",
    "住宅リース",
    "レジリース",
    "工務店",
    "建設",
    "建物",
    "証券",
    "證券",
    "商事",
    "物産",
    "銀行",
    "ハウス",
    "リビング",
    "リース",
    "不動産",
)
_COMPANY_RE = re.compile(
    rf"(?P<company>[{_KANJI_CLASS}{_KATAKANA_CLASS}・]{{2,16}}(?:{'|'.join(_COMPANY_SUFFIXES)}))"
)

# Property codes used in pptx: e.g. "BS北畠", "BF王子神谷".
_PROPERTY_CODE_RE = re.compile(rf"(?P<code>B[SF][{_KANJI_CLASS}]{{1,6}})")

# Discard candidates that are clearly generic or that match the title itself.
_DENYLIST_EXACT = frozenset({
    "部長",
    "次長",
    "課長",
    "係長",
    "氏",
    "様",
    "さん",
    "株式会社",
    "本日",
    "今回",
    "今後",
    "今週",
    "来週",
    "同",  # pronoun "同〜" (e.g., 同部長) false-positive
    "各",
    "貴",  # 貴社
    "当",
    "弊",
    "統括",  # "統括次長" stem — not a surname
    "総括",  # "総括次長" stem — not a surname
    "副",
    "役員",
    "リーダー",
    "プロジェクト",
    "メンバー",
    "管理",
    "業務",
    "運用",
    "営業",
    "事業",
    "企画",
    "経営",
    "製造",
    "研究",
    "開発",
    "調達",
    "物流",
})


@dataclass(frozen=True)
class CandidateHit:
    """One location a candidate string was observed.

    Intentionally excludes raw surrounding text by default. If ``context`` is
    non-empty it is because the caller opted in via ``with_context=True``.
    """

    slide_number: int
    shape_path: str
    paragraph_index: int
    match_start: int
    match_end: int
    context: str = ""


@dataclass
class Candidate:
    text: str
    category: str  # "person_name" | "company_name" | "property_code"
    hits: list[CandidateHit] = field(default_factory=list)

    @property
    def occurrences(self) -> int:
        return len(self.hits)


def scan_pptx(path: Path, *, with_context: bool = False) -> list[Candidate]:
    """Return consolidated Candidate records for a pptx file.

    The scan intentionally over-reports and relies on the confirmation UI to
    reject noise. Raw document text is read but only minimal location data is
    attached to each Candidate unless ``with_context`` is True.
    """
    from pptx import Presentation  # local import: optional dep

    prs = Presentation(str(path))
    by_text: dict[tuple[str, str], Candidate] = {}

    for slide_idx, slide in enumerate(prs.slides, start=1):
        for shape in slide.shapes:
            for shape_path, para_idx, text in _walk_shape(shape, prefix=""):
                for category, match in _scan_paragraph_text(text):
                    key = (category, match["text"])
                    candidate = by_text.setdefault(
                        key, Candidate(text=match["text"], category=category)
                    )
                    candidate.hits.append(
                        CandidateHit(
                            slide_number=slide_idx,
                            shape_path=shape_path,
                            paragraph_index=para_idx,
                            match_start=match["start"],
                            match_end=match["end"],
                            context=text if with_context else "",
                        )
                    )

    return sorted(
        by_text.values(),
        key=lambda c: (c.category, -c.occurrences, c.text),
    )


def _walk_shape(shape, *, prefix: str) -> Iterator[tuple[str, int, str]]:
    """Yield (shape_path, paragraph_index, text) tuples for a shape and children."""
    if getattr(shape, "shape_type", None) == 6:  # GROUP
        for idx, child in enumerate(shape.shapes):
            yield from _walk_shape(child, prefix=f"{prefix}g{idx}/")
        return
    shape_path = f"{prefix}{getattr(shape, 'shape_id', '?')}"
    if getattr(shape, "has_text_frame", False):
        for para_idx, para in enumerate(shape.text_frame.paragraphs):
            text = "".join(r.text for r in para.runs)
            if text.strip():
                yield (shape_path, para_idx, text)
    if getattr(shape, "has_table", False):
        for row_idx, row in enumerate(shape.table.rows):
            for col_idx, cell in enumerate(row.cells):
                for para_idx, para in enumerate(cell.text_frame.paragraphs):
                    text = "".join(r.text for r in para.runs)
                    if text.strip():
                        yield (
                            f"{shape_path}/r{row_idx}c{col_idx}",
                            para_idx,
                            text,
                        )


def _scan_paragraph_text(text: str) -> Iterable[tuple[str, dict]]:
    """Yield (category, match_dict) tuples for each heuristic hit in ``text``."""
    for regex, category, group_name in (
        (_PERSON_KANJI_SURNAME_RE, "person_name", "name"),
        (_PERSON_KATAKANA_SURNAME_RE, "person_name", "name"),
        (_COMPANY_RE, "company_name", "company"),
        (_PROPERTY_CODE_RE, "property_code", "code"),
    ):
        for match in regex.finditer(text):
            value = match.group(group_name)
            if not value or value in _DENYLIST_EXACT:
                continue
            yield category, {
                "text": value,
                "start": match.start(group_name),
                "end": match.start(group_name) + len(value),
            }


def summarize_counts(candidates: list[Candidate]) -> dict[str, int]:
    """Convenience for logging: {'person_name': n, 'company_name': n, ...}."""
    counter: Counter = Counter(cand.category for cand in candidates)
    return dict(counter)
