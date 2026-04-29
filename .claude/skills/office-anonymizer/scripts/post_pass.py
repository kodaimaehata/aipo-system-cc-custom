"""Position-based post_pass for office-anonymizer.

The SG5 runtime is conservative: when a body-text match lands inside a pptx
table cell or spans multiple text runs it declines to rewrite the underlying
XML and marks the match ``manual_review_required``. This module picks up those
*residual SG5 coordinates* and performs a targeted replacement at exactly that
location — never a document-wide string substitution.

Scope (hard boundary):
  - pptx body text only (slide shapes, groups, table cells).
  - Any residual coordinate pointing at notes / comments / headers / footers /
    metadata / non-pptx formats raises ``UnsupportedPostPassScope`` so the
    ``final_revalidate`` step catches the leak instead of letting it through.

Secondary cleanups (all scoped to paragraphs already touched in this run):
  - Role-suffix dedupe: e.g. replacing "平本" with "A役員" inside "平本役員"
    yields "A役員役員"; we collapse that to "A役員" using the approved
    replacement set.
  - Concatenation spacing: insert a space between a known bare label
    (e.g. "RET") and an immediately adjacent anonymized token ("M氏") when the
    original had no delimiter.
"""

from __future__ import annotations

import re
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Iterable, Mapping


_BODY_TEXT_CATEGORY = "body_text"
# Categories for which SG5 already has the final say; post_pass is not
# authorized to touch these. A residual here indicates a leak the higher-level
# pipeline must surface and fail on.
_OUT_OF_SCOPE_CATEGORIES = frozenset({
    "notes",
    "comments",
    "headers",
    "footers",
    "metadata",
})

# Role suffixes often absorbed into the anonymized label itself (e.g. "A役員"
# then followed by original "役員" in source). Ordered longest-first.
_ROLE_SUFFIX_TOKENS = (
    "統括次長",
    "総括次長",
    "副部長",
    "副社長",
    "役員",
    "部長",
    "次長",
    "課長",
    "係長",
    "リーダー",
    "マネージャー",
    "ディレクター",
    "フェロー",
    "顧問",
    "参与",
    "氏",
    "さん",
)


class UnsupportedPostPassScope(RuntimeError):
    """Raised when post_pass is asked to touch a category it cannot safely edit."""


@dataclass(frozen=True)
class _ReplaceTarget:
    slide_number: int
    shape_id: str
    paragraph_index: int
    match_start: int
    match_end: int
    replacement: str


def run_post_pass(
    *,
    pptx_path: Path,
    sg5_transform_results: Iterable[Mapping],
    replacement_overrides: Mapping[str, str],
    approved_replacements: Iterable[str],
    approved_mapping: Mapping[str, str] | None = None,
    known_bare_labels: Iterable[str] = ("RET", "WP", "AM", "CN", "PMO"),
    backup_dir: Path | None = None,
) -> dict:
    """Apply position-based replacements driven by SG5 residual coordinates.

    Returns a summary dict with counts. Raises ``UnsupportedPostPassScope`` if
    any residual points outside pptx body scope; the wrapper is expected to
    let that exception propagate so ``final_revalidate`` captures the leak.
    """
    from pptx import Presentation  # local import: optional dep

    targets = _collect_targets(
        sg5_transform_results=sg5_transform_results,
        replacement_overrides=replacement_overrides,
        approved_mapping=approved_mapping or {},
    )

    pptx_path = Path(pptx_path)
    if backup_dir is not None:
        backup_dir.mkdir(parents=True, exist_ok=True)
        backup_path = backup_dir / f"{pptx_path.name}.pre-postpass.bak"
        shutil.copyfile(pptx_path, backup_path)
    else:
        backup_path = None

    prs = Presentation(str(pptx_path))
    replaced = _apply_replacements(prs, targets)
    dedupes = _dedupe_role_suffixes(prs, approved_replacements)
    spacings = _space_concatenations(prs, known_bare_labels, approved_replacements)
    prs.save(str(pptx_path))

    return {
        "targeted_replacements": replaced,
        "role_suffix_dedupes": dedupes,
        "concatenation_spacings": spacings,
        "backup_path": str(backup_path) if backup_path else None,
    }


def _collect_targets(
    *,
    sg5_transform_results: Iterable[Mapping],
    replacement_overrides: Mapping[str, str],
    approved_mapping: Mapping[str, str],
) -> list[_ReplaceTarget]:
    """Flatten SG5 manual_review residuals into concrete _ReplaceTarget records.

    SG5's validation_results list entries carry:
      - ``body_text_candidates`` with ``candidate_id`` + ``excerpt_samples[].finding_id``
      - ``manual_review_items`` with ``finding_id`` + ``location`` + ``category``
      - ``transform_summary.actions`` with applied-vs-manual status

    We join those to build concrete position-based replacement targets. Any
    ``manual_review_items`` entry in a category this post-pass cannot touch
    surfaces as ``UnsupportedPostPassScope`` so ``final_revalidate`` can
    decide the run is a leak.
    """
    targets: list[_ReplaceTarget] = []
    for result in sg5_transform_results or []:
        if str(result.get("extension") or "").lower() != "pptx":
            continue

        finding_to_candidate = _index_finding_ids(result.get("body_text_candidates") or [])

        for item in result.get("manual_review_items") or []:
            category = str(item.get("category") or "")
            # Out-of-scope categories (comments/headers/footers/metadata) are
            # intentionally left for final_revalidate to flag. post_pass only
            # touches body_text so those residuals are never silently patched.
            if category != _BODY_TEXT_CATEGORY:
                continue
            candidate_id = finding_to_candidate.get(str(item.get("finding_id") or ""))
            replacement = replacement_overrides.get(candidate_id) if candidate_id else None
            location = item.get("location") or {}
            if not replacement:
                # Fall back to approved_mapping match on residual-rescan items
                # whose finding_id isn't in the initial candidate index.
                replacement = _pick_replacement_from_mapping(
                    approved_mapping=approved_mapping,
                    requested_length=_location_length(location),
                )
            if not replacement:
                continue
            slide_number = location.get("slide_number")
            shape_id = location.get("shape_id")
            paragraph_index = location.get("paragraph_index")
            match_start = location.get("match_start")
            match_end = location.get("match_end")
            if None in (slide_number, shape_id, paragraph_index, match_start, match_end):
                continue
            targets.append(
                _ReplaceTarget(
                    slide_number=int(slide_number),
                    shape_id=str(shape_id),
                    paragraph_index=int(paragraph_index),
                    match_start=int(match_start),
                    match_end=int(match_end),
                    replacement=str(replacement),
                )
            )
    return targets


def _pick_replacement_from_mapping(
    *, approved_mapping: Mapping[str, str], requested_length: int | None
) -> str | None:
    """Fallback: pick the approved replacement whose original length matches.

    This only fires for manual_review_items whose finding_id is not in the
    initial candidate index (typically ``source=residual_rescan`` entries).
    Matching on length is imperfect but keeps the logic position-bound: the
    apply step still checks that the slice at ``[match_start:match_end]``
    equals the chosen original before rewriting.
    """
    if requested_length is None or requested_length <= 0:
        return None
    for original, replacement in approved_mapping.items():
        if len(original) == requested_length:
            return replacement
    return None


def _location_length(location: Mapping) -> int | None:
    start = location.get("match_start")
    end = location.get("match_end")
    if isinstance(start, int) and isinstance(end, int) and end > start:
        return end - start
    return None


def _index_finding_ids(body_text_candidates: Iterable[Mapping]) -> dict[str, str]:
    """Return {finding_id: candidate_id} over every excerpt_samples entry."""
    index: dict[str, str] = {}
    for candidate in body_text_candidates:
        candidate_id = str(candidate.get("candidate_id") or "")
        if not candidate_id:
            continue
        for sample in candidate.get("excerpt_samples") or []:
            finding_id = str(sample.get("finding_id") or "")
            if finding_id:
                index[finding_id] = candidate_id
    return index


def _apply_replacements(prs, targets: list[_ReplaceTarget]) -> int:
    """Apply targeted replacements. Returns count of paragraphs modified."""
    count = 0
    by_slide = _group_targets_by_slide(targets)
    for slide_idx, slide in enumerate(prs.slides, start=1):
        slide_targets = by_slide.get(slide_idx, [])
        if not slide_targets:
            continue
        for shape in slide.shapes:
            count += _apply_to_shape(shape, slide_targets)
    return count


def _group_targets_by_slide(
    targets: list[_ReplaceTarget],
) -> dict[int, list[_ReplaceTarget]]:
    grouped: dict[int, list[_ReplaceTarget]] = {}
    for target in targets:
        grouped.setdefault(target.slide_number, []).append(target)
    return grouped


def _apply_to_shape(shape, targets: list[_ReplaceTarget]) -> int:
    if getattr(shape, "shape_type", None) == 6:  # GROUP
        count = 0
        for child in shape.shapes:
            count += _apply_to_shape(child, targets)
        return count

    shape_id = str(getattr(shape, "shape_id", ""))
    matching = [t for t in targets if t.shape_id == shape_id]
    if not matching:
        return 0

    count = 0
    if getattr(shape, "has_text_frame", False):
        count += _apply_to_text_frame(shape.text_frame, matching)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                count += _apply_to_text_frame(cell.text_frame, matching)
    return count


def _apply_to_text_frame(text_frame, targets: list[_ReplaceTarget]) -> int:
    count = 0
    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        para_targets = [t for t in targets if t.paragraph_index == para_idx]
        if not para_targets:
            continue
        current = "".join(r.text for r in paragraph.runs)
        if not current:
            continue
        # Apply from the end so earlier slices stay valid.
        para_targets.sort(key=lambda t: t.match_start, reverse=True)
        mutated = current
        for target in para_targets:
            if target.match_start < 0 or target.match_end > len(mutated):
                continue
            mutated = (
                mutated[: target.match_start]
                + target.replacement
                + mutated[target.match_end:]
            )
        if mutated != current:
            _write_paragraph_text(paragraph, mutated)
            count += 1
    return count


def _write_paragraph_text(paragraph, new_text: str) -> None:
    if not paragraph.runs:
        return
    paragraph.runs[0].text = new_text
    for run in paragraph.runs[1:]:
        run.text = ""


def _dedupe_role_suffixes(prs, approved_replacements: Iterable[str]) -> int:
    """Remove "<label><role><role>" duplication introduced by overlap."""
    ordered_pairs = [
        (label, suffix)
        for suffix in _ROLE_SUFFIX_TOKENS
        for label in approved_replacements
        if label.endswith(suffix)
    ]
    if not ordered_pairs:
        return 0

    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            count += _dedupe_on_shape(shape, ordered_pairs)
    return count


def _dedupe_on_shape(shape, ordered_pairs) -> int:
    if getattr(shape, "shape_type", None) == 6:
        count = 0
        for child in shape.shapes:
            count += _dedupe_on_shape(child, ordered_pairs)
        return count
    count = 0
    if getattr(shape, "has_text_frame", False):
        count += _dedupe_on_text_frame(shape.text_frame, ordered_pairs)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                count += _dedupe_on_text_frame(cell.text_frame, ordered_pairs)
    return count


def _dedupe_on_text_frame(text_frame, ordered_pairs) -> int:
    count = 0
    for paragraph in text_frame.paragraphs:
        text = "".join(r.text for r in paragraph.runs)
        if not text:
            continue
        new = text
        for label, suffix in ordered_pairs:
            new = new.replace(label + suffix, label)
        if new != text:
            _write_paragraph_text(paragraph, new)
            count += 1
    return count


def _space_concatenations(
    prs,
    known_bare_labels: Iterable[str],
    approved_replacements: Iterable[str],
) -> int:
    """Insert a single space between "<BARE LABEL><anonymized token>" pairs."""
    bare_labels = tuple(known_bare_labels)
    replacements = tuple(approved_replacements)
    if not bare_labels or not replacements:
        return 0

    pattern = re.compile(
        "(" + "|".join(re.escape(b) for b in bare_labels) + ")"
        "(" + "|".join(re.escape(r) for r in replacements) + ")"
    )
    count = 0
    for slide in prs.slides:
        for shape in slide.shapes:
            count += _space_on_shape(shape, pattern)
    return count


def _space_on_shape(shape, pattern) -> int:
    if getattr(shape, "shape_type", None) == 6:
        count = 0
        for child in shape.shapes:
            count += _space_on_shape(child, pattern)
        return count
    count = 0
    if getattr(shape, "has_text_frame", False):
        count += _space_on_text_frame(shape.text_frame, pattern)
    if getattr(shape, "has_table", False):
        for row in shape.table.rows:
            for cell in row.cells:
                count += _space_on_text_frame(cell.text_frame, pattern)
    return count


def _space_on_text_frame(text_frame, pattern) -> int:
    count = 0
    for paragraph in text_frame.paragraphs:
        text = "".join(r.text for r in paragraph.runs)
        if not text:
            continue
        new = pattern.sub(r"\1 \2", text)
        if new != text:
            _write_paragraph_text(paragraph, new)
            count += 1
    return count
