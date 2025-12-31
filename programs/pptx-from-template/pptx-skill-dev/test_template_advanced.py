#!/usr/bin/env python3
"""python-pptx テンプレート操作の詳細検証"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import os
import json

OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def inspect_presentation(pptx_path: str):
    """プレゼンテーションの構造を調査"""
    print(f"\n=== プレゼンテーション調査: {pptx_path} ===")

    prs = Presentation(pptx_path)

    print(f"\nスライド数: {len(prs.slides)}")

    for slide_idx, slide in enumerate(prs.slides):
        print(f"\n--- スライド {slide_idx + 1} ---")
        print(f"レイアウト: {slide.slide_layout.name}")

        for shape in slide.shapes:
            print(f"\n  Shape: {shape.shape_type}, name={shape.name}")

            # プレースホルダーの場合
            if shape.is_placeholder:
                ph = shape.placeholder_format
                print(f"    プレースホルダー: idx={ph.idx}, type={ph.type}")

            # テキストフレームがある場合
            if hasattr(shape, "text_frame"):
                text = shape.text_frame.text[:50] if shape.text_frame.text else "(空)"
                print(f"    テキスト: {text}")


def test_template_with_placeholders():
    """プレースホルダー付きテンプレートのテスト"""
    print("\n=== テンプレート＋プレースホルダーテスト ===")

    # テンプレート作成
    prs = Presentation()

    # スライド1: タイトル
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "{{TITLE}}"
    slide1.placeholders[1].text = "{{SUBTITLE}}"

    # スライド2: コンテンツ
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "{{SECTION_TITLE}}"
    slide2.placeholders[1].text_frame.text = "{{CONTENT}}"

    template_path = f"{OUTPUT_DIR}/advanced_template.pptx"
    prs.save(template_path)
    print(f"  テンプレート作成: {template_path}")

    # テンプレートを読み込んで置換
    prs = Presentation(template_path)

    replacements = {
        "{{TITLE}}": "製品紹介プレゼンテーション",
        "{{SUBTITLE}}": "2025年度版 - 株式会社サンプル",
        "{{SECTION_TITLE}}": "主な機能",
        "{{CONTENT}}": "• 高速処理\n• 日本語完全対応\n• クラウド連携",
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        for key, value in replacements.items():
                            if key in run.text:
                                run.text = run.text.replace(key, value)
                    # 段落全体もチェック（runsが空の場合）
                    if not paragraph.runs:
                        for key, value in replacements.items():
                            if key in paragraph.text:
                                paragraph.text = paragraph.text.replace(key, value)

    output_path = f"{OUTPUT_DIR}/advanced_output.pptx"
    prs.save(output_path)
    print(f"  ✓ 出力保存: {output_path}")

    return True


def test_json_data_driven():
    """JSONデータからプレゼンテーション生成"""
    print("\n=== JSONデータ駆動テスト ===")

    # サンプルJSONデータ
    data = {
        "title": "月次レポート",
        "subtitle": "2025年1月",
        "author": "営業部",
        "slides": [
            {
                "type": "content",
                "title": "売上サマリー",
                "content": ["売上高: 1,000万円", "前年比: +15%", "目標達成率: 102%"],
            },
            {
                "type": "content",
                "title": "主要トピック",
                "content": ["新規顧客獲得: 5社", "リピート率向上", "製品Aが好調"],
            },
        ],
    }

    # JSONファイルとして保存
    json_path = f"{OUTPUT_DIR}/sample_data.json"
    with open(json_path, "w", encoding="utf-8") as f:
        json.dump(data, f, ensure_ascii=False, indent=2)
    print(f"  サンプルデータ: {json_path}")

    # プレゼンテーション生成
    prs = Presentation()

    # タイトルスライド
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = data["title"]
    slide.placeholders[1].text = f"{data['subtitle']}\n{data['author']}"

    # コンテンツスライド
    for slide_data in data["slides"]:
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        slide.shapes.title.text = slide_data["title"]

        body = slide.placeholders[1]
        tf = body.text_frame
        tf.text = slide_data["content"][0]

        for item in slide_data["content"][1:]:
            p = tf.add_paragraph()
            p.text = item
            p.level = 0

    output_path = f"{OUTPUT_DIR}/json_driven_output.pptx"
    prs.save(output_path)
    print(f"  ✓ 出力保存: {output_path}")

    return True


def test_preserve_formatting():
    """フォーマット保持テスト"""
    print("\n=== フォーマット保持テスト ===")

    from pptx.dml.color import RGBColor

    # カスタムフォーマットのテンプレート作成
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白

    # フォーマット付きテキストボックス
    txBox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(2))
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = "{{HEADING}}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x00, 0x80)  # Navy

    p2 = tf.add_paragraph()
    p2.text = "{{BODY}}"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(0x33, 0x33, 0x33)  # Dark gray

    template_path = f"{OUTPUT_DIR}/formatted_template.pptx"
    prs.save(template_path)
    print(f"  テンプレート作成: {template_path}")

    # テンプレートを読み込んで置換（フォーマット保持）
    prs = Presentation(template_path)

    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text_frame"):
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # フォーマットを保持したまま置換
                        if "{{HEADING}}" in run.text:
                            run.text = "重要なお知らせ"
                        elif "{{BODY}}" in run.text:
                            run.text = "この度、新サービスを開始いたしました。詳細は以下をご確認ください。"

    output_path = f"{OUTPUT_DIR}/formatted_output.pptx"
    prs.save(output_path)
    print(f"  ✓ 出力保存: {output_path}")

    # 検証
    inspect_presentation(output_path)

    return True


def main():
    print("=" * 60)
    print("python-pptx テンプレート操作 詳細検証")
    print("=" * 60)

    results = []

    tests = [
        test_template_with_placeholders,
        test_json_data_driven,
        test_preserve_formatting,
    ]

    for test in tests:
        try:
            result = test()
            results.append((test.__name__, "PASS" if result else "FAIL"))
        except Exception as e:
            import traceback

            print(f"  ✗ エラー: {e}")
            traceback.print_exc()
            results.append((test.__name__, f"ERROR: {e}"))

    print("\n" + "=" * 60)
    print("テスト結果サマリー")
    print("=" * 60)
    for name, status in results:
        print(f"  {name}: {status}")


if __name__ == "__main__":
    main()
