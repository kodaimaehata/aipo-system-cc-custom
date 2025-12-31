#!/usr/bin/env python3
"""python-pptx 基本動作確認スクリプト"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
import os

OUTPUT_DIR = "output"
os.makedirs(OUTPUT_DIR, exist_ok=True)


def test_1_create_new_presentation():
    """テスト1: 新規プレゼンテーション作成"""
    print("\n=== テスト1: 新規プレゼンテーション作成 ===")

    prs = Presentation()

    # タイトルスライド追加
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    title = slide.shapes.title
    subtitle = slide.placeholders[1]

    title.text = "テストプレゼンテーション"
    subtitle.text = "python-pptxによる自動生成"

    output_path = f"{OUTPUT_DIR}/test1_new_presentation.pptx"
    prs.save(output_path)
    print(f"  ✓ 保存完了: {output_path}")
    return True


def test_2_multiple_slide_types():
    """テスト2: 複数スライドタイプの作成"""
    print("\n=== テスト2: 複数スライドタイプ ===")

    prs = Presentation()

    # スライドレイアウト一覧を表示
    print("  利用可能なレイアウト:")
    for i, layout in enumerate(prs.slide_layouts):
        print(f"    [{i}] {layout.name}")

    # タイトルスライド (レイアウト0)
    slide1 = prs.slides.add_slide(prs.slide_layouts[0])
    slide1.shapes.title.text = "メインタイトル"
    slide1.placeholders[1].text = "サブタイトル - 日本語テスト"

    # タイトルと箇条書き (レイアウト1)
    slide2 = prs.slides.add_slide(prs.slide_layouts[1])
    slide2.shapes.title.text = "箇条書きスライド"
    body = slide2.placeholders[1]
    tf = body.text_frame
    tf.text = "第一項目"
    p = tf.add_paragraph()
    p.text = "第二項目"
    p.level = 0
    p = tf.add_paragraph()
    p.text = "サブ項目"
    p.level = 1

    # 空白スライド (レイアウト6)
    slide3 = prs.slides.add_slide(prs.slide_layouts[6])

    output_path = f"{OUTPUT_DIR}/test2_multiple_slides.pptx"
    prs.save(output_path)
    print(f"  ✓ 保存完了: {output_path}")
    return True


def test_3_placeholder_inspection():
    """テスト3: プレースホルダーの調査"""
    print("\n=== テスト3: プレースホルダー調査 ===")

    prs = Presentation()

    print("  各レイアウトのプレースホルダー:")
    for i, layout in enumerate(prs.slide_layouts[:3]):  # 最初の3つだけ
        print(f"\n  レイアウト[{i}] {layout.name}:")
        for shape in layout.placeholders:
            print(f"    - idx={shape.placeholder_format.idx}, type={shape.placeholder_format.type}, name={shape.name}")

    return True


def test_4_table_creation():
    """テスト4: 表の作成"""
    print("\n=== テスト4: 表の作成 ===")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # タイトルのみ
    slide.shapes.title.text = "表のテスト"

    # 表を追加
    rows, cols = 3, 4
    left = Inches(1)
    top = Inches(2)
    width = Inches(8)
    height = Inches(2)

    table = slide.shapes.add_table(rows, cols, left, top, width, height).table

    # ヘッダー行
    headers = ["名前", "部署", "役職", "入社年"]
    for i, header in enumerate(headers):
        table.cell(0, i).text = header

    # データ行
    data = [
        ["山田太郎", "開発部", "エンジニア", "2020"],
        ["鈴木花子", "営業部", "マネージャー", "2018"],
    ]
    for row_idx, row_data in enumerate(data, start=1):
        for col_idx, value in enumerate(row_data):
            table.cell(row_idx, col_idx).text = value

    output_path = f"{OUTPUT_DIR}/test4_table.pptx"
    prs.save(output_path)
    print(f"  ✓ 保存完了: {output_path}")
    return True


def test_5_text_formatting():
    """テスト5: テキストフォーマット"""
    print("\n=== テスト5: テキストフォーマット ===")

    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白

    # テキストボックス追加
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    left = Inches(1)
    top = Inches(1)
    width = Inches(8)
    height = Inches(1)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame

    p = tf.paragraphs[0]
    p.text = "フォーマットテスト"
    p.font.size = Pt(32)
    p.font.bold = True
    p.font.color.rgb = RGBColor(0x00, 0x66, 0xCC)
    p.alignment = PP_ALIGN.CENTER

    # 追加の段落
    p2 = tf.add_paragraph()
    p2.text = "日本語テキストのフォーマット確認"
    p2.font.size = Pt(18)

    output_path = f"{OUTPUT_DIR}/test5_formatting.pptx"
    prs.save(output_path)
    print(f"  ✓ 保存完了: {output_path}")
    return True


def test_6_template_usage():
    """テスト6: テンプレートとしての使用（既存ファイル読み込み）"""
    print("\n=== テスト6: テンプレート使用 ===")

    # まずテンプレートを作成
    template = Presentation()
    slide = template.slides.add_slide(template.slide_layouts[0])
    slide.shapes.title.text = "{{title}}"
    slide.placeholders[1].text = "{{subtitle}}"

    template_path = f"{OUTPUT_DIR}/template.pptx"
    template.save(template_path)
    print(f"  テンプレート作成: {template_path}")

    # テンプレートを読み込んで編集
    prs = Presentation(template_path)
    slide = prs.slides[0]

    # プレースホルダーテキストを置換
    for shape in slide.shapes:
        if hasattr(shape, "text"):
            if "{{title}}" in shape.text:
                shape.text = "置換されたタイトル"
            elif "{{subtitle}}" in shape.text:
                shape.text = "置換されたサブタイトル"

    output_path = f"{OUTPUT_DIR}/test6_from_template.pptx"
    prs.save(output_path)
    print(f"  ✓ 保存完了: {output_path}")
    return True


def main():
    print("=" * 50)
    print("python-pptx 基本動作確認")
    print("=" * 50)

    results = []

    tests = [
        test_1_create_new_presentation,
        test_2_multiple_slide_types,
        test_3_placeholder_inspection,
        test_4_table_creation,
        test_5_text_formatting,
        test_6_template_usage,
    ]

    for test in tests:
        try:
            result = test()
            results.append((test.__name__, "PASS" if result else "FAIL"))
        except Exception as e:
            print(f"  ✗ エラー: {e}")
            results.append((test.__name__, f"ERROR: {e}"))

    print("\n" + "=" * 50)
    print("テスト結果サマリー")
    print("=" * 50)
    for name, status in results:
        print(f"  {name}: {status}")

    print(f"\n生成ファイル: {OUTPUT_DIR}/ ディレクトリを確認してください")


if __name__ == "__main__":
    main()
