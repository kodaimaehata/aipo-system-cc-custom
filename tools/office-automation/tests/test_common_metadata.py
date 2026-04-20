from __future__ import annotations

from datetime import datetime, timezone
from pathlib import Path

import fitz
import pytest
from docx import Document
from openpyxl import Workbook, load_workbook
from pptx import Presentation

from office_automation.common import metadata as metadata_module
from office_automation.common.metadata import clear_metadata, read_metadata


_SAMPLE_DATETIME = datetime(2024, 1, 2, 3, 4, 5, tzinfo=timezone.utc)


def _create_excel_workbook(path: Path) -> Path:
    workbook = Workbook()
    properties = workbook.properties
    properties.title = "Workbook Title"
    properties.subject = "Workbook Subject"
    properties.creator = "Workbook Creator"
    properties.keywords = "alpha,beta"
    properties.description = "Workbook Description"
    properties.language = "en-US"
    properties.lastModifiedBy = "Workbook Modifier"
    properties.category = "Workbook Category"
    properties.contentStatus = "Workbook Status"
    properties.identifier = "xlsx-id"
    properties.version = "1.0"
    properties.revision = "7"
    properties.created = _SAMPLE_DATETIME
    properties.modified = _SAMPLE_DATETIME
    workbook.active["A1"] = "hello"
    workbook.save(path)
    return path


def _create_word_document(path: Path) -> Path:
    document = Document()
    document.add_paragraph("Metadata sample")
    properties = document.core_properties
    properties.title = "Document Title"
    properties.subject = "Document Subject"
    properties.author = "Document Creator"
    properties.keywords = "gamma,delta"
    properties.comments = "Document Description"
    properties.language = "ja-JP"
    properties.last_modified_by = "Document Modifier"
    properties.category = "Document Category"
    properties.content_status = "Document Status"
    properties.identifier = "docx-id"
    properties.version = "2.0"
    properties.revision = 9
    properties.created = _SAMPLE_DATETIME
    properties.modified = _SAMPLE_DATETIME
    document.save(path)
    return path


def _create_powerpoint(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Metadata sample"
    properties = presentation.core_properties
    properties.title = "Presentation Title"
    properties.subject = "Presentation Subject"
    properties.author = "Presentation Creator"
    properties.keywords = "epsilon,zeta"
    properties.comments = "Presentation Description"
    properties.language = "fr-FR"
    properties.last_modified_by = "Presentation Modifier"
    properties.category = "Presentation Category"
    properties.content_status = "Presentation Status"
    properties.identifier = "pptx-id"
    properties.version = "3.0"
    properties.revision = 11
    properties.created = _SAMPLE_DATETIME
    properties.modified = _SAMPLE_DATETIME
    presentation.save(path)
    return path


def _create_pdf(path: Path) -> Path:
    document = fitz.open()
    document.new_page()
    document.set_metadata(
        {
            "title": "PDF Title",
            "author": "PDF Author",
            "subject": "PDF Subject",
            "keywords": "eta,theta",
            "creator": "PDF Creator",
            "producer": "PDF Producer",
        }
    )
    document.save(path)
    document.close()
    return path


@pytest.mark.parametrize(
    ("builder", "extension", "expectations"),
    [
        (
            _create_excel_workbook,
            "xlsx",
            {
                "title": "Workbook Title",
                "creator": "Workbook Creator",
                "description": "Workbook Description",
            },
        ),
        (
            _create_word_document,
            "docx",
            {
                "title": "Document Title",
                "creator": "Document Creator",
                "description": "Document Description",
            },
        ),
        (
            _create_powerpoint,
            "pptx",
            {
                "title": "Presentation Title",
                "creator": "Presentation Creator",
                "description": "Presentation Description",
            },
        ),
        (
            _create_pdf,
            "pdf",
            {
                "title": "PDF Title",
                "author": "PDF Author",
                "creator": "PDF Creator",
            },
        ),
    ],
)
def test_read_metadata_returns_normalized_fields(
    tmp_path: Path,
    builder,
    extension: str,
    expectations: dict[str, str],
) -> None:
    source = builder(tmp_path / f"sample.{extension}")

    result = read_metadata(source)

    assert result["file_path"] == str(source)
    assert result["extension"] == extension
    assert result["warnings"] == []
    for field_name, expected in expectations.items():
        assert result["fields"][field_name] == expected


@pytest.mark.parametrize(
    ("builder", "extension"),
    [
        (_create_excel_workbook, "xlsx"),
        (_create_word_document, "docx"),
        (_create_powerpoint, "pptx"),
        (_create_pdf, "pdf"),
    ],
)
def test_clear_metadata_blanks_supported_fields_in_place(
    tmp_path: Path,
    builder,
    extension: str,
) -> None:
    source = builder(tmp_path / f"sample.{extension}")

    clear_metadata(source)
    result = read_metadata(source)

    assert result["extension"] == extension
    assert all(value is None for value in result["fields"].values())



def test_clear_metadata_for_xlsm_uses_keep_vba_and_preserves_extension(
    tmp_path: Path,
    monkeypatch: pytest.MonkeyPatch,
) -> None:
    source = _create_excel_workbook(tmp_path / "macro.xlsm")
    observed_keep_vba: list[bool | None] = []
    original_load_workbook = metadata_module.load_workbook

    def recording_load_workbook(*args, **kwargs):
        observed_keep_vba.append(kwargs.get("keep_vba"))
        return original_load_workbook(*args, **kwargs)

    monkeypatch.setattr(metadata_module, "load_workbook", recording_load_workbook)

    clear_metadata(source)

    assert True in observed_keep_vba
    assert source.suffix == ".xlsm"

    workbook = original_load_workbook(source, keep_vba=True)
    try:
        assert workbook.vba_archive is not None
    finally:
        vba_archive = getattr(workbook, "vba_archive", None)
        if vba_archive is not None:
            vba_archive.close()
        workbook.close()

    assert all(value is None for value in read_metadata(source)["fields"].values())


@pytest.mark.parametrize("callable_obj", [read_metadata, clear_metadata])
def test_metadata_helpers_validate_missing_non_file_and_unsupported_paths(
    tmp_path: Path,
    callable_obj,
) -> None:
    directory_path = tmp_path / "folder"
    directory_path.mkdir()
    unsupported = tmp_path / "legacy.xls"
    unsupported.write_text("legacy")

    with pytest.raises(FileNotFoundError, match=r"Metadata source file '.*missing\.docx' does not exist\."):
        callable_obj(tmp_path / "missing.docx")

    with pytest.raises(ValueError, match=r"Metadata source file '.*folder' is not a regular file\."):
        callable_obj(directory_path)

    with pytest.raises(
        ValueError,
        match=r"Metadata source file '.*legacy\.xls' has unsupported extension 'xls'\. Supported extensions: docx, pdf, pptx, xlsm, xlsx\.",
    ):
        callable_obj(unsupported)
