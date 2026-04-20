from __future__ import annotations

from pathlib import Path

import fitz
from docx import Document
from openpyxl import Workbook, load_workbook
from openpyxl.comments import Comment
from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from office_automation.anonymize.candidate_summary import (
    build_body_text_candidate_summary,
    resolve_body_text_confirmation,
)
from office_automation.anonymize.detect import detect
from office_automation.anonymize.transform import transform
from office_automation.common.images import extract_images
from office_automation.common.metadata import read_metadata


STRUCTURAL_CATEGORIES = {"comments", "notes", "headers", "footers"}
_IMAGE_SIZE = (40, 40)



def _policy(
    *,
    comments: dict | None = None,
    notes: dict | None = None,
    headers: dict | None = None,
    footers: dict | None = None,
    metadata: dict | None = None,
    images: dict | None = None,
) -> dict:
    return {
        "comments": comments or {"enabled": True, "action": "remove"},
        "notes": notes or {"enabled": True, "action": "remove"},
        "headers": headers or {"enabled": True, "action": "remove"},
        "footers": footers or {"enabled": True, "action": "remove"},
        "metadata": metadata or {"enabled": True, "action": "clear"},
        "images": images or {"enabled": True, "mode": "report_only", "replacement_path": None},
        "manual_review_required": True,
    }



def _findings_for(folder: Path, file_name: str, *, categories: set[str] | None = None) -> list[dict]:
    return [
        finding
        for finding in detect(folder)
        if finding["relative_path"] == file_name and (categories is None or finding["category"] in categories)
    ]



def _structural_findings_for(folder: Path, file_name: str) -> list[dict]:
    return _findings_for(folder, file_name, categories=STRUCTURAL_CATEGORIES)



def _findings_for_with_body_text(
    folder: Path,
    file_name: str,
    *,
    candidate_inputs: dict,
    categories: set[str] | None = None,
) -> list[dict]:
    return [
        finding
        for finding in detect(folder, body_text_candidate_inputs=candidate_inputs)
        if finding["relative_path"] == file_name and (categories is None or finding["category"] in categories)
    ]



def _create_xlsx_with_note_header_footer(path: Path) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "SheetA"
    worksheet["A1"] = "value"
    worksheet["A1"].comment = Comment("Secret note", "Bob")
    worksheet.oddHeader.left.text = "Left Header"
    worksheet.oddFooter.center.text = "Center Footer"
    workbook.save(path)
    workbook.close()
    return path



def _create_docx_with_comment_header_footer(path: Path) -> Path:
    document = Document()
    section = document.sections[0]
    section.header.paragraphs[0].text = "Document Header"
    section.footer.paragraphs[0].text = "Document Footer"
    paragraph = document.add_paragraph()
    run = paragraph.add_run("Hello document")
    document.add_comment(run, text="Remove this comment", author="Alice", initials="AL")
    document.save(path)
    return path



def _create_pptx_with_notes(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Deck Title"
    slide.notes_slide.notes_text_frame.text = "Presenter note"
    presentation.save(path)
    return path



def _create_pdf_with_annotation(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    annotation = page.add_text_annot((72, 72), "Secret comment")
    annotation.set_info(title="Alice", subject="Review")
    document.save(path)
    document.close()
    return path



def _create_color_image(path: Path, color: tuple[int, int, int]) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", _IMAGE_SIZE, color).save(path, format="PNG")
    return path



def _image_color(path: Path) -> tuple[int, int, int]:
    with Image.open(path) as image:
        return image.convert("RGB").getpixel((0, 0))



def _create_docx_with_metadata_and_image(path: Path, image_path: Path) -> Path:
    document = Document()
    document.add_paragraph("Metadata and image sample")
    document.add_picture(str(image_path))
    properties = document.core_properties
    properties.title = "Sensitive Deck"
    properties.author = "Pat Example"
    properties.subject = "Confidential"
    document.save(path)
    return path



def _create_pdf_with_metadata_and_shared_image(path: Path, image_path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    xref = page.insert_image(fitz.Rect(36, 36, 136, 136), filename=str(image_path))
    page.insert_image(fitz.Rect(156, 36, 256, 136), xref=xref)
    document.set_metadata(
        {
            "title": "Shared Image PDF",
            "author": "Pat Example",
            "subject": "Needs anonymization",
        }
    )
    document.save(path)
    document.close()
    return path



def _create_xlsx_with_body_text_and_header(path: Path) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet["A1"] = "Primary contact: Jane Example"
    worksheet.oddHeader.left.text = "Workbook Header"
    workbook.save(path)
    workbook.close()
    return path



def _create_xlsx_with_multiple_body_text_candidates(path: Path) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet["A1"] = "Primary contact: Jane Example"
    worksheet["A2"] = "Vendor: Example Corp"
    worksheet["A3"] = "Backup contact: Bob Example"
    workbook.save(path)
    workbook.close()
    return path



def _create_docx_with_body_text_targets(path: Path) -> Path:
    document = Document()
    document.add_paragraph("Primary contact: Jane Example")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Vendor: Example Corp"
    document.save(path)
    return path



def _create_pptx_with_body_text_targets(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    textbox = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(5), Inches(1))
    textbox.text_frame.text = "Primary contact: Jane Example"
    table_shape = slide.shapes.add_table(1, 1, Inches(1), Inches(2), Inches(5), Inches(1.2))
    table_shape.table.cell(0, 0).text = "Vendor: Example Corp"
    presentation.save(path)
    return path



def _create_pdf_with_visible_body_text(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Contact john@example.com for the update")
    document.save(path)
    document.close()
    return path



def _build_body_text_policy(
    findings: list[dict],
    *,
    approved_texts: list[str],
    rejected_texts: list[str] | None = None,
    default_replacement_text: str = "[BODY TEXT REDACTED]",
    replacement_overrides: dict[str, str] | None = None,
    include_mode: str = "apply_confirmed",
) -> tuple[dict, dict]:
    summary = build_body_text_candidate_summary(findings)
    candidates_by_text = {candidate["display_text"]: candidate for candidate in summary["candidates"]}
    approved_candidate_ids = [candidates_by_text[text]["candidate_id"] for text in approved_texts]
    rejected_candidate_ids = [candidates_by_text[text]["candidate_id"] for text in rejected_texts or []]
    resolved_overrides = {
        candidates_by_text[text]["candidate_id"]: replacement
        for text, replacement in (replacement_overrides or {}).items()
    }
    resolution = resolve_body_text_confirmation(
        summary,
        {
            "mode": include_mode,
            "approved_candidate_ids": approved_candidate_ids,
            "rejected_candidate_ids": rejected_candidate_ids,
            "replacement_overrides": resolved_overrides,
            "review_notes": [],
        },
    )
    return {
        "enabled": True,
        "mode": include_mode,
        "approved_candidate_ids": list(resolution["approved_candidate_ids"]),
        "rejected_candidate_ids": list(resolution["rejected_candidate_ids"]),
        "undecided_candidate_ids": list(resolution["undecided_candidate_ids"]),
        "approved_finding_ids": list(resolution["approved_finding_ids"]),
        "non_transformable_finding_ids": list(resolution["non_transformable_finding_ids"]),
        "candidate_decisions": dict(resolution["candidate_decisions"]),
        "candidate_summary": {
            "candidates": list(summary["candidates"]),
            "finding_to_candidate": dict(summary["finding_to_candidate"]),
        },
        "replacement_text": default_replacement_text,
        "replacement_map": {},
        "replacement_overrides": resolved_overrides,
        "confirmation_warnings": list(resolution["warnings"]),
    }, candidates_by_text



def test_transform_applies_excel_note_header_and_footer_changes(tmp_path: Path) -> None:
    source = _create_xlsx_with_note_header_footer(tmp_path / "sheet.xlsx")
    findings = _structural_findings_for(tmp_path, source.name)

    results = transform(
        findings,
        _policy(
            notes={"enabled": True, "action": "replace", "replacement_text": "[note removed]"},
            headers={"enabled": True, "action": "clear"},
            footers={"enabled": True, "action": "replace", "replacement_text": "[footer removed]"},
        ),
    )

    assert len(results) == 1
    result = results[0]
    assert result["file_path"] == str(source)
    assert result["status"] == "success"
    assert [action["status"] for action in result["actions"]] == ["applied", "applied", "applied"]

    workbook = load_workbook(source)
    try:
        worksheet = workbook["SheetA"]
        assert worksheet["A1"].comment is not None
        assert worksheet["A1"].comment.text == "[note removed]"
        assert worksheet["A1"].comment.author.strip() == ""
        assert worksheet.oddHeader.left.text == ""
        assert worksheet.oddFooter.center.text == "[footer removed]"
    finally:
        workbook.close()

    rescanned = _structural_findings_for(tmp_path, source.name)
    rescanned_categories = [finding["category"] for finding in rescanned]
    assert rescanned_categories == ["notes", "footers"]
    note_finding = next(finding for finding in rescanned if finding["category"] == "notes")
    footer_finding = next(finding for finding in rescanned if finding["category"] == "footers")
    assert note_finding["payload"]["text"] == "[note removed]"
    assert note_finding["payload"]["author"] is None
    assert footer_finding["payload"]["text"] == "[footer removed]"



def test_transform_reports_docx_comment_manual_review_while_updating_headers_and_footers(tmp_path: Path) -> None:
    source = _create_docx_with_comment_header_footer(tmp_path / "commented.docx")
    findings = _structural_findings_for(tmp_path, source.name)

    results = transform(
        findings,
        _policy(
            comments={"enabled": True, "action": "remove"},
            headers={"enabled": True, "action": "replace", "replacement_text": "Sanitized Header"},
            footers={"enabled": True, "action": "clear"},
        ),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "partial_success"
    assert [action["status"] for action in result["actions"]] == ["manual_review_required", "applied", "applied"]
    assert result["manual_review_items"] == [
        {
            "category": "comments",
            "finding_id": findings[0]["finding_id"],
            "location": findings[0]["location"],
            "reason": "DOCX comments require manual review because this runtime does not yet target Word comment XML with a proved-safe deletion path.",
            "requested_action": "remove",
        }
    ]

    document = Document(source)
    section = document.sections[0]
    assert section.header.paragraphs[0].text == "Sanitized Header"
    assert all(paragraph.text == "" for paragraph in section.footer.paragraphs)

    rescanned = _structural_findings_for(tmp_path, source.name)
    rescanned_categories = [finding["category"] for finding in rescanned]
    assert rescanned_categories == ["comments", "headers"]
    comment_finding = next(finding for finding in rescanned if finding["category"] == "comments")
    header_finding = next(finding for finding in rescanned if finding["category"] == "headers")
    assert comment_finding["payload"]["text"] == "Remove this comment"
    assert header_finding["payload"]["text"] == "Sanitized Header"



def test_transform_returns_structured_error_for_replace_without_replacement_text(tmp_path: Path) -> None:
    source = _create_xlsx_with_note_header_footer(tmp_path / "invalid-policy.xlsx")
    findings = _structural_findings_for(tmp_path, source.name)

    results = transform(
        findings,
        _policy(
            notes={"enabled": True, "action": "replace"},
            headers={"enabled": True, "action": "clear"},
            footers={"enabled": True, "action": "clear"},
        ),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "error"
    note_action = next(action for action in result["actions"] if action["category"] == "notes")
    assert note_action["status"] == "error"
    assert note_action["message"] == "Policy action 'replace' for category 'notes' requires replacement text."

    workbook = load_workbook(source)
    try:
        worksheet = workbook["SheetA"]
        assert worksheet["A1"].comment.text == "Secret note"
        assert worksheet.oddHeader.left.text in (None, "")
        assert worksheet.oddFooter.center.text in (None, "")
    finally:
        workbook.close()



def test_transform_applies_pdf_comment_replacement_and_blanks_annotation_metadata(tmp_path: Path) -> None:
    source = _create_pdf_with_annotation(tmp_path / "annotated.pdf")
    findings = [finding for finding in _structural_findings_for(tmp_path, source.name) if finding["category"] == "comments"]

    results = transform(
        findings,
        _policy(comments={"enabled": True, "action": "replace", "replacement_text": "[annotation removed]"}),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "success"
    assert [action["status"] for action in result["actions"]] == ["applied"]

    document = fitz.open(source)
    try:
        page = document[0]
        annotation = next(page.annots())
        assert annotation.info["content"] == "[annotation removed]"
        assert annotation.info["title"].strip() == ""
        assert annotation.info["subject"].strip() == ""
    finally:
        document.close()

    rescanned = [finding for finding in _structural_findings_for(tmp_path, source.name) if finding["category"] == "comments"]
    assert len(rescanned) == 1
    assert rescanned[0]["payload"]["content"] == "[annotation removed]"
    assert rescanned[0]["payload"]["author"] is None
    assert rescanned[0]["payload"]["subject"] is None



def test_transform_applies_powerpoint_notes_and_preserves_manual_review_items_for_other_surfaces(tmp_path: Path) -> None:
    source = _create_pptx_with_notes(tmp_path / "deck.pptx")
    findings = _structural_findings_for(tmp_path, source.name)

    results = transform(findings, _policy(notes={"enabled": True, "action": "clear"}))

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "partial_success"
    assert [action["status"] for action in result["actions"]] == [
        "manual_review_required",
        "applied",
        "manual_review_required",
        "manual_review_required",
    ]
    assert [item["category"] for item in result["manual_review_items"]] == ["comments", "headers", "footers"]

    rescanned = _structural_findings_for(tmp_path, source.name)
    assert [finding["category"] for finding in rescanned] == ["comments", "headers", "footers"]



def test_transform_clears_metadata_and_masks_docx_images_with_python_only_helpers(tmp_path: Path) -> None:
    image_path = _create_color_image(tmp_path / "inputs" / "blue.png", (0, 0, 255))
    source = _create_docx_with_metadata_and_image(tmp_path / "metadata-image.docx", image_path)
    findings = _findings_for(tmp_path, source.name, categories={"metadata", "images"})

    results = transform(
        findings,
        _policy(
            metadata={"enabled": True, "action": "clear"},
            images={"enabled": True, "mode": "mask", "mask_color": "#11aa22"},
        ),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "success"
    assert [action["category"] for action in result["actions"]] == ["metadata", "images"]
    assert [action["status"] for action in result["actions"]] == ["applied", "applied"]

    metadata_action = next(action for action in result["actions"] if action["category"] == "metadata")
    image_action = next(action for action in result["actions"] if action["category"] == "images")
    for field_name, expected_value in {
        "creator": "Pat Example",
        "subject": "Confidential",
        "title": "Sensitive Deck",
    }.items():
        assert metadata_action["details"]["fields_before"][field_name] == expected_value
    assert metadata_action["details"]["fields_after"] == {}
    assert image_action["applied_action"] == "mask"
    assert image_action["details"]["mask_color"] == "#11aa22"

    assert all(value is None for value in read_metadata(source)["fields"].values())

    extracted = extract_images(source, tmp_path / "extracted" / "docx-masked")
    assert len(extracted) == 1
    assert _image_color(extracted[0]) == (17, 170, 34)



def test_transform_downgrades_remove_to_mask_for_supported_image_slots(tmp_path: Path) -> None:
    image_path = _create_color_image(tmp_path / "inputs" / "red.png", (255, 0, 0))
    source = _create_docx_with_metadata_and_image(tmp_path / "remove-image.docx", image_path)
    findings = _findings_for(tmp_path, source.name, categories={"images"})

    results = transform(
        findings,
        _policy(images={"enabled": True, "mode": "remove", "mask_color": (1, 2, 3)}),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "partial_success"
    assert [action["status"] for action in result["actions"]] == ["applied"]
    action = result["actions"][0]
    assert action["requested_action"] == "remove"
    assert action["applied_action"] == "mask"
    assert action["details"]["mask_color"] == "#010203"
    assert action["warnings"] == [
        "Requested image removal was downgraded to masking because this V1 runtime does not claim a universally safe remove path."
    ]

    extracted = extract_images(source, tmp_path / "extracted" / "docx-remove-downgraded")
    assert [_image_color(path) for path in extracted] == [(1, 2, 3)]



def test_transform_returns_partial_success_when_metadata_clears_but_shared_pdf_images_require_manual_review(tmp_path: Path) -> None:
    image_path = _create_color_image(tmp_path / "inputs" / "green.png", (0, 255, 0))
    replacement_path = _create_color_image(tmp_path / "inputs" / "yellow.png", (255, 255, 0))
    source = _create_pdf_with_metadata_and_shared_image(tmp_path / "shared-image.pdf", image_path)
    findings = _findings_for(tmp_path, source.name, categories={"metadata", "images"})

    results = transform(
        findings,
        _policy(
            metadata={"enabled": True, "action": "clear"},
            images={"enabled": True, "mode": "replace", "replacement_path": str(replacement_path)},
        ),
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "partial_success"
    assert [action["category"] for action in result["actions"]] == ["metadata", "images", "images"]

    metadata_action = next(action for action in result["actions"] if action["category"] == "metadata")
    image_actions = [action for action in result["actions"] if action["category"] == "images"]
    assert metadata_action["status"] == "applied"
    assert all(action["status"] == "manual_review_required" for action in image_actions)
    assert len(result["manual_review_items"]) == 2
    assert all(item["category"] == "images" for item in result["manual_review_items"])
    assert all("uniquely addressed raster image objects" in item["reason"] for item in result["manual_review_items"])

    assert all(value is None for value in read_metadata(source)["fields"].values())

    extracted = extract_images(source, tmp_path / "extracted" / "pdf-shared-review")
    assert [_image_color(path) for path in extracted] == [(0, 255, 0), (0, 255, 0)]



def test_transform_applies_confirmed_excel_body_text_and_preserves_sg3_header_changes(tmp_path: Path) -> None:
    source = _create_xlsx_with_body_text_and_header(tmp_path / "body-header.xlsx")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={"person_names": ["Jane Example"]},
        categories={"body_text", "headers"},
    )
    body_text_policy, candidates_by_text = _build_body_text_policy(
        [finding for finding in findings if finding["category"] == "body_text"],
        approved_texts=["Jane Example"],
        default_replacement_text="[PERSON]",
    )

    results = transform(
        findings,
        {
            **_policy(headers={"enabled": True, "action": "clear"}),
            "body_text": body_text_policy,
        },
    )

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "success"
    assert [action["category"] for action in result["actions"]] == ["headers", "body_text"]
    assert [action["status"] for action in result["actions"]] == ["applied", "applied"]

    body_text_action = next(action for action in result["actions"] if action["category"] == "body_text")
    assert body_text_action["candidate_id"] == candidates_by_text["Jane Example"]["candidate_id"]
    assert body_text_action["decision"] == "approved"
    assert body_text_action["requested_action"] == "replace"
    assert body_text_action["applied_action"] == "replace"
    assert body_text_action["details"]["locator_validated"] is True
    assert body_text_action["details"]["matched_text_before"] == "Jane Example"
    assert body_text_action["details"]["replacement_text"] == "[PERSON]"
    assert body_text_action["details"]["surface_kind"] == "excel_cell"

    workbook = load_workbook(source)
    try:
        worksheet = workbook["Sheet1"]
        assert worksheet["A1"].value == "Primary contact: [PERSON]"
        assert worksheet.oddHeader.left.text in (None, "")
    finally:
        workbook.close()



def test_transform_applies_confirmed_docx_body_text_to_paragraphs_and_table_cells(tmp_path: Path) -> None:
    source = _create_docx_with_body_text_targets(tmp_path / "body.docx")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={"person_names": ["Jane Example"], "company_names": ["Example Corp"]},
        categories={"body_text"},
    )
    body_text_policy, candidates_by_text = _build_body_text_policy(
        findings,
        approved_texts=["Jane Example", "Example Corp"],
        default_replacement_text="[REDACTED]",
        replacement_overrides={"Jane Example": "[PERSON]", "Example Corp": "[COMPANY]"},
    )

    results = transform(findings, {**_policy(), "body_text": body_text_policy})

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "success"
    assert [action["status"] for action in result["actions"]] == ["applied", "applied"]

    paragraph_action = next(action for action in result["actions"] if action["details"]["surface_kind"] == "docx_paragraph")
    table_action = next(action for action in result["actions"] if action["details"]["surface_kind"] == "docx_table_cell")
    assert paragraph_action["candidate_id"] == candidates_by_text["Jane Example"]["candidate_id"]
    assert table_action["candidate_id"] == candidates_by_text["Example Corp"]["candidate_id"]

    document = Document(source)
    assert document.paragraphs[0].text == "Primary contact: [PERSON]"
    assert document.tables[0].cell(0, 0).text == "Vendor: [COMPANY]"



def test_transform_applies_confirmed_pptx_body_text_to_text_frames_and_table_cells(tmp_path: Path) -> None:
    source = _create_pptx_with_body_text_targets(tmp_path / "body.pptx")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={"person_names": ["Jane Example"], "company_names": ["Example Corp"]},
        categories={"body_text"},
    )
    body_text_policy, _ = _build_body_text_policy(
        findings,
        approved_texts=["Jane Example", "Example Corp"],
        default_replacement_text="[REDACTED]",
        replacement_overrides={"Jane Example": "[PERSON]", "Example Corp": "[COMPANY]"},
    )

    results = transform(findings, {**_policy(), "body_text": body_text_policy})

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "success"
    assert [action["status"] for action in result["actions"]] == ["applied", "applied"]
    assert {action["details"]["surface_kind"] for action in result["actions"]} == {"pptx_text_frame", "pptx_table_cell"}

    presentation = Presentation(source)
    slide = presentation.slides[0]
    text_shape = next(shape for shape in slide.shapes if getattr(shape, "has_text_frame", False) and not getattr(shape, "has_table", False))
    table_shape = next(shape for shape in slide.shapes if getattr(shape, "has_table", False))
    assert text_shape.text_frame.paragraphs[0].text == "Primary contact: [PERSON]"
    assert table_shape.table.cell(0, 0).text == "Vendor: [COMPANY]"



def test_transform_skips_rejected_and_undecided_body_text_candidates(tmp_path: Path) -> None:
    source = _create_xlsx_with_multiple_body_text_candidates(tmp_path / "multiple.xlsx")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={
            "person_names": ["Jane Example", "Bob Example"],
            "company_names": ["Example Corp"],
        },
        categories={"body_text"},
    )
    body_text_policy, candidates_by_text = _build_body_text_policy(
        findings,
        approved_texts=["Jane Example"],
        rejected_texts=["Example Corp"],
        default_replacement_text="[REDACTED]",
        replacement_overrides={"Jane Example": "[PERSON]"},
    )

    results = transform(findings, {**_policy(), "body_text": body_text_policy})

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "partial_success"

    actions_by_candidate = {action["candidate_id"]: action for action in result["actions"]}
    assert actions_by_candidate[candidates_by_text["Jane Example"]["candidate_id"]]["status"] == "applied"
    assert actions_by_candidate[candidates_by_text["Example Corp"]["candidate_id"]]["status"] == "skipped"
    assert actions_by_candidate[candidates_by_text["Example Corp"]["candidate_id"]]["decision"] == "rejected"
    assert actions_by_candidate[candidates_by_text["Bob Example"]["candidate_id"]]["status"] == "skipped"
    assert actions_by_candidate[candidates_by_text["Bob Example"]["candidate_id"]]["decision"] == "undecided"

    workbook = load_workbook(source)
    try:
        worksheet = workbook["Sheet1"]
        assert worksheet["A1"].value == "Primary contact: [PERSON]"
        assert worksheet["A2"].value == "Vendor: Example Corp"
        assert worksheet["A3"].value == "Backup contact: Bob Example"
    finally:
        workbook.close()



def test_transform_marks_stale_body_text_locator_for_manual_review(tmp_path: Path) -> None:
    source = _create_xlsx_with_body_text_and_header(tmp_path / "stale.xlsx")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={"person_names": ["Jane Example"]},
        categories={"body_text"},
    )
    body_text_policy, _ = _build_body_text_policy(
        findings,
        approved_texts=["Jane Example"],
        default_replacement_text="[PERSON]",
    )

    workbook = load_workbook(source)
    try:
        worksheet = workbook["Sheet1"]
        worksheet["A1"] = "Updated primary contact: Jane Example"
        workbook.save(source)
    finally:
        workbook.close()

    results = transform(findings, {**_policy(), "body_text": body_text_policy})

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "manual_review_required"
    action = result["actions"][0]
    assert action["status"] == "manual_review_required"
    assert action["applied_action"] == "manual_review"
    assert action["details"]["locator_validated"] is False
    assert "changed since detection" in action["manual_review_reason"]

    workbook = load_workbook(source)
    try:
        worksheet = workbook["Sheet1"]
        assert worksheet["A1"].value == "Updated primary contact: Jane Example"
    finally:
        workbook.close()



def test_transform_keeps_pdf_body_text_manual_review_even_when_candidate_is_approved(tmp_path: Path) -> None:
    source = _create_pdf_with_visible_body_text(tmp_path / "body.pdf")
    findings = _findings_for_with_body_text(
        tmp_path,
        source.name,
        candidate_inputs={},
        categories={"body_text"},
    )
    body_text_policy, _ = _build_body_text_policy(
        findings,
        approved_texts=["john@example.com"],
        default_replacement_text="[EMAIL]",
    )

    results = transform(findings, {**_policy(), "body_text": body_text_policy})

    assert len(results) == 1
    result = results[0]
    assert result["status"] == "manual_review_required"
    action = result["actions"][0]
    assert action["candidate_id"]
    assert action["decision"] == "approved"
    assert action["status"] == "manual_review_required"
    assert action["applied_action"] == "manual_review"
    assert "PDF text-layer matches remain review-first" in action["manual_review_reason"]

    document = fitz.open(source)
    try:
        assert "john@example.com" in document[0].get_text()
    finally:
        document.close()
