from __future__ import annotations

import hashlib
from pathlib import Path

import fitz
import pytest
from docx import Document
from openpyxl import Workbook
from openpyxl.comments import Comment
from openpyxl.drawing.image import Image as OpenPyxlImage
from PIL import Image
from pptx import Presentation

from office_automation.anonymize import detect as detect_module
from office_automation.anonymize.detect import detect


def _sha256(path: Path) -> str:
    return hashlib.sha256(path.read_bytes()).hexdigest()



def _create_png(path: Path, color: tuple[int, int, int], *, size: tuple[int, int] = (32, 24)) -> Path:
    path.parent.mkdir(parents=True, exist_ok=True)
    Image.new("RGB", size, color).save(path, format="PNG")
    return path



def _create_docx_with_comment_and_header(path: Path) -> Path:
    document = Document()
    section = document.sections[0]
    section.header.paragraphs[0].text = "Document Header"
    section.footer.paragraphs[0].text = "Document Footer"
    paragraph = document.add_paragraph()
    run = paragraph.add_run("Hello document")
    document.add_comment(run, text="Remove this comment", author="Alice", initials="AL")
    document.core_properties.author = "Doc Author"
    document.core_properties.title = "Doc Title"
    document.save(path)
    return path



def _create_xlsx_with_note_header_footer_and_image(path: Path) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "SheetA"
    worksheet["A1"] = "cell"
    worksheet["A1"].comment = Comment("Secret note", "Bob")
    worksheet.oddHeader.left.text = "Left Header"
    worksheet.oddFooter.center.text = "Center Footer"
    workbook.properties.creator = "Workbook Creator"
    workbook.properties.title = "Workbook Title"

    image_path = _create_png(path.parent / "xlsx-image.png", (255, 0, 0))
    worksheet.add_image(OpenPyxlImage(str(image_path)), "C3")

    workbook.save(path)
    workbook.close()
    return path



def _create_pptx_with_notes(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    slide.shapes.title.text = "Deck Title"
    slide.notes_slide.notes_text_frame.text = "Presenter note"
    presentation.core_properties.author = "Slides Author"
    presentation.save(path)
    return path



def _create_xlsx_with_body_text(path: Path) -> Path:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Sheet1"
    worksheet["B2"] = "Budget owner: Project Lotus"
    workbook.save(path)
    workbook.close()
    return path



def _create_docx_with_body_text(path: Path) -> Path:
    document = Document()
    document.add_paragraph("Primary contact: Jane Example")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Client: Example Corp"
    document.save(path)
    return path



def _create_pptx_with_body_text(path: Path) -> Path:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    textbox = slide.shapes.add_textbox(left=914400, top=914400, width=3657600, height=914400)
    textbox.text_frame.text = "Call 03 1234 5678 for approval"
    presentation.save(path)
    return path



def _create_pdf_with_body_text(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "Contact john@example.com for the update")
    document.save(path)
    document.close()
    return path



def _create_pdf_with_multi_word_body_text(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text(
        (72, 72),
        (
            "This PDF paragraph includes a long searchable text layer before the sensitive phrase "
            "Jane Example so that excerpts must stay bounded around the actual match instead of "
            "returning the whole page content to reviewers."
        ),
    )
    document.save(path)
    document.close()
    return path



def _create_pdf_with_metadata(path: Path) -> Path:
    document = fitz.open()
    page = document.new_page()
    page.insert_text((72, 72), "PDF body")
    document.set_metadata({"author": "PDF Author", "title": "PDF Title"})
    document.save(path)
    document.close()
    return path



def test_detect_rejects_missing_and_non_directory_target_folder(tmp_path: Path) -> None:
    missing = tmp_path / "missing"
    not_a_directory = tmp_path / "file.docx"
    not_a_directory.write_text("not a folder")

    with pytest.raises(FileNotFoundError, match=r"Target folder '.*missing' does not exist\."):
        detect(str(missing))

    with pytest.raises(NotADirectoryError, match=r"Target folder '.*file\.docx' is not a directory\."):
        detect(str(not_a_directory))



def test_detect_default_extensions_and_override_subset(tmp_path: Path) -> None:
    docx_path = _create_docx_with_comment_and_header(tmp_path / "alpha.docx")
    pdf_path = _create_pdf_with_metadata(tmp_path / "beta.pdf")
    (tmp_path / "ignore.txt").write_text("ignore me")

    default_results = detect(tmp_path)
    default_relative_paths = {finding["relative_path"] for finding in default_results}
    assert default_relative_paths == {docx_path.name, pdf_path.name}

    subset_results = detect(tmp_path, extensions=[".DOCX"])
    assert {finding["relative_path"] for finding in subset_results} == {docx_path.name}

    with pytest.raises(ValueError, match=r"Unsupported extension override: xls"):
        detect(tmp_path, extensions=["docx", "xls"])



def test_detect_returns_deterministic_file_category_and_location_order(tmp_path: Path) -> None:
    docx_path = _create_docx_with_comment_and_header(tmp_path / "A.docx")
    xlsx_path = _create_xlsx_with_note_header_footer_and_image(tmp_path / "b.xlsx")

    first_run = detect(tmp_path)
    second_run = detect(tmp_path)

    assert [finding["finding_id"] for finding in first_run] == [finding["finding_id"] for finding in second_run]
    assert [(finding["relative_path"], finding["category"]) for finding in first_run] == [
        (docx_path.name, "comments"),
        (docx_path.name, "headers"),
        (docx_path.name, "footers"),
        (docx_path.name, "metadata"),
        (xlsx_path.name, "notes"),
        (xlsx_path.name, "headers"),
        (xlsx_path.name, "footers"),
        (xlsx_path.name, "metadata"),
        (xlsx_path.name, "images"),
    ]



def test_detect_uses_shared_metadata_and_image_helpers(tmp_path: Path, monkeypatch: pytest.MonkeyPatch) -> None:
    source = _create_docx_with_comment_and_header(tmp_path / "sample.docx")
    observed_metadata_paths: list[Path] = []
    observed_image_paths: list[Path] = []

    def fake_read_metadata(file_path: str | Path) -> dict:
        observed_metadata_paths.append(Path(file_path))
        return {
            "file_path": str(file_path),
            "extension": "docx",
            "fields": {"creator": "Pat", "title": "Plan", "subject": None},
            "warnings": [],
        }

    def fake_extract_images(file_path: str | Path, output_dir: str | Path) -> list[Path]:
        observed_image_paths.append(Path(file_path))
        created = _create_png(Path(output_dir) / "fake-image.png", (0, 255, 0), size=(10, 20))
        return [created]

    monkeypatch.setattr(detect_module, "read_metadata", fake_read_metadata)
    monkeypatch.setattr(detect_module, "extract_images", fake_extract_images)

    findings = detect(tmp_path, extensions=["docx"])

    metadata_finding = next(finding for finding in findings if finding["category"] == "metadata")
    image_finding = next(finding for finding in findings if finding["category"] == "images")

    assert observed_metadata_paths == [source]
    assert observed_image_paths == [source]
    assert metadata_finding["payload"]["fields"] == {"creator": "Pat", "title": "Plan"}
    assert image_finding["location"]["image_index"] == 0
    assert image_finding["payload"]["width"] == 10
    assert image_finding["payload"]["height"] == 20



def test_detect_finds_representative_comment_note_header_footer_paths_without_mutating_sources(tmp_path: Path) -> None:
    docx_path = _create_docx_with_comment_and_header(tmp_path / "commented.docx")
    xlsx_path = _create_xlsx_with_note_header_footer_and_image(tmp_path / "sheet.xlsx")
    before_hashes = {path: _sha256(path) for path in (docx_path, xlsx_path)}

    findings = detect(tmp_path)

    comments_finding = next(
        finding for finding in findings if finding["relative_path"] == docx_path.name and finding["category"] == "comments"
    )
    docx_header_finding = next(
        finding
        for finding in findings
        if finding["relative_path"] == docx_path.name and finding["category"] == "headers"
    )
    xlsx_note_finding = next(
        finding for finding in findings if finding["relative_path"] == xlsx_path.name and finding["category"] == "notes"
    )
    xlsx_footer_finding = next(
        finding
        for finding in findings
        if finding["relative_path"] == xlsx_path.name and finding["category"] == "footers"
    )

    assert comments_finding["payload"]["text"] == "Remove this comment"
    assert comments_finding["payload"]["author"] == "Alice"
    assert docx_header_finding["payload"]["text"] == "Document Header"
    assert xlsx_note_finding["payload"]["text"] == "Secret note"
    assert xlsx_note_finding["payload"]["author"] == "Bob"
    assert xlsx_footer_finding["payload"]["text"] == "Center Footer"

    after_hashes = {path: _sha256(path) for path in (docx_path, xlsx_path)}
    assert after_hashes == before_hashes



def test_detect_emits_low_confidence_manual_review_findings_for_pptx_surfaces(tmp_path: Path) -> None:
    pptx_path = _create_pptx_with_notes(tmp_path / "deck.pptx")

    findings = detect(tmp_path)

    manual_review_findings = [
        finding
        for finding in findings
        if finding["relative_path"] == pptx_path.name and finding["confidence"] == "low"
    ]
    assert [finding["category"] for finding in manual_review_findings] == ["comments", "headers", "footers"]
    assert all(finding["action_hint"] == "review" for finding in manual_review_findings)
    assert all(finding["manual_review_reason"] for finding in manual_review_findings)

    notes_finding = next(
        finding for finding in findings if finding["relative_path"] == pptx_path.name and finding["category"] == "notes"
    )
    assert notes_finding["confidence"] == "high"
    assert notes_finding["payload"]["text"] == "Presenter note"



def test_detect_emits_excel_body_text_findings_from_exact_user_hints(tmp_path: Path) -> None:
    workbook_path = _create_xlsx_with_body_text(tmp_path / "body.xlsx")

    findings = detect(
        tmp_path,
        extensions=["xlsx"],
        body_text_candidate_inputs={"exact_phrases": ["Project Lotus"]},
    )

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert len(body_text_findings) == 1

    finding = body_text_findings[0]
    assert finding["relative_path"] == workbook_path.name
    assert finding["location"] == {"cell": "B2", "match_end": 27, "match_start": 14, "sheet": "Sheet1"}
    assert finding["payload"] == {
        "excerpt": "Budget owner: Project Lotus",
        "matched_text": "Project Lotus",
        "normalized_text": "project lotus",
        "surface_type": "excel_cell",
    }
    assert finding["action_hint"] == "candidate_confirmation_required"
    assert finding["confidence"] == "high"
    assert finding["manual_review_reason"] is None
    assert finding["source"] == "user_hint"
    assert finding["reason_tags"] == ["exact_phrase"]
    assert finding["finding_id"] == "body.xlsx::body_text::{\"cell\":\"B2\",\"match_end\":27,\"match_start\":14,\"sheet\":\"Sheet1\"}"



def test_detect_emits_docx_body_text_findings_for_paragraphs_and_table_cells(tmp_path: Path) -> None:
    docx_path = _create_docx_with_body_text(tmp_path / "body.docx")

    findings = detect(
        tmp_path,
        extensions=["docx"],
        body_text_candidate_inputs={
            "person_names": ["Jane Example"],
            "company_names": ["Example Corp"],
        },
    )

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert {finding["payload"]["matched_text"] for finding in body_text_findings} == {"Jane Example", "Example Corp"}

    paragraph_finding = next(finding for finding in body_text_findings if finding["payload"]["matched_text"] == "Jane Example")
    assert paragraph_finding["relative_path"] == docx_path.name
    assert paragraph_finding["location"] == {
        "match_end": 29,
        "match_start": 17,
        "paragraph_index": 0,
        "surface": "paragraph",
    }
    assert paragraph_finding["payload"]["surface_type"] == "docx_paragraph"
    assert paragraph_finding["payload"]["normalized_text"] == "jane example"
    assert paragraph_finding["source"] == "user_hint"
    assert paragraph_finding["reason_tags"] == ["person_hint"]

    table_cell_finding = next(finding for finding in body_text_findings if finding["payload"]["matched_text"] == "Example Corp")
    assert table_cell_finding["location"] == {
        "column_index": 0,
        "match_end": 20,
        "match_start": 8,
        "paragraph_index": 0,
        "row_index": 0,
        "surface": "table_cell",
        "table_index": 0,
    }
    assert table_cell_finding["payload"]["surface_type"] == "docx_table_cell"
    assert table_cell_finding["payload"]["normalized_text"] == "example corp"
    assert table_cell_finding["source"] == "user_hint"
    assert table_cell_finding["reason_tags"] == ["company_hint"]



def test_detect_emits_pptx_body_text_findings_from_phone_patterns(tmp_path: Path) -> None:
    pptx_path = _create_pptx_with_body_text(tmp_path / "body.pptx")

    findings = detect(tmp_path, extensions=["pptx"], body_text_candidate_inputs={})

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert len(body_text_findings) == 1

    finding = body_text_findings[0]
    assert finding["relative_path"] == pptx_path.name
    assert finding["location"]["slide_number"] == 1
    assert finding["location"]["paragraph_index"] == 0
    assert isinstance(finding["location"]["shape_id"], int)
    assert finding["location"]["match_start"] == 5
    assert finding["location"]["match_end"] == 17
    assert finding["payload"] == {
        "excerpt": "Call 03 1234 5678 for approval",
        "matched_text": "03 1234 5678",
        "normalized_text": "03 1234 5678",
        "surface_type": "pptx_text_frame",
    }
    assert finding["action_hint"] == "candidate_confirmation_required"
    assert finding["confidence"] == "medium"
    assert finding["manual_review_reason"] is None
    assert finding["source"] == "pattern"
    assert finding["reason_tags"] == ["phone_pattern"]



def test_detect_emits_conservative_pdf_body_text_review_findings(tmp_path: Path) -> None:
    pdf_path = _create_pdf_with_body_text(tmp_path / "body.pdf")

    findings = detect(tmp_path, extensions=["pdf"], body_text_candidate_inputs={})

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert len(body_text_findings) == 1

    finding = body_text_findings[0]
    assert finding["relative_path"] == pdf_path.name
    assert finding["location"]["page_number"] == 1
    assert finding["location"]["span_index"] == 0
    assert finding["location"]["match_start"] == 8
    assert finding["location"]["match_end"] == 24
    assert finding["location"]["bbox"]
    assert finding["payload"] == {
        "excerpt": "Contact john@example.com for the update",
        "matched_text": "john@example.com",
        "normalized_text": "john@example.com",
        "surface_type": "pdf_text_span",
    }
    assert finding["action_hint"] == "review"
    assert finding["confidence"] == "low"
    assert "PDF text-layer matches remain review-first" in finding["manual_review_reason"]
    assert finding["source"] == "pattern"
    assert finding["reason_tags"] == ["email_pattern", "pdf_text_layer"]



def test_detect_preserves_legacy_sg3_behavior_when_body_text_detection_is_not_requested(tmp_path: Path) -> None:
    _create_pptx_with_body_text(tmp_path / "body.pptx")
    _create_pdf_with_body_text(tmp_path / "body.pdf")

    findings = detect(tmp_path, extensions=["pptx", "pdf"])

    assert all(finding["category"] != "body_text" for finding in findings)



def test_detect_preserves_sg3_findings_and_body_text_order_deterministically(tmp_path: Path) -> None:
    docx_path = _create_docx_with_comment_and_header(tmp_path / "alpha.docx")
    xlsx_path = _create_xlsx_with_note_header_footer_and_image(tmp_path / "beta.xlsx")

    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "Body"
    worksheet["A1"] = "Confidential Project Atlas"
    workbook.save(tmp_path / "gamma.xlsx")
    workbook.close()

    first_run = detect(
        tmp_path,
        body_text_candidate_inputs={"exact_phrases": ["Project Atlas"]},
    )
    second_run = detect(
        tmp_path,
        body_text_candidate_inputs={"exact_phrases": ["Project Atlas"]},
    )

    assert [finding["finding_id"] for finding in first_run] == [finding["finding_id"] for finding in second_run]
    assert [(finding["relative_path"], finding["category"]) for finding in first_run] == [
        (docx_path.name, "comments"),
        (docx_path.name, "headers"),
        (docx_path.name, "footers"),
        (docx_path.name, "metadata"),
        (xlsx_path.name, "notes"),
        (xlsx_path.name, "headers"),
        (xlsx_path.name, "footers"),
        (xlsx_path.name, "metadata"),
        (xlsx_path.name, "images"),
        ("gamma.xlsx", "metadata"),
        ("gamma.xlsx", "body_text"),
    ]



def test_detect_rejects_malformed_normalized_body_text_runtime_payloads(tmp_path: Path) -> None:
    _create_xlsx_with_body_text(tmp_path / "body.xlsx")

    with pytest.raises(TypeError, match=r"body_text_candidate_inputs field 'emails' must be a list of strings"):
        detect(tmp_path, extensions=["xlsx"], body_text_candidate_inputs={"emails": "user@example.com"})

    with pytest.raises(ValueError, match=r"Unsupported body_text_candidate_inputs key 'has_any_values'\."):
        detect(tmp_path, extensions=["xlsx"], body_text_candidate_inputs={"has_any_values": True})



def test_detect_emits_context_assisted_heuristic_body_text_findings(tmp_path: Path) -> None:
    docx_path = _create_docx_with_body_text(tmp_path / "body.docx")

    findings = detect(
        tmp_path,
        extensions=["docx"],
        body_text_candidate_inputs={"context_terms": ["contact"]},
    )

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert len(body_text_findings) == 1

    finding = body_text_findings[0]
    assert finding["relative_path"] == docx_path.name
    assert finding["location"] == {
        "match_end": 29,
        "match_start": 17,
        "paragraph_index": 0,
        "surface": "paragraph",
    }
    assert finding["payload"] == {
        "excerpt": "Primary contact: Jane Example",
        "matched_text": "Jane Example",
        "normalized_text": "jane example",
        "surface_type": "docx_paragraph",
    }
    assert finding["action_hint"] == "candidate_confirmation_required"
    assert finding["confidence"] == "medium"
    assert finding["manual_review_reason"] is None
    assert finding["source"] == "heuristic"
    assert finding["reason_tags"] == ["context_assisted_phrase", "context_term"]



def test_detect_emits_pdf_multi_word_hint_findings_with_bounded_excerpt(tmp_path: Path) -> None:
    pdf_path = _create_pdf_with_multi_word_body_text(tmp_path / "body.pdf")

    findings = detect(
        tmp_path,
        extensions=["pdf"],
        body_text_candidate_inputs={"person_names": ["Jane Example"]},
    )

    body_text_findings = [finding for finding in findings if finding["category"] == "body_text"]
    assert len(body_text_findings) == 1

    finding = body_text_findings[0]
    assert finding["relative_path"] == pdf_path.name
    assert finding["location"]["page_number"] == 1
    assert finding["location"]["bbox"]
    assert finding["payload"]["matched_text"] == "Jane Example"
    assert finding["payload"]["normalized_text"] == "jane example"
    assert finding["payload"]["surface_type"] == "pdf_text_span"
    assert finding["payload"]["excerpt"] == "… text layer before the sensitive phrase Jane Example so that"
    assert len(finding["payload"]["excerpt"]) < 120
    assert finding["action_hint"] == "review"
    assert finding["confidence"] == "low"
    assert "PDF text-layer matches remain review-first" in finding["manual_review_reason"]
    assert finding["source"] == "user_hint"
    assert finding["reason_tags"] == ["pdf_text_layer", "person_hint"]
