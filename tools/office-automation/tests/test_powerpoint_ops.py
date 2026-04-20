from __future__ import annotations

from pathlib import Path
from zipfile import ZIP_DEFLATED, ZipFile

import pytest
from pptx import Presentation
from pptx.util import Inches

from office_automation.powerpoint_ops import edit, read


def test_read_reports_slides_text_shapes_and_tables(tmp_path: Path) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)

    result = read(source)

    assert result["format"] == "pptx"
    assert result["file_path"] == str(source)
    assert result["slide_count"] == 1
    assert result["warnings"] == []

    slide = result["slides"][0]
    assert slide["slide_index"] == 0
    assert slide["slide_number"] == 1
    assert slide["title_text"] == "Quarterly Update"

    text_by_name = {shape["shape_name"]: shape for shape in slide["text_shapes"]}
    assert text_by_name["Title 1"]["is_title"] is True
    assert text_by_name["Title 1"]["text"] == "Quarterly Update"
    assert text_by_name["TextBox 2"]["text"] == "Revenue is up"

    table = slide["tables"][0]
    assert table["table_index"] == 0
    assert table["row_count"] == 2
    assert table["column_count"] == 2
    assert table["rows"] == [["Metric", "Value"], ["Revenue", "100"]]
    assert table["cells"] == [
        {"row_index": 0, "column_index": 0, "text": "Metric"},
        {"row_index": 0, "column_index": 1, "text": "Value"},
        {"row_index": 1, "column_index": 0, "text": "Revenue"},
        {"row_index": 1, "column_index": 1, "text": "100"},
    ]


def test_edit_updates_title_shapes_and_tables_and_honors_output_path(tmp_path: Path) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)
    source_snapshot = read(source)
    slide_data = source_snapshot["slides"][0]
    body_shape = next(shape for shape in slide_data["text_shapes"] if shape["shape_name"] == "TextBox 2")

    explicit_output = tmp_path / "exports" / "custom-output.pptx"
    ignored_output_dir = tmp_path / "ignored-dir"

    result = edit(
        source,
        {
            "operations": [
                {
                    "type": "replace_title_text",
                    "slide_index": 0,
                    "new_text": "Updated Quarterly Update",
                },
                {
                    "type": "replace_shape_text",
                    "slide_number": 1,
                    "shape_index": body_shape["shape_index"],
                    "new_text": "Revenue is significantly up",
                },
                {
                    "type": "replace_table_cell",
                    "slide_index": 0,
                    "table_index": 0,
                    "row_index": 1,
                    "column_index": 1,
                    "new_text": "125",
                },
                {
                    "type": "append_table_row",
                    "slide_index": 0,
                    "table_index": 0,
                    "values": ["Profit", "45"],
                },
            ],
            "output_path": explicit_output,
            "output_dir": ignored_output_dir,
            "copy_before_edit": True,
            "options": {},
        },
    )

    assert result == explicit_output
    assert explicit_output.exists()
    assert not ignored_output_dir.exists()

    original_after_edit = read(source)
    assert original_after_edit["slides"][0]["title_text"] == "Quarterly Update"
    assert original_after_edit["slides"][0]["tables"][0]["rows"] == [["Metric", "Value"], ["Revenue", "100"]]

    edited_after = read(explicit_output)
    edited_slide = edited_after["slides"][0]
    assert edited_slide["title_text"] == "Updated Quarterly Update"
    edited_text_map = {shape["shape_name"]: shape["text"] for shape in edited_slide["text_shapes"]}
    assert edited_text_map["TextBox 2"] == "Revenue is significantly up"
    assert edited_slide["tables"][0]["rows"] == [
        ["Metric", "Value"],
        ["Revenue", "125"],
        ["Profit", "45"],
    ]


@pytest.mark.parametrize(
    "instructions",
    [
        {
            "operations": [
                {
                    "type": "replace_title_text",
                    "slide_index": 0,
                    "new_text": "Missing Output Path",
                }
            ],
            "options": {},
        },
        {
            "operations": [
                {
                    "type": "replace_title_text",
                    "slide_index": 0,
                    "new_text": "Missing Output Path",
                }
            ],
            "output_path": "",
            "options": {},
        },
        {
            "operations": [
                {
                    "type": "replace_title_text",
                    "slide_index": 0,
                    "new_text": "Missing Output Path",
                }
            ],
            "output_dir": "ignored-without-output-path",
            "options": {},
        },
    ],
)
def test_edit_requires_explicit_output_path(tmp_path: Path, instructions: dict) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)

    with pytest.raises(
        ValueError,
        match=r"PowerPoint edit instructions must include 'output_path'; runtime save-path fallback is not supported\.",
    ):
        edit(source, instructions)


def test_edit_rejects_same_output_path_when_copy_before_edit_is_enabled(tmp_path: Path) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)

    with pytest.raises(
        ValueError,
        match=r"copy_before_edit=True requires output_path to be different from the source presentation path\.",
    ):
        edit(
            source,
            {
                "operations": [
                    {
                        "type": "replace_title_text",
                        "slide_index": 0,
                        "new_text": "Should Not Save",
                    }
                ],
                "output_path": str(source),
                "copy_before_edit": True,
                "options": {},
            },
        )


@pytest.mark.parametrize("operation_type", ["rewrite_animation", "edit_smartart"])
def test_edit_rejects_unsupported_operation_types(tmp_path: Path, operation_type: str) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)

    with pytest.raises(NotImplementedError, match=r"Unsupported PowerPoint operation"):
        edit(
            source,
            {
                "operations": [
                    {
                        "type": operation_type,
                        "slide_index": 0,
                    }
                ],
                "output_path": tmp_path / "out.pptx",
                "options": {},
            },
        )


def test_edit_rejects_non_pptx_output_paths(tmp_path: Path) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)

    with pytest.raises(ValueError, match=r"Output path '.*out\.txt' has unsupported extension 'txt'"):
        edit(
            source,
            {
                "operations": [
                    {
                        "type": "replace_title_text",
                        "slide_index": 0,
                        "new_text": "Wrong Destination",
                    }
                ],
                "output_path": tmp_path / "out.txt",
                "options": {},
            },
        )


def test_read_and_edit_surface_animation_warnings(tmp_path: Path) -> None:
    source = tmp_path / "quarterly-update.pptx"
    _create_sample_presentation(source)
    _inject_animation_timing(source)

    read_result = read(source)
    assert any("animation timing data" in message for message in read_result["warnings"])

    with pytest.warns(UserWarning, match=r"animation timing data"):
        result = edit(
            source,
            {
                "operations": [
                    {
                        "type": "replace_title_text",
                        "slide_index": 0,
                        "new_text": "Animated But Editable",
                    }
                ],
                "output_path": tmp_path / "animated-edited.pptx",
                "options": {},
            },
        )

    assert result.exists()
    assert read(result)["slides"][0]["title_text"] == "Animated But Editable"


def _create_sample_presentation(path: Path) -> None:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    assert slide.shapes.title is not None
    slide.shapes.title.text = "Quarterly Update"

    text_box = slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(1))
    text_box.text = "Revenue is up"

    table_shape = slide.shapes.add_table(2, 2, Inches(1), Inches(3), Inches(4.5), Inches(1.5))
    table = table_shape.table
    table.cell(0, 0).text = "Metric"
    table.cell(0, 1).text = "Value"
    table.cell(1, 0).text = "Revenue"
    table.cell(1, 1).text = "100"

    path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(path)


def _inject_animation_timing(path: Path) -> None:
    updated_entries: dict[str, bytes] = {}
    with ZipFile(path) as archive:
        for info in archive.infolist():
            data = archive.read(info.filename)
            if info.filename == "ppt/slides/slide1.xml" and b"<p:timing" not in data:
                data = data.replace(b"</p:sld>", b"<p:timing/></p:sld>")
            updated_entries[info.filename] = data

    with ZipFile(path, "w", compression=ZIP_DEFLATED) as archive:
        for name, data in updated_entries.items():
            archive.writestr(name, data)
