"""PowerPoint runtime helpers for SG2 V1 `pptx` read/edit flows.

Supported edit instructions use the SG2 envelope shape:

{
    "operations": [
        {
            "type": "replace_title_text",
            "slide_index": 0,
            "new_text": "Updated title",
        },
        {
            "type": "replace_shape_text",
            "slide_index": 0,
            "shape_index": 2,
            "new_text": "Updated body text",
        },
        {
            "type": "replace_table_cell",
            "slide_index": 0,
            "table_index": 0,
            "row_index": 1,
            "column_index": 0,
            "new_text": "Updated cell",
        },
        {
            "type": "append_table_row",
            "slide_index": 0,
            "table_index": 0,
            "values": ["left", "right"],
        },
    ],
    "output_path": "/path/to/deck-edited.pptx",
    "output_dir": "/ignored/by/runtime/when/output_path/is/present",
    "copy_before_edit": True,
    "options": {},
}

V1 target resolution rules:
- slide targets may use `slide_index` (0-based) or `slide_number` (1-based)
- shape targets may use `shape_index`, `shape_id`, or `shape_name`
- table targets may use `table_index`, or a shape selector that resolves to a table shape

Out of scope in V1:
- SmartArt-specific editing
- advanced animation rewriting
- slide master / layout surgery
- arbitrary redesign beyond existing text/basic-shape/table updates
"""

from __future__ import annotations

from collections.abc import Iterable, Sequence
from copy import deepcopy
from os import PathLike
from pathlib import Path
import re
import shutil
import warnings
from xml.etree import ElementTree as ET
from zipfile import ZipFile

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from office_automation.common.files import copy_original

__all__ = ["read", "edit"]

_SUPPORTED_EXTENSION = ".pptx"
_SUPPORTED_EXTENSION_TEXT = "pptx"
_SUPPORTED_OPERATIONS = (
    "replace_shape_text",
    "replace_table_cell",
    "replace_title_text",
    "append_table_row",
)
_PPT_NAMESPACES = {
    "p": "http://schemas.openxmlformats.org/presentationml/2006/main",
    "rel": "http://schemas.openxmlformats.org/package/2006/relationships",
}
_SLIDE_XML_PATTERN = re.compile(r"^ppt/slides/slide(\d+)\.xml$")
_SLIDE_RELS_PATTERN = re.compile(r"^ppt/slides/_rels/slide(\d+)\.xml\.rels$")


def read(file_path: str | Path) -> dict:
    """Read a `pptx` presentation and return SG2-oriented structure data."""
    presentation_path = _validate_presentation_path(file_path)
    presentation = Presentation(str(presentation_path))

    presentation_warnings = _collect_presentation_feature_warnings(presentation_path)
    slides: list[dict] = []
    warning_messages = list(presentation_warnings)

    for slide_index, slide in enumerate(presentation.slides):
        slide_data = _serialize_slide(slide, slide_index)
        slides.append(slide_data)
        warning_messages.extend(slide_data["warnings"])

    return {
        "format": _SUPPORTED_EXTENSION_TEXT,
        "file_path": str(presentation_path),
        "slide_count": len(presentation.slides),
        "slides": slides,
        "warnings": _deduplicate_messages(warning_messages),
    }


def edit(file_path: str | Path, instructions: dict) -> Path:
    """Apply supported `pptx` edits and return the saved output path."""
    presentation_path = _validate_presentation_path(file_path)
    operation_list, output_path, copy_before_edit = _normalize_edit_request(instructions)

    for message in _collect_presentation_feature_warnings(presentation_path):
        warnings.warn(message, UserWarning, stacklevel=2)

    load_path = _prepare_edit_load_path(
        presentation_path,
        output_path,
        copy_before_edit=copy_before_edit,
    )
    presentation = Presentation(str(load_path))

    for operation in operation_list:
        _apply_operation(presentation, operation)

    output_path.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(output_path))
    return output_path


def _serialize_slide(slide, slide_index: int) -> dict:
    title_shape = slide.shapes.title
    title_text = None
    if title_shape is not None and getattr(title_shape, "has_text_frame", False):
        title_text = _extract_text_frame_text(title_shape.text_frame)

    text_shapes: list[dict] = []
    tables: list[dict] = []
    slide_warnings: list[str] = []
    table_index = 0

    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            slide_warnings.append(
                f"Slide {slide_index + 1} shape {shape_index} is a grouped shape; grouped-shape editing is unsupported in V1."
            )
            continue

        if getattr(shape, "has_text_frame", False):
            text_shapes.append(
                {
                    "shape_index": shape_index,
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "shape_type": _shape_type_name(shape),
                    "is_title": bool(title_shape is not None and shape == title_shape),
                    "text": _extract_text_frame_text(shape.text_frame),
                }
            )

        if getattr(shape, "has_table", False):
            table = shape.table
            tables.append(
                {
                    "table_index": table_index,
                    "shape_index": shape_index,
                    "shape_id": shape.shape_id,
                    "shape_name": shape.name,
                    "row_count": len(table.rows),
                    "column_count": len(table.columns),
                    "rows": _table_rows_as_matrix(table),
                    "cells": _table_cells_as_records(table),
                }
            )
            table_index += 1

    return {
        "slide_index": slide_index,
        "slide_number": slide_index + 1,
        "title_text": title_text,
        "shape_count": len(slide.shapes),
        "text_shapes": text_shapes,
        "tables": tables,
        "warnings": _deduplicate_messages(slide_warnings),
    }


def _table_rows_as_matrix(table) -> list[list[str]]:
    return [
        [table.cell(row_index, column_index).text for column_index in range(len(table.columns))]
        for row_index in range(len(table.rows))
    ]


def _table_cells_as_records(table) -> list[dict]:
    cells: list[dict] = []
    for row_index in range(len(table.rows)):
        for column_index in range(len(table.columns)):
            cells.append(
                {
                    "row_index": row_index,
                    "column_index": column_index,
                    "text": table.cell(row_index, column_index).text,
                }
            )
    return cells


def _collect_presentation_feature_warnings(presentation_path: Path) -> list[str]:
    warning_messages: list[str] = []
    with ZipFile(presentation_path) as archive:
        for name in archive.namelist():
            slide_match = _SLIDE_XML_PATTERN.match(name)
            if slide_match:
                slide_number = int(slide_match.group(1))
                xml_bytes = archive.read(name)
                if b"<p:timing" in xml_bytes:
                    warning_messages.append(
                        f"Slide {slide_number} contains animation timing data; advanced animation rewriting is unsupported in V1."
                    )
                if b":dgm" in xml_bytes or b"drawingml/2006/diagram" in xml_bytes:
                    warning_messages.append(
                        f"Slide {slide_number} appears to contain SmartArt/diagram content; SmartArt-specific editing is unsupported in V1."
                    )
                continue

            rels_match = _SLIDE_RELS_PATTERN.match(name)
            if rels_match:
                slide_number = int(rels_match.group(1))
                xml_root = ET.fromstring(archive.read(name))
                for relationship in xml_root.findall("rel:Relationship", _PPT_NAMESPACES):
                    relationship_type = relationship.attrib.get("Type", "").lower()
                    target = relationship.attrib.get("Target", "").lower()
                    if "diagram" in relationship_type or "smartart" in relationship_type or "diagram" in target:
                        warning_messages.append(
                            f"Slide {slide_number} has SmartArt/diagram relationships; SmartArt-specific editing is unsupported in V1."
                        )
                        break

    return _deduplicate_messages(warning_messages)


def _normalize_edit_request(
    instructions: dict,
) -> tuple[list[dict], Path, bool]:
    if not isinstance(instructions, dict):
        raise TypeError("instructions must be a dict.")

    operations = instructions.get("operations")
    if not isinstance(operations, list) or not operations:
        raise ValueError("instructions['operations'] must be a non-empty list.")
    if not all(isinstance(operation, dict) for operation in operations):
        raise TypeError("instructions['operations'] entries must be dict objects.")

    options = instructions.get("options", {})
    if not isinstance(options, dict):
        raise TypeError("instructions['options'] must be a dict when provided.")

    output_path_value = instructions.get("output_path")
    if output_path_value in (None, ""):
        raise ValueError(
            "PowerPoint edit instructions must include 'output_path'; runtime save-path fallback is not supported."
        )

    output_path = Path(output_path_value)
    if output_path.suffix.lower() != _SUPPORTED_EXTENSION:
        display_extension = output_path.suffix.lower().lstrip(".") or "<none>"
        raise ValueError(
            f"Output path '{output_path}' has unsupported extension '{display_extension}'. Supported extensions: {_SUPPORTED_EXTENSION_TEXT}."
        )

    copy_before_edit = instructions.get("copy_before_edit", True)
    if not isinstance(copy_before_edit, bool):
        raise TypeError("instructions['copy_before_edit'] must be a bool when provided.")

    return operations, output_path, copy_before_edit


def _prepare_edit_load_path(
    source_path: Path,
    output_path: Path,
    *,
    copy_before_edit: bool,
) -> Path:
    output_path.parent.mkdir(parents=True, exist_ok=True)

    if _paths_refer_to_same_location(source_path, output_path):
        if copy_before_edit:
            raise ValueError(
                "copy_before_edit=True requires output_path to be different from the source presentation path."
            )
        return source_path

    if output_path.exists():
        raise FileExistsError(
            f"Output path '{output_path}' already exists; the runtime will not silently overwrite it."
        )

    if not copy_before_edit:
        return source_path

    staged_copy = copy_original(source_path, output_path.parent)
    if staged_copy != output_path:
        shutil.move(str(staged_copy), str(output_path))
    return output_path


def _apply_operation(presentation: Presentation, operation: dict) -> None:
    operation_type = operation.get("type") or operation.get("op") or operation.get("operation")
    if operation_type not in _SUPPORTED_OPERATIONS:
        supported_text = ", ".join(_SUPPORTED_OPERATIONS)
        raise NotImplementedError(
            f"Unsupported PowerPoint operation '{operation_type}'. Supported operations: {supported_text}."
        )

    if operation_type == "replace_title_text":
        slide = _resolve_slide(presentation, operation)
        title_shape = slide.shapes.title
        if title_shape is None or not getattr(title_shape, "has_text_frame", False):
            slide_number = _slide_number_from_operation(operation)
            raise ValueError(f"Slide {slide_number} does not have an editable title shape.")
        _set_text_frame_text(title_shape.text_frame, _resolve_new_text(operation))
        return

    if operation_type == "replace_shape_text":
        slide = _resolve_slide(presentation, operation)
        shape = _resolve_shape_target(slide, operation)
        if not getattr(shape, "has_text_frame", False):
            raise ValueError(
                f"Target shape '{shape.name}' on slide {_slide_number_from_operation(operation)} does not have an editable text frame."
            )
        if getattr(shape, "has_table", False):
            raise ValueError(
                f"Target shape '{shape.name}' on slide {_slide_number_from_operation(operation)} is a table. Use replace_table_cell or append_table_row instead."
            )
        _set_text_frame_text(shape.text_frame, _resolve_new_text(operation))
        return

    if operation_type == "replace_table_cell":
        slide = _resolve_slide(presentation, operation)
        table = _resolve_table_target(slide, operation)
        row_index = _require_non_negative_int(operation, "row_index")
        column_index = _require_non_negative_int(operation, "column_index")
        _validate_table_coordinates(table, row_index, column_index)
        _set_text_frame_text(
            table.cell(row_index, column_index).text_frame,
            _resolve_new_text(operation),
        )
        return

    if operation_type == "append_table_row":
        slide = _resolve_slide(presentation, operation)
        table = _resolve_table_target(slide, operation)
        row_values = _normalize_row_values(operation, expected_columns=len(table.columns))
        _append_table_row(table, row_values)
        return

    raise AssertionError(f"Unhandled PowerPoint operation '{operation_type}'.")


def _resolve_slide(presentation: Presentation, operation: dict):
    if "slide_index" in operation:
        slide_index = _require_non_negative_int(operation, "slide_index")
    elif "slide_number" in operation:
        slide_number = _require_positive_int(operation, "slide_number")
        slide_index = slide_number - 1
    else:
        raise ValueError("PowerPoint operations must include slide_index or slide_number.")

    if slide_index >= len(presentation.slides):
        raise IndexError(
            f"slide_index {slide_index} is out of range for presentation with {len(presentation.slides)} slides."
        )
    return presentation.slides[slide_index]


def _resolve_shape_target(slide, operation: dict):
    selectors = [
        key for key in ("shape_index", "shape_id", "shape_name") if key in operation and operation[key] is not None
    ]
    if not selectors:
        raise ValueError(
            "Shape-targeted PowerPoint operations must include one of shape_index, shape_id, or shape_name."
        )

    if "shape_index" in operation and operation["shape_index"] is not None:
        shape_index = _require_non_negative_int(operation, "shape_index")
        if shape_index >= len(slide.shapes):
            raise IndexError(
                f"shape_index {shape_index} is out of range for slide with {len(slide.shapes)} shapes."
            )
        return slide.shapes[shape_index]

    if "shape_id" in operation and operation["shape_id"] is not None:
        shape_id = _require_positive_int(operation, "shape_id")
        for shape in slide.shapes:
            if shape.shape_id == shape_id:
                return shape
        raise ValueError(f"No shape with shape_id {shape_id} exists on the target slide.")

    shape_name = operation["shape_name"]
    if not isinstance(shape_name, str) or not shape_name.strip():
        raise ValueError("shape_name must be a non-empty string when provided.")

    matching_shapes = [shape for shape in slide.shapes if shape.name == shape_name]
    if not matching_shapes:
        raise ValueError(f"No shape with shape_name '{shape_name}' exists on the target slide.")
    if len(matching_shapes) > 1:
        raise ValueError(
            f"shape_name '{shape_name}' is ambiguous on the target slide; use shape_index or shape_id instead."
        )
    return matching_shapes[0]


def _resolve_table_target(slide, operation: dict):
    if "table_index" in operation and operation["table_index"] is not None:
        table_index = _require_non_negative_int(operation, "table_index")
        table_shapes = [shape for shape in slide.shapes if getattr(shape, "has_table", False)]
        if table_index >= len(table_shapes):
            raise IndexError(
                f"table_index {table_index} is out of range for slide with {len(table_shapes)} tables."
            )
        return table_shapes[table_index].table

    shape = _resolve_shape_target(slide, operation)
    if not getattr(shape, "has_table", False):
        raise ValueError(
            f"Target shape '{shape.name}' on slide {_slide_number_from_operation(operation)} is not a table."
        )
    return shape.table


def _validate_table_coordinates(table, row_index: int, column_index: int) -> None:
    row_count = len(table.rows)
    column_count = len(table.columns)
    if row_index >= row_count:
        raise IndexError(f"row_index {row_index} is out of range for table with {row_count} rows.")
    if column_index >= column_count:
        raise IndexError(
            f"column_index {column_index} is out of range for table with {column_count} columns."
        )


def _normalize_row_values(operation: dict, *, expected_columns: int) -> list[str]:
    values = operation.get("values")
    if isinstance(values, Sequence) and not isinstance(values, (str, bytes, bytearray)):
        normalized_values = [str(value) for value in values]
    else:
        raise ValueError("append_table_row requires a sequence of values under 'values'.")

    if len(normalized_values) != expected_columns:
        raise ValueError(
            f"append_table_row requires exactly {expected_columns} values for the target table."
        )
    return normalized_values


def _append_table_row(table, row_values: Sequence[str]) -> None:
    new_row = deepcopy(table._tbl.tr_lst[-1])
    table._tbl.append(new_row)
    new_row_index = len(table.rows) - 1
    for column_index, value in enumerate(row_values):
        _set_text_frame_text(table.cell(new_row_index, column_index).text_frame, value)


def _resolve_new_text(operation: dict) -> str:
    for key in ("new_text", "text", "replacement_text"):
        if key in operation:
            value = operation[key]
            if not isinstance(value, str):
                raise TypeError(f"{key} must be a string.")
            return value
    raise ValueError("The operation is missing replacement text. Use new_text, text, or replacement_text.")


def _slide_number_from_operation(operation: dict) -> int:
    if "slide_number" in operation and operation["slide_number"] is not None:
        return _require_positive_int(operation, "slide_number")
    return _require_non_negative_int(operation, "slide_index") + 1


def _require_positive_int(payload: dict, key: str) -> int:
    value = payload.get(key)
    if not isinstance(value, int):
        raise TypeError(f"{key} must be an integer.")
    if value <= 0:
        raise ValueError(f"{key} must be greater than 0.")
    return value


def _require_non_negative_int(payload: dict, key: str) -> int:
    value = payload.get(key)
    if not isinstance(value, int):
        raise TypeError(f"{key} must be an integer.")
    if value < 0:
        raise ValueError(f"{key} must be greater than or equal to 0.")
    return value


def _set_text_frame_text(text_frame, text: str) -> None:
    text_frame.clear()
    text_frame.text = text


def _extract_text_frame_text(text_frame) -> str:
    return text_frame.text


def _shape_type_name(shape) -> str:
    shape_type = getattr(shape, "shape_type", None)
    if shape_type is None:
        return "UNKNOWN"
    return getattr(shape_type, "name", str(shape_type))


def _validate_presentation_path(file_path: str | PathLike[str]) -> Path:
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"Presentation file '{path}' does not exist.")
    if path.is_symlink() or not path.is_file():
        raise ValueError(f"Presentation file '{path}' is not a regular file.")
    if path.suffix.lower() != _SUPPORTED_EXTENSION:
        display_extension = path.suffix.lower().lstrip(".") or "<none>"
        raise ValueError(
            f"Presentation file '{path}' has unsupported extension '{display_extension}'. Supported extensions: {_SUPPORTED_EXTENSION_TEXT}."
        )
    return path


def _paths_refer_to_same_location(left: Path, right: Path) -> bool:
    return left.resolve() == right.resolve(strict=False)


def _deduplicate_messages(messages: Iterable[str]) -> list[str]:
    deduplicated: list[str] = []
    seen: set[str] = set()
    for message in messages:
        if message not in seen:
            deduplicated.append(message)
            seen.add(message)
    return deduplicated
