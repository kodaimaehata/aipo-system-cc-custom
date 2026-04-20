"""Anonymization transforms for SG3 V1.

This module implements the shared-runtime callable ``transform(detected, policy)``
for the approved V1 extensions ``xlsx``, ``xlsm``, ``docx``, ``pptx``, and
``pdf``.

Supported policy vocabulary is intentionally explicit and wrapper-friendly:
- structural categories (``comments``, ``notes``, ``headers``, ``footers``)
  accept ``remove`` / ``clear``, ``replace``, and ``skip``
- ``metadata`` accepts ``clear`` / ``remove`` and ``skip`` / ``report_only``
- ``images`` accepts run-level modes ``replace``, ``mask``, ``remove``,
  ``skip``, and ``report_only``
- text replacement requires one of ``replacement_text``, ``replacement``, or
  ``text``
- image replacement requires ``replacement_path`` (or a compatible alias)
- image masking stays Python-only by generating a solid-color replacement asset
  with Pillow and then routing through ``office_automation.common.images``

Result model established for later SG3 tasks:
- returned value is a list of per-file dictionaries
- each file record includes ``file_path``, ``extension``, ``output_path``,
  ``status``, ``actions``, ``warnings``, and ``manual_review_items``
- action status vocabulary is ``applied``, ``skipped``,
  ``manual_review_required``, and ``error``
- file status vocabulary is ``success``, ``partial_success``,
  ``manual_review_required``, and ``error``
- T005 extends action records with optional ``details`` so wrapper/report logic
  can distinguish clean success from warning-bearing partial cleanup

Important V1 non-scope / conservative behaviors:
- Word comment deletion/rewrite is not claimed safe here, so DOCX comment
  findings become structured manual-review items instead of silent success.
- PowerPoint comment, header, and footer surfaces remain manual-review-only when
  detection already surfaced them as such.
- PDF header/footer findings remain manual-review-only when detection marks them
  as heuristic-sensitive.
- Embedded objects, OLE containers, attachments, and OCR-dependent cleanup are
  excluded from image anonymization scope and must surface as warnings/manual
  review instead of false success.
- True in-place image removal is only used when a proved-safe path exists. In
  this V1 runtime, ``remove`` conservatively downgrades to masking unless later
  format-specific support is added.
"""

from __future__ import annotations

import json
import tempfile
from collections import OrderedDict
from pathlib import Path

import fitz
from docx import Document
from openpyxl import load_workbook
from openpyxl.comments import Comment
from PIL import Image, ImageColor
from pptx import Presentation

from office_automation.common.images import extract_images, replace_image
from office_automation.common.metadata import clear_metadata, read_metadata

__all__ = ["transform"]

_SUPPORTED_EXTENSIONS = {"xlsx", "xlsm", "docx", "pptx", "pdf"}
_STRUCTURAL_CATEGORIES = {"comments", "notes", "headers", "footers"}
_BODY_TEXT_CATEGORIES = {"body_text"}
_METADATA_CATEGORIES = {"metadata"}
_IMAGE_CATEGORIES = {"images"}
_ALLOWED_POLICY_ACTIONS = {"remove", "clear", "replace", "skip"}
_ALLOWED_METADATA_ACTIONS = {"clear", "remove", "skip", "report_only"}
_ALLOWED_IMAGE_ACTIONS = {"replace", "mask", "remove", "skip", "report_only"}
_DEFAULT_IMAGE_MASK_COLOR = "#000000"
_EXCLUDED_IMAGE_CONTAINER_MARKERS = (
    "attachment",
    "embedded",
    "embedded_object",
    "ole",
    "ocr",
)
_DOCX_HEADER_VARIANTS = (
    ("default", "header"),
    ("first", "first_page_header"),
    ("even", "even_page_header"),
)
_DOCX_FOOTER_VARIANTS = (
    ("default", "footer"),
    ("first", "first_page_footer"),
    ("even", "even_page_footer"),
)
_EXCEL_HEADER_FOOTER_ATTRS = {
    ("headers", "odd"): "oddHeader",
    ("headers", "even"): "evenHeader",
    ("headers", "first"): "firstHeader",
    ("footers", "odd"): "oddFooter",
    ("footers", "even"): "evenFooter",
    ("footers", "first"): "firstFooter",
}


class _FileMutationState:
    def __init__(self, file_path: Path) -> None:
        self.file_path = file_path
        self.changed = False


class _ResultCollector:
    def __init__(self, file_result: dict) -> None:
        self.file_result = file_result

    def add_action(self, action: dict) -> None:
        action.setdefault("warnings", [])
        self.file_result["actions"].append(_sorted_copy(action))
        for warning in action["warnings"]:
            self.add_warning(warning)
        manual_review_reason = action.get("manual_review_reason")
        if manual_review_reason:
            item = {
                "finding_id": action.get("finding_id"),
                "category": action.get("category"),
                "location": _sorted_copy(action.get("location") or {}),
                "requested_action": action.get("requested_action"),
                "reason": manual_review_reason,
            }
            self.file_result["manual_review_items"].append(item)

    def add_warning(self, warning: str) -> None:
        if warning and warning not in self.file_result["warnings"]:
            self.file_result["warnings"].append(warning)


class _Plan:
    def __init__(self, *, requested_action: str, replacement_text: str | None = None, policy_error: str | None = None) -> None:
        self.requested_action = requested_action
        self.replacement_text = replacement_text
        self.policy_error = policy_error


class _ImagePlan:
    def __init__(
        self,
        *,
        requested_action: str,
        replacement_path: str | None = None,
        mask_color: str = _DEFAULT_IMAGE_MASK_COLOR,
        policy_error: str | None = None,
    ) -> None:
        self.requested_action = requested_action
        self.replacement_path = replacement_path
        self.mask_color = mask_color
        self.policy_error = policy_error


class _BodyTextPlan:
    def __init__(
        self,
        *,
        requested_action: str,
        candidate_id: str | None,
        decision: str,
        replacement_text: str | None = None,
        policy_error: str | None = None,
        message: str | None = None,
        warnings: list[str] | None = None,
        manual_review_reason: str | None = None,
    ) -> None:
        self.requested_action = requested_action
        self.candidate_id = candidate_id
        self.decision = decision
        self.replacement_text = replacement_text
        self.policy_error = policy_error
        self.message = message
        self.warnings = list(warnings or [])
        self.manual_review_reason = manual_review_reason



def transform(detected: list[dict], policy: dict) -> list[dict]:
    """Apply structural anonymization policy to detected findings."""
    grouped = _group_findings(detected)
    results: list[dict] = []
    for file_result, findings in grouped:
        collector = _ResultCollector(file_result)
        file_path = Path(file_result["file_path"])
        extension = file_result["extension"]

        if extension not in _SUPPORTED_EXTENSIONS:
            for finding in findings:
                collector.add_action(
                    _base_action(
                        finding,
                        requested_action="skip",
                        applied_action="none",
                        status="error",
                        message=f"Unsupported extension '{extension}' for transform runtime.",
                    )
                )
            file_result["status"] = _finalize_file_status(file_result)
            results.append(_sorted_copy(file_result))
            continue

        if not file_path.exists():
            for finding in findings:
                collector.add_action(
                    _base_action(
                        finding,
                        requested_action="skip",
                        applied_action="none",
                        status="error",
                        message=f"Source file '{file_path}' does not exist.",
                    )
                )
            file_result["status"] = _finalize_file_status(file_result)
            results.append(_sorted_copy(file_result))
            continue

        try:
            _transform_file(file_path, extension, findings, policy, collector)
        except Exception as exc:  # pragma: no cover - exercised by future fault-injection tests
            collector.add_warning(f"Fatal transform failure for '{file_path}': {exc}")
            if not file_result["actions"]:
                file_result["actions"].append(
                    _sorted_copy(
                        {
                            "finding_id": None,
                            "category": None,
                            "location": {},
                            "payload": {},
                            "confidence": None,
                            "requested_action": "skip",
                            "applied_action": "none",
                            "status": "error",
                            "message": f"Fatal transform failure for '{file_path}': {exc}",
                            "warnings": [],
                            "manual_review_reason": None,
                            "details": {},
                        }
                    )
                )

        file_result["status"] = _finalize_file_status(file_result)
        results.append(_sorted_copy(file_result))
    return results



def _transform_file(file_path: Path, extension: str, findings: list[dict], policy: dict, collector: _ResultCollector) -> None:
    pre_save_findings = [
        finding
        for finding in findings
        if str(finding.get("category") or "") in _STRUCTURAL_CATEGORIES | _BODY_TEXT_CATEGORIES
    ]
    post_save_findings = [
        finding
        for finding in findings
        if str(finding.get("category") or "") not in _STRUCTURAL_CATEGORIES | _BODY_TEXT_CATEGORIES
    ]

    if pre_save_findings:
        if extension in {"xlsx", "xlsm"}:
            workbook = load_workbook(file_path, keep_vba=extension == "xlsm")
            state = _FileMutationState(file_path)
            try:
                for finding in pre_save_findings:
                    _apply_excel_finding(workbook, finding, policy, collector, state)
                if state.changed:
                    _save_openpyxl_workbook(workbook, file_path)
            finally:
                _close_excel_workbook(workbook)
        elif extension == "docx":
            document = Document(file_path)
            state = _FileMutationState(file_path)
            for finding in pre_save_findings:
                _apply_docx_finding(document, finding, policy, collector, state)
            if state.changed:
                _save_docx_document(document, file_path)
        elif extension == "pptx":
            presentation = Presentation(file_path)
            state = _FileMutationState(file_path)
            for finding in pre_save_findings:
                _apply_pptx_finding(presentation, finding, policy, collector, state)
            if state.changed:
                _save_presentation(presentation, file_path)
        elif extension == "pdf":
            document = fitz.open(file_path)
            try:
                if getattr(document, "needs_pass", False):
                    raise NotImplementedError(f"Encrypted PDF files are outside transform V1 scope: '{file_path}'.")
                state = _FileMutationState(file_path)
                for finding in pre_save_findings:
                    _apply_pdf_finding(document, finding, policy, collector, state)
                if state.changed:
                    _save_pdf_document(document, file_path)
            finally:
                document.close()
        else:
            raise ValueError(f"Unsupported extension '{extension}' for '{file_path}'.")

    for finding in post_save_findings:
        category = str(finding.get("category") or "")
        if category == "metadata":
            _apply_metadata_finding(file_path, finding, policy, collector)
            continue
        if category == "images":
            _apply_image_finding(file_path, extension, finding, policy, collector)
            continue
        collector.add_action(_deferred_scope_action(finding))



def _apply_excel_finding(workbook, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    category = str(finding.get("category") or "")
    if category == "body_text":
        _apply_excel_body_text_finding(workbook, finding, policy, collector, state)
        return

    plan = _plan_for_finding(finding, policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_STRUCTURAL_CATEGORIES):
        collector.add_action(action)
        return

    if category == "notes":
        worksheet = workbook[_required_location_value(finding, "sheet")]
        cell = worksheet[_required_location_value(finding, "cell")]
        if plan.requested_action in {"remove", "clear"}:
            cell.comment = None
            state.changed = True
            collector.add_action(
                _base_action(
                    finding,
                    requested_action=plan.requested_action,
                    applied_action="remove",
                    status="applied",
                    message="Removed Excel cell comment/note.",
                )
            )
            return

        replacement_text = plan.replacement_text or ""
        cell.comment = Comment(replacement_text, " ")
        state.changed = True
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="replace",
                status="applied",
                message="Replaced Excel cell comment/note text and blanked the author field.",
            )
        )
        return

    if category in {"headers", "footers"}:
        worksheet = workbook[_required_location_value(finding, "sheet")]
        variant = _required_location_value(finding, "variant")
        part = _required_location_value(finding, "part")
        section_attr = _EXCEL_HEADER_FOOTER_ATTRS[(category, variant)]
        text_part = getattr(getattr(worksheet, section_attr), part)
        text_part.text = "" if plan.requested_action in {"remove", "clear"} else (plan.replacement_text or "")
        state.changed = True
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="clear" if plan.requested_action in {"remove", "clear"} else "replace",
                status="applied",
                message=f"Updated Excel {category[:-1]} text.",
            )
        )
        return

    collector.add_action(_deferred_scope_action(finding))



def _apply_docx_finding(document: Document, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    category = str(finding.get("category") or "")
    if category == "body_text":
        _apply_docx_body_text_finding(document, finding, policy, collector, state)
        return

    plan = _plan_for_finding(finding, policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_STRUCTURAL_CATEGORIES):
        collector.add_action(action)
        return

    if category == "comments":
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="DOCX comment rewrite/removal is not claimed safe in this V1 runtime.",
                manual_review_reason="DOCX comments require manual review because this runtime does not yet target Word comment XML with a proved-safe deletion path.",
            )
        )
        return

    if category in {"headers", "footers"}:
        container = _docx_container_by_part(document, category=category, part_name=_required_location_value(finding, "part"))
        replacement_text = "" if plan.requested_action in {"remove", "clear"} else (plan.replacement_text or "")
        _rewrite_docx_story_container(container, replacement_text)
        state.changed = True
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="clear" if plan.requested_action in {"remove", "clear"} else "replace",
                status="applied",
                message=f"Updated DOCX {category[:-1]} text.",
            )
        )
        return

    collector.add_action(_deferred_scope_action(finding))



def _apply_pptx_finding(presentation: Presentation, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    category = str(finding.get("category") or "")
    if category == "body_text":
        _apply_pptx_body_text_finding(presentation, finding, policy, collector, state)
        return

    plan = _plan_for_finding(finding, policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_STRUCTURAL_CATEGORIES):
        collector.add_action(action)
        return

    if category == "notes":
        slide_number = int(_required_location_value(finding, "slide_number"))
        shape_id = int(_required_location_value(finding, "shape_id"))
        slide = presentation.slides[slide_number - 1]
        target_shape = next(
            (shape for shape in slide.notes_slide.shapes if getattr(shape, "shape_id", None) == shape_id),
            None,
        )
        if target_shape is None or getattr(target_shape, "text_frame", None) is None:
            raise ValueError(f"Could not locate PowerPoint notes shape {shape_id} on slide {slide_number}.")
        target_shape.text_frame.text = "" if plan.requested_action in {"remove", "clear"} else (plan.replacement_text or "")
        state.changed = True
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="clear" if plan.requested_action in {"remove", "clear"} else "replace",
                status="applied",
                message="Updated PowerPoint speaker notes text.",
            )
        )
        return

    collector.add_action(_deferred_scope_action(finding))



def _apply_pdf_finding(document: fitz.Document, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    category = str(finding.get("category") or "")
    if category == "body_text":
        _apply_pdf_body_text_finding(document, finding, policy, collector, state)
        return

    plan = _plan_for_finding(finding, policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_STRUCTURAL_CATEGORIES):
        collector.add_action(action)
        return

    if category == "comments":
        page_number = int(_required_location_value(finding, "page_number"))
        xref = int(_required_location_value(finding, "xref"))
        page = document.load_page(page_number - 1)
        annotation = next((annot for annot in page.annots() or [] if int(annot.xref) == xref), None)
        if annotation is None:
            raise ValueError(f"Could not locate PDF annotation xref {xref} on page {page_number}.")
        if plan.requested_action in {"remove", "clear"}:
            page.delete_annot(annotation)
            state.changed = True
            collector.add_action(
                _base_action(
                    finding,
                    requested_action=plan.requested_action,
                    applied_action="remove",
                    status="applied",
                    message="Removed PDF annotation/comment.",
                )
            )
            return

        annotation.set_info(content=plan.replacement_text or "", title=" ", subject=" ")
        annotation.update()
        state.changed = True
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="replace",
                status="applied",
                message="Replaced PDF annotation content and blanked annotation author/subject fields.",
            )
        )
        return

    collector.add_action(_deferred_scope_action(finding))



def _apply_excel_body_text_finding(workbook, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    plan = _body_text_plan_for_finding(finding, policy)
    if action := _precomputed_body_text_action_if_needed(finding, plan):
        collector.add_action(action)
        return

    try:
        worksheet = workbook[_required_location_value(finding, "sheet")]
        cell = worksheet[_required_location_value(finding, "cell")]
    except Exception as exc:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Excel body-text locator could not be resolved safely.",
                manual_review_reason=f"Excel cell locator could not be resolved safely: {exc}",
                details=_body_text_details(finding, plan=plan, locator_validated=False),
            )
        )
        return

    if cell.data_type == "f":
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Excel formula cells are left for manual review.",
                manual_review_reason="Excel body-text rewrite is restricted to plain-string cells; formula cells are manual-review-only in this V1 runtime.",
                details=_body_text_details(finding, plan=plan, locator_validated=False),
            )
        )
        return

    if not isinstance(cell.value, str):
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Excel body-text rewrite requires a current plain-string cell value.",
                manual_review_reason="Excel body-text rewrite is restricted to plain-string cells that still expose the detected text directly.",
                details=_body_text_details(finding, plan=plan, locator_validated=False),
            )
        )
        return

    locator_validated, matched_text_before, reason = _validate_body_text_locator(cell.value, finding)
    if not locator_validated:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Excel body-text locator no longer matches the current cell text.",
                manual_review_reason=reason,
                details=_body_text_details(
                    finding,
                    plan=plan,
                    locator_validated=False,
                    matched_text_before=matched_text_before,
                ),
            )
        )
        return

    cell.value = _replace_body_text_span(cell.value, finding, replacement_text=plan.replacement_text or "")
    state.changed = True
    collector.add_action(
        _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="replace",
            status="applied",
            message="Replaced approved Excel body text through a validated cell locator.",
            details=_body_text_details(
                finding,
                plan=plan,
                locator_validated=True,
                matched_text_before=matched_text_before,
            ),
        )
    )



def _apply_docx_body_text_finding(document: Document, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    plan = _body_text_plan_for_finding(finding, policy)
    if action := _precomputed_body_text_action_if_needed(finding, plan):
        collector.add_action(action)
        return

    try:
        paragraph = _docx_body_text_paragraph(document, finding)
    except Exception as exc:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="DOCX body-text locator could not be resolved safely.",
                manual_review_reason=f"DOCX body-text locator could not be resolved safely: {exc}",
                details=_body_text_details(finding, plan=plan, locator_validated=False),
            )
        )
        return

    locator_validated, matched_text_before, reason = _validate_body_text_locator(paragraph.text, finding)
    if not locator_validated:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="DOCX body-text locator no longer matches the current paragraph text.",
                manual_review_reason=reason,
                details=_body_text_details(
                    finding,
                    plan=plan,
                    locator_validated=False,
                    matched_text_before=matched_text_before,
                ),
            )
        )
        return

    if not _paragraph_supports_safe_inline_rewrite(paragraph):
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="DOCX body-text rewrite was downgraded to manual review.",
                manual_review_reason=(
                    "DOCX body-text rewrite is restricted to single-run paragraphs/table cells in this V1 runtime to avoid formatting-risky cross-run edits."
                ),
                details=_body_text_details(
                    finding,
                    plan=plan,
                    locator_validated=True,
                    matched_text_before=matched_text_before,
                ),
            )
        )
        return

    paragraph.runs[0].text = _replace_body_text_span(paragraph.text, finding, replacement_text=plan.replacement_text or "")
    state.changed = True
    collector.add_action(
        _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="replace",
            status="applied",
            message="Replaced approved DOCX body text through a validated paragraph locator.",
            details=_body_text_details(
                finding,
                plan=plan,
                locator_validated=True,
                matched_text_before=matched_text_before,
            ),
        )
    )



def _apply_pptx_body_text_finding(presentation: Presentation, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    plan = _body_text_plan_for_finding(finding, policy)
    if action := _precomputed_body_text_action_if_needed(finding, plan):
        collector.add_action(action)
        return

    try:
        paragraph = _pptx_body_text_paragraph(presentation, finding)
    except Exception as exc:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="PowerPoint body-text locator could not be resolved safely.",
                manual_review_reason=f"PowerPoint body-text locator could not be resolved safely: {exc}",
                details=_body_text_details(finding, plan=plan, locator_validated=False),
            )
        )
        return

    locator_validated, matched_text_before, reason = _validate_body_text_locator(paragraph.text, finding)
    if not locator_validated:
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="PowerPoint body-text locator no longer matches the current paragraph text.",
                manual_review_reason=reason,
                details=_body_text_details(
                    finding,
                    plan=plan,
                    locator_validated=False,
                    matched_text_before=matched_text_before,
                ),
            )
        )
        return

    if not _paragraph_supports_safe_inline_rewrite(paragraph):
        collector.add_action(
            _body_text_action(
                finding,
                candidate_id=plan.candidate_id,
                decision=plan.decision,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="PowerPoint body-text rewrite was downgraded to manual review.",
                manual_review_reason=(
                    "PowerPoint body-text rewrite is restricted to single-run paragraphs/table cells in this V1 runtime to avoid formatting-risky cross-run edits."
                ),
                details=_body_text_details(
                    finding,
                    plan=plan,
                    locator_validated=True,
                    matched_text_before=matched_text_before,
                ),
            )
        )
        return

    paragraph.runs[0].text = _replace_body_text_span(paragraph.text, finding, replacement_text=plan.replacement_text or "")
    state.changed = True
    collector.add_action(
        _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="replace",
            status="applied",
            message="Replaced approved PowerPoint body text through a validated paragraph locator.",
            details=_body_text_details(
                finding,
                plan=plan,
                locator_validated=True,
                matched_text_before=matched_text_before,
            ),
        )
    )



def _apply_pdf_body_text_finding(document: fitz.Document, finding: dict, policy: dict, collector: _ResultCollector, state: _FileMutationState) -> None:
    del document, state
    plan = _body_text_plan_for_finding(finding, policy)
    if action := _precomputed_body_text_action_if_needed(finding, plan):
        collector.add_action(action)
        return

    collector.add_action(
        _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="manual_review",
            status="manual_review_required",
            message="PDF body-text rewrite remains conservative/manual-review-heavy in this V1 runtime.",
            manual_review_reason=(
                "PDF body-text rewrite remains conservative/manual-review-heavy unless an exact stable text-layer mutation path is proven safe for the current finding."
            ),
            details=_body_text_details(finding, plan=plan, locator_validated=False),
        )
    )



def _body_text_plan_for_finding(finding: dict, policy: dict) -> _BodyTextPlan:
    category_policy = policy.get("body_text")
    candidate_id = _body_text_candidate_id(finding, category_policy if isinstance(category_policy, dict) else {})
    if not isinstance(category_policy, dict):
        return _BodyTextPlan(
            requested_action="skip",
            candidate_id=candidate_id,
            decision="undecided",
            policy_error="Policy entry for category 'body_text' must be a dict.",
        )

    candidate_decisions = category_policy.get("candidate_decisions")
    decision_payload = candidate_decisions.get(candidate_id) if isinstance(candidate_decisions, dict) and candidate_id else None
    decision = "undecided"
    if isinstance(decision_payload, dict):
        decision = str(decision_payload.get("decision", "undecided")).strip().lower() or "undecided"

    enabled = category_policy.get("enabled", True)
    mode = str(category_policy.get("mode", "preview_only")).strip().lower() or "preview_only"
    replacement_text = _body_text_replacement_text(finding, category_policy, candidate_id=candidate_id, decision_payload=decision_payload)

    if enabled is False or mode != "apply_confirmed":
        return _BodyTextPlan(
            requested_action="skip",
            candidate_id=candidate_id,
            decision=decision,
            replacement_text=replacement_text,
            message=f"Body-text policy mode '{mode}' does not apply confirmed mutations.",
            warnings=[
                f"Body-text finding was left in place because mode '{mode}' does not permit confirmed runtime mutation."
            ],
        )

    if candidate_id is None:
        return _BodyTextPlan(
            requested_action="skip",
            candidate_id=None,
            decision="undecided",
            replacement_text=replacement_text,
            message="Body-text finding could not be matched back to a confirmed candidate decision.",
            warnings=[
                "Body-text finding was left in place because the runtime could not resolve its candidate_id from the SG5 confirmation payload."
            ],
        )

    if not isinstance(decision_payload, dict):
        return _BodyTextPlan(
            requested_action="skip",
            candidate_id=candidate_id,
            decision="undecided",
            replacement_text=replacement_text,
            message="Confirmed body-text policy did not contain a candidate decision for this finding.",
            warnings=[
                "Body-text finding was left in place because candidate_decisions did not contain a decision for this candidate_id."
            ],
        )

    if decision == "approved":
        finding_id = str(finding.get("finding_id"))
        transformable_ids = _string_set(decision_payload.get("transformable_finding_ids"))
        approved_finding_ids = _string_set(category_policy.get("approved_finding_ids"))
        if transformable_ids and finding_id not in transformable_ids:
            return _BodyTextPlan(
                requested_action="skip",
                candidate_id=candidate_id,
                decision=decision,
                replacement_text=replacement_text,
                message="Approved candidate did not authorize this specific body-text finding for mutation.",
                warnings=[
                    "Body-text finding was left unchanged because its finding_id was not listed in the candidate's transformable_finding_ids."
                ],
            )
        if approved_finding_ids and finding_id not in approved_finding_ids:
            return _BodyTextPlan(
                requested_action="skip",
                candidate_id=candidate_id,
                decision=decision,
                replacement_text=replacement_text,
                message="Approved candidate data for this body-text finding is stale.",
                warnings=[
                    "Body-text finding was left unchanged because its finding_id was not listed in approved_finding_ids for this confirmed run."
                ],
            )
        if replacement_text is None:
            return _BodyTextPlan(
                requested_action="replace",
                candidate_id=candidate_id,
                decision=decision,
                policy_error="Approved body-text candidates require replacement text before transform can run.",
            )
        return _BodyTextPlan(
            requested_action="replace",
            candidate_id=candidate_id,
            decision=decision,
            replacement_text=replacement_text,
        )

    if decision == "manual_review" or bool(decision_payload.get("manual_review_required", False)):
        return _BodyTextPlan(
            requested_action="replace",
            candidate_id=candidate_id,
            decision=decision,
            replacement_text=replacement_text,
            manual_review_reason="Candidate was kept manual-review-only in the SG5 confirmation payload.",
        )

    return _BodyTextPlan(
        requested_action="skip",
        candidate_id=candidate_id,
        decision=decision,
        replacement_text=replacement_text,
        message=f"Candidate decision '{decision}' does not permit automatic body-text mutation.",
        warnings=[f"Body-text finding was left in place because the candidate decision resolved to '{decision}'."],
    )



def _precomputed_body_text_action_if_needed(finding: dict, plan: _BodyTextPlan) -> dict | None:
    if plan.policy_error:
        return _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="none",
            status="error",
            message=plan.policy_error,
            details=_body_text_details(finding, plan=plan, locator_validated=False),
        )

    manual_review_reason = str(finding.get("manual_review_reason") or "").strip() or plan.manual_review_reason
    if manual_review_reason:
        return _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action=plan.requested_action,
            applied_action="manual_review",
            status="manual_review_required",
            message=manual_review_reason,
            manual_review_reason=manual_review_reason,
            warnings=plan.warnings,
            details=_body_text_details(finding, plan=plan, locator_validated=False),
        )

    if plan.requested_action == "skip":
        return _body_text_action(
            finding,
            candidate_id=plan.candidate_id,
            decision=plan.decision,
            requested_action="skip",
            applied_action="none",
            status="skipped",
            message=plan.message or "Body-text finding was left unchanged.",
            warnings=plan.warnings,
            details=_body_text_details(finding, plan=plan, locator_validated=False),
        )

    return None



def _body_text_action(
    finding: dict,
    *,
    candidate_id: str | None,
    decision: str,
    requested_action: str,
    applied_action: str,
    status: str,
    message: str,
    warnings: list[str] | None = None,
    manual_review_reason: str | None = None,
    details: dict | None = None,
) -> dict:
    action = _base_action(
        finding,
        requested_action=requested_action,
        applied_action=applied_action,
        status=status,
        message=message,
        warnings=warnings,
        manual_review_reason=manual_review_reason,
        details=details,
    )
    action["candidate_id"] = candidate_id
    action["decision"] = decision
    return action



def _body_text_details(
    finding: dict,
    *,
    plan: _BodyTextPlan,
    locator_validated: bool,
    matched_text_before: str | None = None,
) -> dict:
    return {
        "decision": plan.decision,
        "locator_validated": locator_validated,
        "matched_text_before": matched_text_before,
        "replacement_text": plan.replacement_text,
        "surface_kind": _body_text_surface_kind(finding),
    }



def _body_text_candidate_id(finding: dict, category_policy: dict) -> str | None:
    explicit_candidate_id = finding.get("candidate_id") or (finding.get("payload") or {}).get("candidate_id")
    if isinstance(explicit_candidate_id, str) and explicit_candidate_id.strip():
        return explicit_candidate_id.strip()

    finding_id = str(finding.get("finding_id") or "")
    candidate_summary = category_policy.get("candidate_summary") if isinstance(category_policy, dict) else None
    if isinstance(candidate_summary, dict):
        finding_to_candidate = candidate_summary.get("finding_to_candidate")
        if isinstance(finding_to_candidate, dict):
            candidate_id = finding_to_candidate.get(finding_id)
            if isinstance(candidate_id, str) and candidate_id.strip():
                return candidate_id.strip()

    candidate_decisions = category_policy.get("candidate_decisions") if isinstance(category_policy, dict) else None
    if isinstance(candidate_decisions, dict):
        for candidate_id, decision_payload in candidate_decisions.items():
            if not isinstance(candidate_id, str) or not isinstance(decision_payload, dict):
                continue
            if finding_id in _string_set(decision_payload.get("transformable_finding_ids")):
                return candidate_id
            if finding_id in _string_set(decision_payload.get("non_transformable_finding_ids")):
                return candidate_id
    return None



def _body_text_replacement_text(
    finding: dict,
    category_policy: dict,
    *,
    candidate_id: str | None,
    decision_payload: dict | None,
) -> str | None:
    if isinstance(decision_payload, dict):
        replacement_text = decision_payload.get("replacement_text")
        if replacement_text not in (None, ""):
            return str(replacement_text)

    if candidate_id and isinstance(category_policy.get("replacement_overrides"), dict):
        override = category_policy["replacement_overrides"].get(candidate_id)
        if override not in (None, ""):
            return str(override)

    payload = finding.get("payload") or {}
    replacement_map = category_policy.get("replacement_map") if isinstance(category_policy.get("replacement_map"), dict) else {}
    for key in (
        payload.get("matched_text"),
        payload.get("normalized_text"),
    ):
        if isinstance(key, str) and key in replacement_map and replacement_map[key] not in (None, ""):
            return str(replacement_map[key])

    for key in (
        payload.get("suggested_replacement"),
        payload.get("recommended_replacement"),
        payload.get("default_replacement_text"),
        category_policy.get("default_replacement_text"),
        category_policy.get("replacement_text"),
    ):
        if key not in (None, ""):
            return str(key)
    return None



def _validate_body_text_locator(container_text: object, finding: dict) -> tuple[bool, str | None, str]:
    if not isinstance(container_text, str):
        return False, None, "Current text container is not a plain string, so the SG5 locator cannot be validated safely."

    matched_text = str((finding.get("payload") or {}).get("matched_text") or "")
    location = finding.get("location") or {}
    start = location.get("match_start")
    end = location.get("match_end")
    if isinstance(start, bool) or isinstance(end, bool) or not isinstance(start, int) or not isinstance(end, int):
        return False, None, "Body-text finding is missing integer match_start/match_end locator values."
    if start < 0 or end < start or end > len(container_text):
        return False, None, "Body-text locator offsets are no longer valid for the current text container."

    current_slice = container_text[start:end]
    if current_slice != matched_text:
        return (
            False,
            current_slice,
            "Body-text container changed since detection, so the stored locator no longer matches the current text safely.",
        )
    return True, current_slice, ""



def _replace_body_text_span(container_text: str, finding: dict, *, replacement_text: str) -> str:
    location = finding.get("location") or {}
    start = int(location["match_start"])
    end = int(location["match_end"])
    return f"{container_text[:start]}{replacement_text}{container_text[end:]}"



def _docx_body_text_paragraph(document: Document, finding: dict):
    location = finding.get("location") or {}
    surface = _required_location_value(finding, "surface")
    paragraph_index = int(_required_location_value(finding, "paragraph_index"))
    if surface == "paragraph":
        return document.paragraphs[paragraph_index]
    if surface == "table_cell":
        table = document.tables[int(_required_location_value(finding, "table_index"))]
        cell = table.rows[int(_required_location_value(finding, "row_index"))].cells[
            int(_required_location_value(finding, "column_index"))
        ]
        return cell.paragraphs[paragraph_index]
    raise ValueError(f"Unsupported DOCX body-text surface '{location.get('surface')}'.")



def _pptx_body_text_paragraph(presentation: Presentation, finding: dict):
    slide_number = int(_required_location_value(finding, "slide_number"))
    shape_id = int(_required_location_value(finding, "shape_id"))
    paragraph_index = int(_required_location_value(finding, "paragraph_index"))
    slide = presentation.slides[slide_number - 1]
    shape = next((item for item in slide.shapes if getattr(item, "shape_id", None) == shape_id), None)
    if shape is None:
        raise ValueError(f"Could not locate PowerPoint shape {shape_id} on slide {slide_number}.")

    location = finding.get("location") or {}
    if "row_index" in location or "column_index" in location:
        if not getattr(shape, "has_table", False):
            raise ValueError("Resolved PowerPoint body-text locator expected a table shape, but the shape no longer exposes a table.")
        cell = shape.table.rows[int(_required_location_value(finding, "row_index"))].cells[
            int(_required_location_value(finding, "column_index"))
        ]
        return cell.text_frame.paragraphs[paragraph_index]

    if not getattr(shape, "has_text_frame", False):
        raise ValueError("Resolved PowerPoint body-text locator expected a text frame, but the shape no longer exposes one.")
    return shape.text_frame.paragraphs[paragraph_index]



def _paragraph_supports_safe_inline_rewrite(paragraph) -> bool:
    runs = list(getattr(paragraph, "runs", []))
    return len(runs) == 1 and getattr(runs[0], "text", None) == getattr(paragraph, "text", None)



def _body_text_surface_kind(finding: dict) -> str | None:
    payload = finding.get("payload") or {}
    value = payload.get("surface_kind", payload.get("surface_type"))
    return str(value) if value not in (None, "") else None



def _string_set(value: object) -> set[str]:
    if not isinstance(value, list):
        return set()
    return {str(item) for item in value if str(item).strip()}



def _apply_metadata_finding(file_path: Path, finding: dict, policy: dict, collector: _ResultCollector) -> None:
    category = str(finding.get("category") or "")
    plan = _metadata_plan_for_finding(policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_METADATA_CATEGORIES):
        collector.add_action(action)
        return

    before = read_metadata(file_path)
    before_fields = _non_empty_mapping(before.get("fields") or {})
    before_warnings = list(before.get("warnings") or [])
    if not before_fields:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="skip",
                status="skipped",
                message="No document-level metadata remained by transform time.",
                warnings=["Metadata finding had no remaining non-empty fields when transform ran."],
                details={
                    "fields_after": {},
                    "fields_before": {},
                    "helper_warnings_before": before_warnings,
                },
            )
        )
        return

    try:
        clear_metadata(file_path)
    except NotImplementedError as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message=str(exc),
                manual_review_reason=str(exc),
                details={
                    "fields_before": before_fields,
                    "helper_warnings_before": before_warnings,
                },
            )
        )
        return
    except Exception as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="none",
                status="error",
                message=str(exc),
                details={
                    "fields_before": before_fields,
                    "helper_warnings_before": before_warnings,
                },
            )
        )
        return

    after = read_metadata(file_path)
    after_fields = _non_empty_mapping(after.get("fields") or {})
    after_warnings = list(after.get("warnings") or [])
    cleared_fields = sorted(set(before_fields) - set(after_fields))
    residual_fields = sorted(after_fields)
    warnings = list(dict.fromkeys(before_warnings + after_warnings))

    if residual_fields:
        warnings.append(
            "Metadata cleanup left residual fields that require manual review: "
            + ", ".join(residual_fields)
            + "."
        )
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="clear",
                status="applied",
                message="Cleared supported metadata fields, but residual metadata still requires manual review.",
                warnings=warnings,
                manual_review_reason=(
                    "Residual metadata fields remain after helper-driven cleanup: "
                    + ", ".join(residual_fields)
                    + "."
                ),
                details={
                    "fields_after": after_fields,
                    "fields_before": before_fields,
                    "fields_cleared": cleared_fields,
                    "helper_warnings_after": after_warnings,
                    "helper_warnings_before": before_warnings,
                },
            )
        )
        return

    collector.add_action(
        _base_action(
            finding,
            requested_action=plan.requested_action,
            applied_action="clear",
            status="applied",
            message="Cleared supported document-level metadata.",
            warnings=warnings,
            details={
                "fields_after": {},
                "fields_before": before_fields,
                "fields_cleared": cleared_fields,
                "helper_warnings_after": after_warnings,
                "helper_warnings_before": before_warnings,
            },
        )
    )



def _apply_image_finding(file_path: Path, extension: str, finding: dict, policy: dict, collector: _ResultCollector) -> None:
    category = str(finding.get("category") or "")
    plan = _image_plan_for_finding(policy)
    if action := _precomputed_action_if_needed(finding, category, plan, handled_categories=_IMAGE_CATEGORIES):
        collector.add_action(action)
        return

    if _has_excluded_image_scope(finding):
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Embedded-object, attachment, OLE, or OCR-dependent image handling is outside V1 scope.",
                manual_review_reason=(
                    "Embedded objects, attachments, OLE content, and OCR-dependent cleanup are excluded from image anonymization V1 scope."
                ),
            )
        )
        return

    image_index = _finding_image_index(finding)
    if image_index is None:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message="Image finding did not include a stable image_index locator.",
                manual_review_reason="Image finding lacks a deterministic image_index locator, so this runtime cannot safely target it.",
            )
        )
        return

    try:
        before_context = _image_before_context(file_path, image_index=image_index)
    except IndexError as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message=str(exc),
                manual_review_reason=str(exc),
                details={"image_index": image_index},
            )
        )
        return
    except NotImplementedError as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message=str(exc),
                manual_review_reason=str(exc),
                details={"image_index": image_index},
            )
        )
        return
    except Exception as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=plan.requested_action,
                applied_action="none",
                status="error",
                message=str(exc),
                details={"image_index": image_index},
            )
        )
        return

    requested_action = plan.requested_action
    effective_action = requested_action
    warnings: list[str] = []
    details = {
        "before": before_context,
        "container_extension": extension,
        "image_index": image_index,
    }

    try:
        if requested_action == "replace":
            replace_image(file_path, image_index, plan.replacement_path)
            details["replacement_path"] = str(plan.replacement_path)
        else:
            if requested_action == "remove":
                effective_action = "mask"
                warnings.append(
                    "Requested image removal was downgraded to masking because this V1 runtime does not claim a universally safe remove path."
                )
            with tempfile.TemporaryDirectory(prefix="office-automation-transform-mask-") as temp_dir:
                mask_path = _create_mask_image_asset(
                    Path(temp_dir),
                    image_index=image_index,
                    size=(before_context["width"], before_context["height"]),
                    color=plan.mask_color,
                )
                replace_image(file_path, image_index, mask_path)
                details["mask_color"] = plan.mask_color
    except NotImplementedError as exc:
        manual_review_reason = str(exc)
        if requested_action == "remove":
            manual_review_reason = (
                "Requested image removal could not be completed automatically. "
                f"{exc}"
            )
        collector.add_action(
            _base_action(
                finding,
                requested_action=requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message=str(exc),
                warnings=warnings,
                manual_review_reason=manual_review_reason,
                details=details,
            )
        )
        return
    except IndexError as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=requested_action,
                applied_action="manual_review",
                status="manual_review_required",
                message=str(exc),
                warnings=warnings,
                manual_review_reason=str(exc),
                details=details,
            )
        )
        return
    except Exception as exc:
        collector.add_action(
            _base_action(
                finding,
                requested_action=requested_action,
                applied_action="none",
                status="error",
                message=str(exc),
                warnings=warnings,
                details=details,
            )
        )
        return

    collector.add_action(
        _base_action(
            finding,
            requested_action=requested_action,
            applied_action=effective_action,
            status="applied",
            message=_image_success_message(requested_action=requested_action, applied_action=effective_action),
            warnings=warnings,
            details=details,
        )
    )



def _precomputed_action_if_needed(
    finding: dict,
    category: str,
    plan: _Plan | _ImagePlan,
    *,
    handled_categories: set[str],
) -> dict | None:
    if plan.policy_error:
        return _base_action(
            finding,
            requested_action=plan.requested_action,
            applied_action="none",
            status="error",
            message=plan.policy_error,
        )

    manual_review_reason = finding.get("manual_review_reason")
    if manual_review_reason:
        return _base_action(
            finding,
            requested_action=plan.requested_action,
            applied_action="manual_review",
            status="manual_review_required",
            message=str(manual_review_reason),
            manual_review_reason=str(manual_review_reason),
        )

    if category not in handled_categories:
        return _deferred_scope_action(finding)

    if plan.requested_action in {"skip", "report_only"}:
        mode_text = "report_only" if plan.requested_action == "report_only" else "skip"
        return _base_action(
            finding,
            requested_action=plan.requested_action,
            applied_action="skip",
            status="skipped",
            message=f"Policy requested {mode_text}; finding left in place.",
            warnings=[f"Finding was left in place because the confirmed policy requested {mode_text}."],
        )

    return None



def _deferred_scope_action(finding: dict) -> dict:
    category = str(finding.get("category") or "unknown")
    return _base_action(
        finding,
        requested_action="skip",
        applied_action="deferred",
        status="skipped",
        message=f"Category '{category}' is outside this transform runtime scope and was left unchanged.",
        warnings=[f"Category '{category}' was deferred because no supported V1 transform path is implemented for it."],
    )



def _plan_for_finding(finding: dict, policy: dict) -> _Plan:
    category = str(finding.get("category") or "")
    if category not in _STRUCTURAL_CATEGORIES:
        return _Plan(requested_action="skip")

    category_policy = policy.get(category)
    if not isinstance(category_policy, dict):
        return _Plan(
            requested_action="skip",
            policy_error=f"Policy entry for category '{category}' must be a dict.",
        )

    enabled = category_policy.get("enabled", True)
    requested_action = "skip" if enabled is False else str(category_policy.get("action", "skip")).strip().lower()
    if requested_action not in _ALLOWED_POLICY_ACTIONS:
        return _Plan(
            requested_action=requested_action or "skip",
            policy_error=f"Unsupported action '{requested_action}' for category '{category}'. Supported actions: clear, remove, replace, skip.",
        )

    replacement_text = _replacement_text_from_policy(category_policy)
    if requested_action == "replace" and replacement_text is None:
        return _Plan(
            requested_action=requested_action,
            policy_error=f"Policy action 'replace' for category '{category}' requires replacement text.",
        )

    return _Plan(requested_action=requested_action, replacement_text=replacement_text)



def _metadata_plan_for_finding(policy: dict) -> _Plan:
    category_policy = policy.get("metadata")
    if not isinstance(category_policy, dict):
        return _Plan(
            requested_action="skip",
            policy_error="Policy entry for category 'metadata' must be a dict.",
        )

    enabled = category_policy.get("enabled", True)
    requested_action = "skip" if enabled is False else str(category_policy.get("action", "skip")).strip().lower()
    if requested_action not in _ALLOWED_METADATA_ACTIONS:
        return _Plan(
            requested_action=requested_action or "skip",
            policy_error=(
                f"Unsupported action '{requested_action}' for category 'metadata'. "
                "Supported actions: clear, remove, report_only, skip."
            ),
        )
    return _Plan(requested_action=requested_action)



def _image_plan_for_finding(policy: dict) -> _ImagePlan:
    category_policy = policy.get("images")
    if not isinstance(category_policy, dict):
        return _ImagePlan(
            requested_action="skip",
            policy_error="Policy entry for category 'images' must be a dict.",
        )

    enabled = category_policy.get("enabled", True)
    requested_action = "report_only" if enabled is False else str(
        category_policy.get("mode", category_policy.get("action", "report_only"))
    ).strip().lower()
    if requested_action not in _ALLOWED_IMAGE_ACTIONS:
        return _ImagePlan(
            requested_action=requested_action or "report_only",
            policy_error=(
                f"Unsupported image mode '{requested_action}'. "
                "Supported modes: mask, remove, replace, report_only, skip."
            ),
        )

    replacement_path = _replacement_path_from_policy(category_policy)
    if requested_action == "replace" and replacement_path is None:
        return _ImagePlan(
            requested_action=requested_action,
            policy_error="Image mode 'replace' requires replacement_path.",
        )

    mask_color, mask_color_error = _mask_color_from_policy(category_policy)
    if mask_color_error:
        return _ImagePlan(
            requested_action=requested_action,
            replacement_path=replacement_path,
            policy_error=mask_color_error,
        )

    return _ImagePlan(
        requested_action=requested_action,
        replacement_path=replacement_path,
        mask_color=mask_color,
    )



def _replacement_text_from_policy(category_policy: dict) -> str | None:
    for key in ("replacement_text", "replacement", "text"):
        value = category_policy.get(key)
        if value is not None:
            return str(value)
    return None



def _replacement_path_from_policy(category_policy: dict) -> str | None:
    for key in ("replacement_path", "replacement_image", "replacement_asset", "path", "replacement"):
        value = category_policy.get(key)
        if value not in (None, ""):
            return str(value)
    return None



def _mask_color_from_policy(category_policy: dict) -> tuple[str, str | None]:
    color_value = category_policy.get("mask_color", category_policy.get("color", _DEFAULT_IMAGE_MASK_COLOR))
    try:
        return _normalize_mask_color(color_value), None
    except ValueError as exc:
        return _DEFAULT_IMAGE_MASK_COLOR, str(exc)



def _group_findings(detected: list[dict]) -> list[tuple[dict, list[dict]]]:
    grouped: OrderedDict[tuple[str, str], dict] = OrderedDict()
    for raw_finding in detected:
        finding = _sorted_copy(raw_finding if isinstance(raw_finding, dict) else {"payload": {"raw_finding": raw_finding}})
        file_path_value = finding.get("file_path")
        file_path = str(file_path_value) if file_path_value not in (None, "") else "<unknown>"
        extension = str(finding.get("extension") or _path_extension(Path(file_path)) if file_path != "<unknown>" else "").lower()
        relative_path = finding.get("relative_path")
        key = (file_path, extension)
        if key not in grouped:
            grouped[key] = {
                "result": {
                    "file_path": file_path,
                    "relative_path": relative_path,
                    "extension": extension,
                    "output_path": file_path if file_path != "<unknown>" else None,
                    "status": "success",
                    "actions": [],
                    "warnings": [],
                    "manual_review_items": [],
                },
                "findings": [],
            }
        grouped[key]["findings"].append(finding)
    return [(value["result"], value["findings"]) for value in grouped.values()]



def _base_action(
    finding: dict,
    *,
    requested_action: str,
    applied_action: str,
    status: str,
    message: str,
    warnings: list[str] | None = None,
    manual_review_reason: str | None = None,
    details: dict | None = None,
) -> dict:
    return {
        "finding_id": finding.get("finding_id"),
        "category": finding.get("category"),
        "location": _sorted_copy(finding.get("location") or {}),
        "payload": _sorted_copy(finding.get("payload") or {}),
        "confidence": finding.get("confidence"),
        "requested_action": requested_action,
        "applied_action": applied_action,
        "status": status,
        "message": message,
        "warnings": list(warnings or []),
        "manual_review_reason": manual_review_reason,
        "details": _sorted_copy(details or {}),
    }



def _finalize_file_status(file_result: dict) -> str:
    actions = file_result["actions"]
    if any(action.get("status") == "error" for action in actions):
        return "error"
    if file_result["manual_review_items"]:
        if any(action.get("status") == "applied" for action in actions):
            return "partial_success"
        return "manual_review_required"
    if file_result["warnings"] or any(action.get("status") == "skipped" for action in actions):
        return "partial_success"
    return "success"



def _docx_container_by_part(document: Document, *, category: str, part_name: str):
    variants = _DOCX_HEADER_VARIANTS if category == "headers" else _DOCX_FOOTER_VARIANTS
    seen_parts: set[str] = set()
    for section in document.sections:
        for _variant_name, attr_name in variants:
            container = getattr(section, attr_name)
            current_part = str(container.part.partname)
            if current_part in seen_parts:
                continue
            seen_parts.add(current_part)
            if current_part == part_name:
                return container
    raise ValueError(f"Could not locate DOCX {category[:-1]} part '{part_name}'.")



def _rewrite_docx_story_container(container, replacement_text: str) -> None:
    for child in list(container._element):
        container._element.remove(child)
    container.add_paragraph(replacement_text)



def _non_empty_mapping(values: dict) -> dict:
    return {
        str(key): value
        for key, value in values.items()
        if value not in (None, "")
    }



def _finding_image_index(finding: dict) -> int | None:
    location = finding.get("location") or {}
    payload = finding.get("payload") or {}
    candidate = location.get("image_index", payload.get("image_index"))
    if isinstance(candidate, bool) or not isinstance(candidate, int) or candidate < 0:
        return None
    return candidate



def _image_before_context(file_path: Path, *, image_index: int) -> dict:
    with tempfile.TemporaryDirectory(prefix="office-automation-transform-images-") as temp_dir:
        extracted_paths = extract_images(file_path, Path(temp_dir))
        if image_index >= len(extracted_paths):
            raise IndexError(
                f"Image index {image_index} is out of range for '{file_path}'. Available image count: {len(extracted_paths)}."
            )
        image_path = extracted_paths[image_index]
        width, height = _image_dimensions_for_path(image_path)
        return {
            "extracted_filename": image_path.name,
            "height": height,
            "width": width,
        }



def _create_mask_image_asset(directory: Path, *, image_index: int, size: tuple[int, int], color: str) -> Path:
    width = int(size[0]) if size[0] else 1
    height = int(size[1]) if size[1] else 1
    width = max(width, 1)
    height = max(height, 1)
    rgb = ImageColor.getrgb(color)
    path = directory / f"image-mask-{image_index:04d}.png"
    Image.new("RGB", (width, height), rgb).save(path, format="PNG")
    return path



def _image_success_message(*, requested_action: str, applied_action: str) -> str:
    if requested_action == "replace":
        return "Replaced image content with the policy-provided replacement asset."
    if requested_action == "mask":
        return "Masked image content by replacing the target slot with a solid-color image."
    if requested_action == "remove" and applied_action == "mask":
        return "Requested image removal was fulfilled by a conservative mask replacement because a safe remove path is not claimed in V1."
    return "Updated image content."



def _has_excluded_image_scope(finding: dict) -> bool:
    combined_values = list(_iter_string_values(finding.get("location") or {})) + list(_iter_string_values(finding.get("payload") or {}))
    return any(marker in value for value in combined_values for marker in _EXCLUDED_IMAGE_CONTAINER_MARKERS)



def _iter_string_values(value):
    if isinstance(value, dict):
        for key, nested in value.items():
            yield from _iter_string_values(key)
            yield from _iter_string_values(nested)
        return
    if isinstance(value, list):
        for item in value:
            yield from _iter_string_values(item)
        return
    if value is None:
        return
    yield str(value).strip().lower()



def _image_dimensions_for_path(path: Path) -> tuple[int, int]:
    with Image.open(path) as image:
        width, height = image.size
    return width, height



def _normalize_mask_color(value) -> str:
    if isinstance(value, str):
        color = value.strip()
        if not color:
            raise ValueError("Image mask color must not be empty.")
        rgb = ImageColor.getrgb(color)
        return "#" + "".join(f"{channel:02x}" for channel in rgb)
    if isinstance(value, (list, tuple)) and len(value) in {3, 4}:
        channels = value[:3]
        if any(not isinstance(channel, int) or channel < 0 or channel > 255 for channel in channels):
            raise ValueError("Image mask color tuples must contain integer RGB values from 0 to 255.")
        return "#" + "".join(f"{channel:02x}" for channel in channels)
    raise ValueError("Image mask color must be a Pillow-compatible color string or RGB tuple.")



def _save_openpyxl_workbook(workbook, path: Path) -> None:
    temp_path = _temporary_output_path(path)
    try:
        workbook.save(temp_path)
        temp_path.replace(path)
    finally:
        if temp_path.exists():
            temp_path.unlink()



def _save_docx_document(document: Document, path: Path) -> None:
    temp_path = _temporary_output_path(path)
    try:
        document.save(temp_path)
        temp_path.replace(path)
    finally:
        if temp_path.exists():
            temp_path.unlink()



def _save_presentation(presentation: Presentation, path: Path) -> None:
    temp_path = _temporary_output_path(path)
    try:
        presentation.save(temp_path)
        temp_path.replace(path)
    finally:
        if temp_path.exists():
            temp_path.unlink()



def _save_pdf_document(document: fitz.Document, path: Path) -> None:
    temp_path = _temporary_output_path(path)
    try:
        document.save(temp_path, garbage=4)
        temp_path.replace(path)
    finally:
        if temp_path.exists():
            temp_path.unlink()



def _temporary_output_path(path: Path) -> Path:
    handle = tempfile.NamedTemporaryFile(prefix=f".{path.stem}-", suffix=path.suffix, dir=path.parent, delete=False)
    handle.close()
    temp_path = Path(handle.name)
    temp_path.unlink(missing_ok=True)
    return temp_path



def _required_location_value(finding: dict, key: str):
    location = finding.get("location")
    if not isinstance(location, dict) or key not in location:
        raise ValueError(f"Finding '{finding.get('finding_id')}' is missing required location key '{key}'.")
    return location[key]



def _close_excel_workbook(workbook) -> None:
    vba_archive = getattr(workbook, "vba_archive", None)
    if vba_archive is not None:
        vba_archive.close()
    close = getattr(workbook, "close", None)
    if callable(close):
        close()



def _path_extension(path: Path) -> str:
    return path.suffix.lower().lstrip(".")



def _sorted_copy(value):
    if isinstance(value, dict):
        return {key: _sorted_copy(value[key]) for key in sorted(value)}
    if isinstance(value, list):
        return [_sorted_copy(item) for item in value]
    return value



def _location_key(location: dict) -> str:
    return json.dumps(location, ensure_ascii=False, separators=(",", ":"), sort_keys=True)
