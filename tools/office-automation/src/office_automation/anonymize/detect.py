"""Folder detection helpers for V1 anonymization targets."""

from __future__ import annotations

from collections.abc import Mapping
import json
from pathlib import Path
import re
from tempfile import TemporaryDirectory

import fitz
from docx import Document
from openpyxl import load_workbook
from PIL import Image
from pptx import Presentation

from office_automation.common.files import list_office_files
from office_automation.common.images import extract_images
from office_automation.common.metadata import read_metadata

__all__ = ["detect"]

_HEADER_FOOTER_VARIANTS = (
    ("odd", "oddHeader", "oddFooter"),
    ("even", "evenHeader", "evenFooter"),
    ("first", "firstHeader", "firstFooter"),
)
_DOCX_HEADER_VARIANTS = (
    ("default", "header"),
    ("first", "first_page_header"),
    ("even", "even_page_header"),
)
_DOCX_FOOTER_VARIANTS = (
    ("default", "footer"),
    ("first", "first_page_footer"),
    ("even", "even_page_footer"),
)
_PPT_MANUAL_REVIEW_REASONS = {
    "comments": "python-pptx does not expose PowerPoint comment threads reliably in this V1 runtime.",
    "headers": "PowerPoint header surfaces are not exposed reliably enough to prove a presentation is clean in this V1 runtime.",
    "footers": "PowerPoint footer surfaces are not exposed reliably enough to prove a presentation is clean in this V1 runtime.",
}
_PDF_MANUAL_REVIEW_REASONS = {
    "headers": "PDF header detection is heuristic-sensitive and this V1 runtime does not claim reliable clean-state proof for headers.",
    "footers": "PDF footer detection is heuristic-sensitive and this V1 runtime does not claim reliable clean-state proof for footers.",
}
_CATEGORY_ORDER = {
    "comments": 0,
    "notes": 1,
    "headers": 2,
    "footers": 3,
    "metadata": 4,
    "body_text": 5,
    "images": 6,
}
_BODY_TEXT_LIST_FIELDS = (
    "person_names",
    "company_names",
    "emails",
    "phones",
    "addresses",
    "domains",
    "exact_phrases",
    "context_terms",
)
_BODY_TEXT_ALLOWED_KEYS = {*_BODY_TEXT_LIST_FIELDS, "replacement_text", "replacement_map"}
_EMAIL_PATTERN = re.compile(r"\b[A-Z0-9._%+-]+@[A-Z0-9.-]+\.[A-Z]{2,}\b", re.IGNORECASE)
_DOMAIN_PATTERN = re.compile(r"(?<!@)\b(?:[A-Z0-9](?:[A-Z0-9-]{0,61}[A-Z0-9])?\.)+[A-Z]{2,}\b", re.IGNORECASE)
_PHONE_PATTERN = re.compile(r"(?<!\w)(?:\+?\d[\d()\- ]{8,}\d)(?!\w)")
_CONTEXT_ASSISTED_PHRASE_TEMPLATE = r"\b{context_term}\b\s*[:\-]\s*([A-Z][\w.&'/-]*(?:\s+[A-Z][\w.&'/-]*){{1,3}})"


def detect(
    target_folder: str,
    extensions: list[str] | None = None,
    body_text_candidate_inputs: dict | None = None,
) -> list[dict]:
    """Scan target_folder for comments, notes, headers, footers, metadata, images, and SG5 body text."""
    body_text_detection_enabled = body_text_candidate_inputs is not None
    normalized_body_text_inputs = _validate_body_text_candidate_inputs(body_text_candidate_inputs)
    folder = Path(target_folder)
    _validate_target_folder(folder)

    files = list_office_files(folder, extensions=extensions)
    findings: list[dict] = []
    for file_path in files:
        relative_path = file_path.relative_to(folder).as_posix()
        findings.extend(
            _scan_file(
                file_path,
                relative_path=relative_path,
                body_text_detection_enabled=body_text_detection_enabled,
                body_text_candidate_inputs=normalized_body_text_inputs,
            )
        )

    return sorted(findings, key=_finding_sort_key)



def _scan_file(
    file_path: Path,
    *,
    relative_path: str,
    body_text_detection_enabled: bool,
    body_text_candidate_inputs: dict,
) -> list[dict]:
    extension = _path_extension(file_path)
    if extension in {"xlsx", "xlsm"}:
        return _scan_excel_file(
            file_path,
            relative_path=relative_path,
            body_text_detection_enabled=body_text_detection_enabled,
            body_text_candidate_inputs=body_text_candidate_inputs,
        )
    if extension == "docx":
        return _scan_docx_file(
            file_path,
            relative_path=relative_path,
            body_text_detection_enabled=body_text_detection_enabled,
            body_text_candidate_inputs=body_text_candidate_inputs,
        )
    if extension == "pptx":
        return _scan_pptx_file(
            file_path,
            relative_path=relative_path,
            body_text_detection_enabled=body_text_detection_enabled,
            body_text_candidate_inputs=body_text_candidate_inputs,
        )
    if extension == "pdf":
        return _scan_pdf_file(
            file_path,
            relative_path=relative_path,
            body_text_detection_enabled=body_text_detection_enabled,
            body_text_candidate_inputs=body_text_candidate_inputs,
        )
    raise ValueError(f"Unsupported extension '{extension}' for '{file_path}'.")



def _scan_excel_file(
    file_path: Path,
    *,
    relative_path: str,
    body_text_detection_enabled: bool,
    body_text_candidate_inputs: dict,
) -> list[dict]:
    workbook = load_workbook(file_path, keep_vba=_path_extension(file_path) == "xlsm")
    findings: list[dict] = []
    try:
        findings.extend(_scan_excel_notes(workbook, file_path=file_path, relative_path=relative_path))
        findings.extend(_scan_excel_headers_and_footers(workbook, file_path=file_path, relative_path=relative_path))
        if body_text_detection_enabled:
            findings.extend(
                _scan_excel_body_text(
                    workbook,
                    file_path=file_path,
                    relative_path=relative_path,
                    body_text_candidate_inputs=body_text_candidate_inputs,
                )
            )
    finally:
        _close_excel_workbook(workbook)

    findings.extend(_metadata_findings(file_path, relative_path=relative_path))
    findings.extend(_image_findings(file_path, relative_path=relative_path))
    return findings



def _scan_docx_file(
    file_path: Path,
    *,
    relative_path: str,
    body_text_detection_enabled: bool,
    body_text_candidate_inputs: dict,
) -> list[dict]:
    document = Document(file_path)
    findings: list[dict] = []

    try:
        findings.extend(_scan_docx_comments(document, file_path=file_path, relative_path=relative_path))
    except Exception as exc:
        findings.append(
            _manual_review_finding(
                file_path=file_path,
                relative_path=relative_path,
                category="comments",
                reason=f"Word comments could not be inspected safely: {exc}",
                location={"scope": "document"},
            )
        )

    try:
        findings.extend(_scan_docx_headers_and_footers(document, file_path=file_path, relative_path=relative_path))
    except Exception as exc:
        findings.extend(
            [
                _manual_review_finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    category="headers",
                    reason=f"Word headers could not be inspected safely: {exc}",
                    location={"scope": "document"},
                ),
                _manual_review_finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    category="footers",
                    reason=f"Word footers could not be inspected safely: {exc}",
                    location={"scope": "document"},
                ),
            ]
        )

    findings.extend(_metadata_findings(file_path, relative_path=relative_path))
    if body_text_detection_enabled:
        findings.extend(
            _scan_docx_body_text(
                document,
                file_path=file_path,
                relative_path=relative_path,
                body_text_candidate_inputs=body_text_candidate_inputs,
            )
        )
    findings.extend(_image_findings(file_path, relative_path=relative_path))
    return findings



def _scan_pptx_file(
    file_path: Path,
    *,
    relative_path: str,
    body_text_detection_enabled: bool,
    body_text_candidate_inputs: dict,
) -> list[dict]:
    presentation = Presentation(file_path)
    findings: list[dict] = []

    findings.append(
        _manual_review_finding(
            file_path=file_path,
            relative_path=relative_path,
            category="comments",
            reason=_PPT_MANUAL_REVIEW_REASONS["comments"],
            location={"scope": "presentation"},
        )
    )

    try:
        findings.extend(_scan_pptx_notes(presentation, file_path=file_path, relative_path=relative_path))
    except Exception as exc:
        findings.append(
            _manual_review_finding(
                file_path=file_path,
                relative_path=relative_path,
                category="notes",
                reason=f"PowerPoint notes could not be inspected safely: {exc}",
                location={"scope": "presentation"},
            )
        )

    findings.append(
        _manual_review_finding(
            file_path=file_path,
            relative_path=relative_path,
            category="headers",
            reason=_PPT_MANUAL_REVIEW_REASONS["headers"],
            location={"scope": "presentation"},
        )
    )
    findings.append(
        _manual_review_finding(
            file_path=file_path,
            relative_path=relative_path,
            category="footers",
            reason=_PPT_MANUAL_REVIEW_REASONS["footers"],
            location={"scope": "presentation"},
        )
    )

    findings.extend(_metadata_findings(file_path, relative_path=relative_path))
    if body_text_detection_enabled:
        findings.extend(
            _scan_pptx_body_text(
                presentation,
                file_path=file_path,
                relative_path=relative_path,
                body_text_candidate_inputs=body_text_candidate_inputs,
            )
        )
    findings.extend(_image_findings(file_path, relative_path=relative_path))
    return findings



def _scan_pdf_file(
    file_path: Path,
    *,
    relative_path: str,
    body_text_detection_enabled: bool,
    body_text_candidate_inputs: dict,
) -> list[dict]:
    findings: list[dict] = []
    findings.extend(_scan_pdf_comments(file_path, relative_path=relative_path))
    findings.append(
        _manual_review_finding(
            file_path=file_path,
            relative_path=relative_path,
            category="headers",
            reason=_PDF_MANUAL_REVIEW_REASONS["headers"],
            location={"scope": "document"},
        )
    )
    findings.append(
        _manual_review_finding(
            file_path=file_path,
            relative_path=relative_path,
            category="footers",
            reason=_PDF_MANUAL_REVIEW_REASONS["footers"],
            location={"scope": "document"},
        )
    )
    findings.extend(_metadata_findings(file_path, relative_path=relative_path))
    if body_text_detection_enabled:
        findings.extend(
            _scan_pdf_body_text(
                file_path,
                relative_path=relative_path,
                body_text_candidate_inputs=body_text_candidate_inputs,
            )
        )
    findings.extend(_image_findings(file_path, relative_path=relative_path, image_locations=_pdf_image_locations(file_path)))
    return findings



def _scan_excel_notes(workbook, *, file_path: Path, relative_path: str) -> list[dict]:
    findings: list[dict] = []
    for worksheet in workbook.worksheets:
        comment_cells = []
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.comment is None:
                    continue
                comment_cells.append(cell)
        for cell in sorted(comment_cells, key=lambda item: (item.coordinate, item.row, item.column)):
            text = _normalized_text(cell.comment.text)
            author = _normalized_text(cell.comment.author)
            findings.append(
                _finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    category="notes",
                    location={"sheet": worksheet.title, "cell": cell.coordinate},
                    payload={
                        "text": text,
                        "author": author,
                    },
                    action_hint="remove",
                    confidence="high",
                    manual_review_reason=None,
                )
            )
    return findings



def _scan_excel_headers_and_footers(workbook, *, file_path: Path, relative_path: str) -> list[dict]:
    findings: list[dict] = []
    for worksheet in workbook.worksheets:
        for variant_name, header_attr, footer_attr in _HEADER_FOOTER_VARIANTS:
            findings.extend(
                _excel_header_footer_findings(
                    worksheet,
                    file_path=file_path,
                    relative_path=relative_path,
                    category="headers",
                    section_attr=header_attr,
                    variant_name=variant_name,
                )
            )
            findings.extend(
                _excel_header_footer_findings(
                    worksheet,
                    file_path=file_path,
                    relative_path=relative_path,
                    category="footers",
                    section_attr=footer_attr,
                    variant_name=variant_name,
                )
            )
    return findings



def _excel_header_footer_findings(
    worksheet,
    *,
    file_path: Path,
    relative_path: str,
    category: str,
    section_attr: str,
    variant_name: str,
) -> list[dict]:
    findings: list[dict] = []
    section = getattr(worksheet, section_attr)
    for part_name in ("left", "center", "right"):
        part = getattr(section, part_name)
        text = _normalized_text(part.text)
        if text is None:
            continue
        findings.append(
            _finding(
                file_path=file_path,
                relative_path=relative_path,
                category=category,
                location={
                    "sheet": worksheet.title,
                    "variant": variant_name,
                    "part": part_name,
                },
                payload={"text": text},
                action_hint="remove",
                confidence="high",
                manual_review_reason=None,
            )
        )
    return findings



def _scan_docx_comments(document: Document, *, file_path: Path, relative_path: str) -> list[dict]:
    findings: list[dict] = []
    comments = sorted(document.comments, key=lambda comment: int(comment.comment_id))
    for comment in comments:
        findings.append(
            _finding(
                file_path=file_path,
                relative_path=relative_path,
                category="comments",
                location={"comment_id": int(comment.comment_id)},
                payload={
                    "text": _normalized_text(comment.text),
                    "author": _normalized_text(comment.author),
                    "initials": _normalized_text(comment.initials),
                    "timestamp": _normalized_datetime(getattr(comment, "timestamp", None)),
                },
                action_hint="remove",
                confidence="high",
                manual_review_reason=None,
            )
        )
    return findings



def _scan_docx_headers_and_footers(document: Document, *, file_path: Path, relative_path: str) -> list[dict]:
    findings: list[dict] = []
    seen_header_parts: set[str] = set()
    seen_footer_parts: set[str] = set()

    for section_index, section in enumerate(document.sections, start=1):
        findings.extend(
            _docx_header_footer_findings(
                section,
                section_index=section_index,
                file_path=file_path,
                relative_path=relative_path,
                category="headers",
                variants=_DOCX_HEADER_VARIANTS,
                seen_parts=seen_header_parts,
            )
        )
        findings.extend(
            _docx_header_footer_findings(
                section,
                section_index=section_index,
                file_path=file_path,
                relative_path=relative_path,
                category="footers",
                variants=_DOCX_FOOTER_VARIANTS,
                seen_parts=seen_footer_parts,
            )
        )
    return findings



def _docx_header_footer_findings(
    section,
    *,
    section_index: int,
    file_path: Path,
    relative_path: str,
    category: str,
    variants: tuple[tuple[str, str], ...],
    seen_parts: set[str],
) -> list[dict]:
    findings: list[dict] = []
    for variant_name, attr_name in variants:
        container = getattr(section, attr_name)
        part_name = str(container.part.partname)
        if part_name in seen_parts:
            continue
        seen_parts.add(part_name)

        text = _docx_container_text(container)
        if text is None:
            continue
        findings.append(
            _finding(
                file_path=file_path,
                relative_path=relative_path,
                category=category,
                location={
                    "section_index": section_index,
                    "variant": variant_name,
                    "part": part_name,
                },
                payload={"text": text},
                action_hint="remove",
                confidence="high",
                manual_review_reason=None,
            )
        )
    return findings



def _scan_pptx_notes(presentation: Presentation, *, file_path: Path, relative_path: str) -> list[dict]:
    findings: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        if not slide.has_notes_slide:
            continue
        note_shapes = []
        for shape in slide.notes_slide.shapes:
            text_frame = getattr(shape, "text_frame", None)
            if text_frame is None:
                continue
            shape_name = _normalized_text(getattr(shape, "name", None)) or ""
            if "slide image" in shape_name.casefold() or "slide number" in shape_name.casefold():
                continue
            text = _normalized_text(text_frame.text)
            if text is None:
                continue
            note_shapes.append((shape.shape_id, shape_name, text))
        for shape_id, shape_name, text in sorted(note_shapes, key=lambda item: (item[0], item[1])):
            findings.append(
                _finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    category="notes",
                    location={"slide_number": slide_number, "shape_id": shape_id},
                    payload={
                        "shape_name": shape_name or None,
                        "text": text,
                    },
                    action_hint="remove",
                    confidence="high",
                    manual_review_reason=None,
                )
            )
    return findings



def _scan_pdf_comments(file_path: Path, *, relative_path: str) -> list[dict]:
    document = fitz.open(file_path)
    try:
        if getattr(document, "needs_pass", False):
            raise NotImplementedError(f"Encrypted PDF files are outside detection V1 scope: '{file_path}'.")

        findings: list[dict] = []
        for page_number in range(1, document.page_count + 1):
            page = document.load_page(page_number - 1)
            annotations = list(page.annots() or [])
            for annotation_index, annotation in enumerate(annotations):
                annotation_type = annotation.type[1] if isinstance(annotation.type, tuple) else str(annotation.type)
                info = annotation.info or {}
                findings.append(
                    _finding(
                        file_path=file_path,
                        relative_path=relative_path,
                        category="comments",
                        location={
                            "page_number": page_number,
                            "annotation_index": annotation_index,
                            "xref": annotation.xref,
                        },
                        payload={
                            "annotation_type": annotation_type,
                            "content": _normalized_text(info.get("content")),
                            "author": _normalized_text(info.get("title")),
                            "subject": _normalized_text(info.get("subject")),
                            "annotation_id": _normalized_text(info.get("id")),
                            "rect": [
                                float(annotation.rect.x0),
                                float(annotation.rect.y0),
                                float(annotation.rect.x1),
                                float(annotation.rect.y1),
                            ],
                        },
                        action_hint="remove",
                        confidence="high",
                        manual_review_reason=None,
                    )
                )
        return findings
    finally:
        document.close()



def _validate_body_text_candidate_inputs(value: Mapping | None) -> dict:
    if value is None:
        return _empty_body_text_candidate_inputs()
    if not isinstance(value, Mapping):
        raise TypeError("body_text_candidate_inputs must be a mapping or None.")

    unknown_keys = sorted(str(key) for key in value if str(key) not in _BODY_TEXT_ALLOWED_KEYS)
    if unknown_keys:
        raise ValueError(f"Unsupported body_text_candidate_inputs key '{unknown_keys[0]}'.")

    normalized = _empty_body_text_candidate_inputs()
    for field_name in _BODY_TEXT_LIST_FIELDS:
        field_value = value.get(field_name, [])
        if not isinstance(field_value, list) or not all(isinstance(item, str) for item in field_value):
            raise TypeError(f"body_text_candidate_inputs field '{field_name}' must be a list of strings.")
        normalized[field_name] = [item for item in field_value if item]

    replacement_text = value.get("replacement_text")
    if replacement_text is not None and not isinstance(replacement_text, str):
        raise TypeError("body_text_candidate_inputs field 'replacement_text' must be a string or None.")
    normalized["replacement_text"] = replacement_text

    replacement_map = value.get("replacement_map", {})
    if not isinstance(replacement_map, dict) or not all(
        isinstance(key, str) and isinstance(item, str) for key, item in replacement_map.items()
    ):
        raise TypeError("body_text_candidate_inputs field 'replacement_map' must be a mapping of strings to strings.")
    normalized["replacement_map"] = dict(replacement_map)
    return normalized



def _empty_body_text_candidate_inputs() -> dict:
    return {
        "person_names": [],
        "company_names": [],
        "emails": [],
        "phones": [],
        "addresses": [],
        "domains": [],
        "exact_phrases": [],
        "context_terms": [],
        "replacement_text": None,
        "replacement_map": {},
    }



def _scan_excel_body_text(workbook, *, file_path: Path, relative_path: str, body_text_candidate_inputs: dict) -> list[dict]:
    findings: list[dict] = []
    for worksheet in workbook.worksheets:
        for row in worksheet.iter_rows():
            for cell in row:
                if cell.data_type == "f":
                    continue
                text = _normalized_text(cell.value)
                if text is None:
                    continue
                for match in _collect_body_text_matches(text, body_text_candidate_inputs):
                    findings.append(
                        _body_text_finding(
                            file_path=file_path,
                            relative_path=relative_path,
                            location={
                                "sheet": worksheet.title,
                                "cell": cell.coordinate,
                                "match_start": match["match_start"],
                                "match_end": match["match_end"],
                            },
                            payload={
                                "matched_text": match["matched_text"],
                                "normalized_text": match["normalized_text"],
                                "excerpt": _excerpt(text, match["match_start"], match["match_end"]),
                                "surface_type": "excel_cell",
                            },
                            action_hint="candidate_confirmation_required",
                            confidence=match["confidence"],
                            manual_review_reason=None,
                            source=match["source"],
                            reason_tags=match["reason_tags"],
                        )
                    )
    return findings



def _scan_docx_body_text(document: Document, *, file_path: Path, relative_path: str, body_text_candidate_inputs: dict) -> list[dict]:
    findings: list[dict] = []
    for paragraph_index, paragraph in enumerate(document.paragraphs):
        text = _normalized_text(paragraph.text)
        if text is None:
            continue
        for match in _collect_body_text_matches(text, body_text_candidate_inputs):
            findings.append(
                _body_text_finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    location={
                        "surface": "paragraph",
                        "paragraph_index": paragraph_index,
                        "match_start": match["match_start"],
                        "match_end": match["match_end"],
                    },
                    payload={
                        "matched_text": match["matched_text"],
                        "normalized_text": match["normalized_text"],
                        "excerpt": _excerpt(text, match["match_start"], match["match_end"]),
                        "surface_type": "docx_paragraph",
                    },
                    action_hint="candidate_confirmation_required",
                    confidence=match["confidence"],
                    manual_review_reason=None,
                    source=match["source"],
                    reason_tags=match["reason_tags"],
                )
            )

    for table_index, table in enumerate(document.tables):
        for row_index, row in enumerate(table.rows):
            for column_index, cell in enumerate(row.cells):
                for paragraph_index, paragraph in enumerate(cell.paragraphs):
                    text = _normalized_text(paragraph.text)
                    if text is None:
                        continue
                    for match in _collect_body_text_matches(text, body_text_candidate_inputs):
                        findings.append(
                            _body_text_finding(
                                file_path=file_path,
                                relative_path=relative_path,
                                location={
                                    "surface": "table_cell",
                                    "table_index": table_index,
                                    "row_index": row_index,
                                    "column_index": column_index,
                                    "paragraph_index": paragraph_index,
                                    "match_start": match["match_start"],
                                    "match_end": match["match_end"],
                                },
                                payload={
                                    "matched_text": match["matched_text"],
                                    "normalized_text": match["normalized_text"],
                                    "excerpt": _excerpt(text, match["match_start"], match["match_end"]),
                                    "surface_type": "docx_table_cell",
                                },
                                action_hint="candidate_confirmation_required",
                                confidence=match["confidence"],
                                manual_review_reason=None,
                                source=match["source"],
                                reason_tags=match["reason_tags"],
                            )
                        )
    return findings



def _scan_pptx_body_text(presentation: Presentation, *, file_path: Path, relative_path: str, body_text_candidate_inputs: dict) -> list[dict]:
    findings: list[dict] = []
    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs):
                    text = _normalized_text(paragraph.text)
                    if text is None:
                        continue
                    for match in _collect_body_text_matches(text, body_text_candidate_inputs):
                        findings.append(
                            _body_text_finding(
                                file_path=file_path,
                                relative_path=relative_path,
                                location={
                                    "slide_number": slide_number,
                                    "shape_id": int(shape.shape_id),
                                    "paragraph_index": paragraph_index,
                                    "match_start": match["match_start"],
                                    "match_end": match["match_end"],
                                },
                                payload={
                                    "matched_text": match["matched_text"],
                                    "normalized_text": match["normalized_text"],
                                    "excerpt": _excerpt(text, match["match_start"], match["match_end"]),
                                    "surface_type": "pptx_text_frame",
                                },
                                action_hint="candidate_confirmation_required",
                                confidence=match["confidence"],
                                manual_review_reason=None,
                                source=match["source"],
                                reason_tags=match["reason_tags"],
                            )
                        )
            if getattr(shape, "has_table", False):
                for row_index, row in enumerate(shape.table.rows):
                    for column_index, cell in enumerate(row.cells):
                        for paragraph_index, paragraph in enumerate(cell.text_frame.paragraphs):
                            text = _normalized_text(paragraph.text)
                            if text is None:
                                continue
                            for match in _collect_body_text_matches(text, body_text_candidate_inputs):
                                findings.append(
                                    _body_text_finding(
                                        file_path=file_path,
                                        relative_path=relative_path,
                                        location={
                                            "slide_number": slide_number,
                                            "shape_id": int(shape.shape_id),
                                            "row_index": row_index,
                                            "column_index": column_index,
                                            "paragraph_index": paragraph_index,
                                            "match_start": match["match_start"],
                                            "match_end": match["match_end"],
                                        },
                                        payload={
                                            "matched_text": match["matched_text"],
                                            "normalized_text": match["normalized_text"],
                                            "excerpt": _excerpt(text, match["match_start"], match["match_end"]),
                                            "surface_type": "pptx_table_cell",
                                        },
                                        action_hint="candidate_confirmation_required",
                                        confidence=match["confidence"],
                                        manual_review_reason=None,
                                        source=match["source"],
                                        reason_tags=match["reason_tags"],
                                    )
                                )
    return findings



def _scan_pdf_body_text(file_path: Path, *, relative_path: str, body_text_candidate_inputs: dict) -> list[dict]:
    document = fitz.open(file_path)
    try:
        if getattr(document, "needs_pass", False):
            raise NotImplementedError(f"Encrypted PDF files are outside detection V1 scope: '{file_path}'.")

        findings: list[dict] = []
        for page_number in range(1, document.page_count + 1):
            page = document.load_page(page_number - 1)
            for span_index, span in enumerate(_pdf_text_spans(page)):
                for match in _collect_body_text_matches(span["text"], body_text_candidate_inputs):
                    findings.append(
                        _body_text_finding(
                            file_path=file_path,
                            relative_path=relative_path,
                            location={
                                "page_number": page_number,
                                "span_index": span_index,
                                "match_start": match["match_start"],
                                "match_end": match["match_end"],
                                "bbox": _pdf_match_bbox(
                                    page,
                                    matched_text=match["matched_text"],
                                    fallback_bbox=span["bbox"],
                                    clip_rect=fitz.Rect(span["bbox"]),
                                ),
                            },
                            payload={
                                "matched_text": match["matched_text"],
                                "normalized_text": match["normalized_text"],
                                "excerpt": _excerpt(span["text"], match["match_start"], match["match_end"]),
                                "surface_type": "pdf_text_span",
                            },
                            action_hint="review",
                            confidence="low",
                            manual_review_reason="PDF text-layer matches remain review-first because layout-safe rewrite is not proven in this runtime.",
                            source=match["source"],
                            reason_tags=sorted({*match["reason_tags"], "pdf_text_layer"}),
                        )
                    )
        return findings
    finally:
        document.close()



def _collect_body_text_matches(text: str, body_text_candidate_inputs: dict) -> list[dict]:
    matches_by_key: dict[tuple[int, int, str], dict] = {}

    for field_name, reason_tag in (
        ("exact_phrases", "exact_phrase"),
        ("person_names", "person_hint"),
        ("company_names", "company_hint"),
        ("emails", "email_hint"),
        ("phones", "phone_hint"),
        ("addresses", "address_hint"),
        ("domains", "domain_hint"),
    ):
        confidence = "high" if field_name != "addresses" else "medium"
        for hint in body_text_candidate_inputs[field_name]:
            for start, end, matched_text in _find_case_insensitive_occurrences(text, hint):
                _merge_body_text_match(
                    matches_by_key,
                    start=start,
                    end=end,
                    matched_text=matched_text,
                    source="user_hint",
                    reason_tag=reason_tag,
                    confidence=confidence,
                )

    for regex, reason_tag, confidence in (
        (_EMAIL_PATTERN, "email_pattern", "high"),
        (_DOMAIN_PATTERN, "domain_pattern", "high"),
        (_PHONE_PATTERN, "phone_pattern", "medium"),
    ):
        for match in regex.finditer(text):
            _merge_body_text_match(
                matches_by_key,
                start=match.start(),
                end=match.end(),
                matched_text=match.group(0),
                source="pattern",
                reason_tag=reason_tag,
                confidence=confidence,
            )

    for context_term in body_text_candidate_inputs["context_terms"]:
        for start, end, matched_text in _find_context_assisted_phrase_occurrences(text, context_term):
            _merge_body_text_match(
                matches_by_key,
                start=start,
                end=end,
                matched_text=matched_text,
                source="heuristic",
                reason_tag="context_assisted_phrase",
                confidence="medium",
            )

    if any(_contains_case_insensitive(text, context_term) for context_term in body_text_candidate_inputs["context_terms"]):
        for match in matches_by_key.values():
            match["reason_tags"].add("context_term")

    return [
        {
            "match_start": match["start"],
            "match_end": match["end"],
            "matched_text": match["matched_text"],
            "normalized_text": _normalized_match_text(match["matched_text"]),
            "source": _resolve_match_source(match["sources"]),
            "reason_tags": sorted(match["reason_tags"]),
            "confidence": match["confidence"],
        }
        for _, match in sorted(matches_by_key.items(), key=lambda item: item[0])
    ]



def _merge_body_text_match(
    matches_by_key: dict[tuple[int, int, str], dict],
    *,
    start: int,
    end: int,
    matched_text: str,
    source: str,
    reason_tag: str,
    confidence: str,
) -> None:
    key = (start, end, matched_text)
    existing = matches_by_key.get(key)
    if existing is None:
        matches_by_key[key] = {
            "start": start,
            "end": end,
            "matched_text": matched_text,
            "sources": {source},
            "reason_tags": {reason_tag},
            "confidence": confidence,
        }
        return
    existing["sources"].add(source)
    existing["reason_tags"].add(reason_tag)
    existing["confidence"] = _stronger_confidence(existing["confidence"], confidence)



def _body_text_finding(
    *,
    file_path: Path,
    relative_path: str,
    location: dict,
    payload: dict,
    action_hint: str,
    confidence: str,
    manual_review_reason: str | None,
    source: str,
    reason_tags: list[str],
) -> dict:
    finding = _finding(
        file_path=file_path,
        relative_path=relative_path,
        category="body_text",
        location=location,
        payload=payload,
        action_hint=action_hint,
        confidence=confidence,
        manual_review_reason=manual_review_reason,
    )
    finding["source"] = source
    finding["reason_tags"] = list(reason_tags)
    return finding



def _metadata_findings(file_path: Path, *, relative_path: str) -> list[dict]:
    metadata = read_metadata(file_path)
    non_empty_fields = {
        field_name: value
        for field_name, value in metadata["fields"].items()
        if value not in (None, "")
    }
    if not non_empty_fields:
        return []

    return [
        _finding(
            file_path=file_path,
            relative_path=relative_path,
            category="metadata",
            location={"scope": "document"},
            payload={"fields": non_empty_fields},
            action_hint="clear",
            confidence="high",
            manual_review_reason=None,
        )
    ]



def _image_findings(
    file_path: Path,
    *,
    relative_path: str,
    image_locations: list[dict] | None = None,
) -> list[dict]:
    with TemporaryDirectory(prefix="office-automation-detect-images-") as temp_dir:
        extracted_paths = extract_images(file_path, Path(temp_dir))
        findings: list[dict] = []
        for image_index, image_path in enumerate(extracted_paths):
            width, height = _image_dimensions(image_path)
            location = {"image_index": image_index}
            if image_locations is not None and image_index < len(image_locations):
                location.update(image_locations[image_index])
            findings.append(
                _finding(
                    file_path=file_path,
                    relative_path=relative_path,
                    category="images",
                    location=location,
                    payload={
                        "image_index": image_index,
                        "extracted_filename": image_path.name,
                        "width": width,
                        "height": height,
                    },
                    action_hint="replace_or_mask",
                    confidence="high",
                    manual_review_reason=None,
                )
            )
        return findings



def _pdf_image_locations(file_path: Path) -> list[dict]:
    document = fitz.open(file_path)
    try:
        if getattr(document, "needs_pass", False):
            raise NotImplementedError(f"Encrypted PDF files are outside detection V1 scope: '{file_path}'.")

        locations: list[dict] = []
        for page_index in range(document.page_count):
            page = document.load_page(page_index)
            for page_image_index, image_info in enumerate(page.get_image_info(xrefs=True)):
                xref = int(image_info["xref"])
                extracted = document.extract_image(xref)
                extension = str(extracted.get("ext", "")).lower().lstrip(".")
                if extension not in {"png", "jpg", "jpeg", "gif", "bmp", "tif", "tiff"}:
                    continue
                locations.append(
                    {
                        "page_number": page_index + 1,
                        "page_image_index": page_image_index,
                    }
                )
        return locations
    finally:
        document.close()



def _manual_review_finding(
    *,
    file_path: Path,
    relative_path: str,
    category: str,
    reason: str,
    location: dict,
) -> dict:
    return _finding(
        file_path=file_path,
        relative_path=relative_path,
        category=category,
        location=location,
        payload={"status": "manual_review_required"},
        action_hint="review",
        confidence="low",
        manual_review_reason=reason,
    )



def _finding(
    *,
    file_path: Path,
    relative_path: str,
    category: str,
    location: dict,
    payload: dict,
    action_hint: str,
    confidence: str,
    manual_review_reason: str | None,
) -> dict:
    normalized_location = _sorted_copy(location)
    finding = {
        "finding_id": f"{relative_path}::{category}::{_location_key(normalized_location)}",
        "file_path": str(file_path),
        "relative_path": relative_path,
        "extension": _path_extension(file_path),
        "category": category,
        "location": normalized_location,
        "payload": _sorted_copy(payload),
        "action_hint": action_hint,
        "confidence": confidence,
        "manual_review_reason": manual_review_reason,
    }
    return finding



def _finding_sort_key(finding: dict) -> tuple[str, str, int, str, str]:
    relative_path = str(finding["relative_path"])
    return (
        relative_path.casefold(),
        relative_path,
        _CATEGORY_ORDER.get(str(finding["category"]), 99),
        _location_key(finding["location"]),
        str(finding["finding_id"]),
    )



def _find_case_insensitive_occurrences(text: str, needle: str) -> list[tuple[int, int, str]]:
    if not needle:
        return []
    matches: list[tuple[int, int, str]] = []
    haystack = text.casefold()
    target = needle.casefold()
    start = 0
    while True:
        index = haystack.find(target, start)
        if index < 0:
            return matches
        end = index + len(needle)
        matches.append((index, end, text[index:end]))
        start = end



def _contains_case_insensitive(text: str, needle: str) -> bool:
    return bool(needle) and needle.casefold() in text.casefold()



def _find_context_assisted_phrase_occurrences(text: str, context_term: str) -> list[tuple[int, int, str]]:
    if not context_term:
        return []
    pattern = re.compile(
        _CONTEXT_ASSISTED_PHRASE_TEMPLATE.format(context_term=re.escape(context_term)),
        re.IGNORECASE,
    )
    matches: list[tuple[int, int, str]] = []
    for match in pattern.finditer(text):
        matched_text = match.group(1).strip()
        if len(matched_text.split()) < 2:
            continue
        start = match.start(1)
        end = match.end(1)
        matches.append((start, end, text[start:end]))
    return matches



def _normalized_match_text(value: str) -> str:
    return re.sub(r"\s+", " ", value.strip()).lower()



def _resolve_match_source(sources: set[str]) -> str:
    if sources == {"user_hint"}:
        return "user_hint"
    if sources == {"pattern"}:
        return "pattern"
    if sources == {"heuristic"}:
        return "heuristic"
    return "mixed"



def _stronger_confidence(left: str, right: str) -> str:
    rank = {"low": 0, "medium": 1, "high": 2}
    return left if rank[left] >= rank[right] else right



def _excerpt(text: str, start: int, end: int, *, radius: int = 40) -> str:
    excerpt_start = max(0, start - radius)
    excerpt_end = min(len(text), end + radius)
    excerpt = text[excerpt_start:excerpt_end]
    if excerpt_start > 0:
        excerpt = f"…{excerpt}"
    if excerpt_end < len(text):
        excerpt = f"{excerpt}…"
    return excerpt



def _pdf_text_spans(page) -> list[dict]:
    spans: list[dict] = []
    text_dict = page.get_text("dict")
    for block in text_dict.get("blocks", []):
        for line in block.get("lines", []):
            for span in line.get("spans", []):
                text = _normalized_text(span.get("text"))
                if text is None:
                    continue
                bbox = span.get("bbox")
                if not bbox:
                    continue
                spans.append(
                    {
                        "text": text,
                        "bbox": [float(bbox[0]), float(bbox[1]), float(bbox[2]), float(bbox[3])],
                    }
                )
    return spans



def _pdf_match_bbox(page, *, matched_text: str, fallback_bbox: list[float], clip_rect) -> list[float]:
    if not matched_text:
        return list(fallback_bbox)
    try:
        rects = page.search_for(matched_text, clip=clip_rect)
    except Exception:
        rects = []
    if not rects:
        return list(fallback_bbox)
    rect = rects[0]
    return [float(rect.x0), float(rect.y0), float(rect.x1), float(rect.y1)]



def _validate_target_folder(path: Path) -> None:
    if not path.exists():
        raise FileNotFoundError(f"Target folder '{path}' does not exist.")
    if path.is_symlink() or not path.is_dir():
        raise NotADirectoryError(f"Target folder '{path}' is not a directory.")



def _path_extension(path: Path) -> str:
    return path.suffix.lower().lstrip(".")



def _close_excel_workbook(workbook) -> None:
    vba_archive = getattr(workbook, "vba_archive", None)
    if vba_archive is not None:
        vba_archive.close()
    close = getattr(workbook, "close", None)
    if callable(close):
        close()



def _docx_container_text(container) -> str | None:
    parts: list[str] = []
    for paragraph in container.paragraphs:
        text = _normalized_text(paragraph.text)
        if text is not None:
            parts.append(text)
    for table in container.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    text = _normalized_text(paragraph.text)
                    if text is not None:
                        parts.append(text)
    if not parts:
        return None
    return "\n".join(parts)



def _image_dimensions(path: Path) -> tuple[int | None, int | None]:
    with Image.open(path) as image:
        width, height = image.size
    return width, height



def _location_key(location: dict) -> str:
    return json.dumps(location, ensure_ascii=False, separators=(",", ":"), sort_keys=True)



def _normalized_datetime(value) -> str | None:
    if value in (None, ""):
        return None
    isoformat = getattr(value, "isoformat", None)
    if callable(isoformat):
        return isoformat()
    return _normalized_text(value)



def _normalized_text(value) -> str | None:
    if value in (None, ""):
        return None
    text = str(value).replace("\r\n", "\n").replace("\r", "\n")
    lines = [line.strip() for line in text.split("\n")]
    cleaned_lines = [line for line in lines if line]
    if not cleaned_lines:
        return None
    return "\n".join(cleaned_lines)



def _sorted_copy(value):
    if isinstance(value, dict):
        return {key: _sorted_copy(value[key]) for key in sorted(value)}
    if isinstance(value, list):
        return [_sorted_copy(item) for item in value]
    return value
